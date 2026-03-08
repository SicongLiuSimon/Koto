#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
文件处理器 - 统一处理各种文件类型的解析和内容提取
支持格式: .docx, .pdf, .xlsx, .txt, 图片等
"""

import os
import mimetypes
from typing import Dict, Any, Optional, Tuple


class FileProcessor:
    """文件处理器 - 提取文件内容和元数据"""
    
    @staticmethod
    def process_file(filepath: str) -> Dict[str, Any]:
        """
        处理文件，返回处理结果
        
        Args:
            filepath: 文件路径
            
        Returns:
            {
                'success': bool,
                'mime_type': str,
                'filename': str,
                'text_content': str,  # 提取的文本内容
                'binary_data': bytes,  # 二进制数据（图片/PDF）
                'error': str,  # 错误信息
                'metadata': dict  # 额外元数据
            }
        """
        filename = os.path.basename(filepath)
        mime_type, _ = mimetypes.guess_type(filepath)
        if not mime_type:
            mime_type = "application/octet-stream"
        
        result = {
            'success': False,
            'mime_type': mime_type,
            'filename': filename,
            'text_content': '',
            'binary_data': None,
            'error': '',
            'metadata': {}
        }
        
        try:
            # 根据文件类型分发处理
            if mime_type.startswith('image'):
                return FileProcessor._process_image(filepath, result)
            elif filename.endswith('.pdf') or mime_type == 'application/pdf':
                return FileProcessor._process_pdf(filepath, result)
            elif filename.endswith(('.doc', '.docx')):
                return FileProcessor._process_word(filepath, result)
            elif filename.endswith(('.ppt', '.pptx')):
                return FileProcessor._process_powerpoint(filepath, result)
            elif filename.endswith(('.xls', '.xlsx')):
                return FileProcessor._process_excel(filepath, result)
            elif mime_type.startswith('text') or filename.endswith('.txt'):
                return FileProcessor._process_text(filepath, result)
            else:
                # 未知类型，尝试作为文本读取
                return FileProcessor._process_text(filepath, result)
                
        except Exception as e:
            result['error'] = f"处理文件失败: {str(e)}"
            result['success'] = False
            return result
    
    @staticmethod
    def _process_image(filepath: str, result: Dict) -> Dict:
        """处理图片文件"""
        try:
            with open(filepath, 'rb') as f:
                result['binary_data'] = f.read()
            
            # 获取图片尺寸（可选）
            try:
                from PIL import Image
                with Image.open(filepath) as img:
                    result['metadata']['dimensions'] = f"{img.width}x{img.height}"
                    result['metadata']['format'] = img.format
            except ImportError:
                pass  # PIL未安装，跳过元数据
            
            result['success'] = True
            print(f"[FileProcessor] 成功处理图片: {result['filename']}")
            return result
            
        except Exception as e:
            result['error'] = f"读取图片失败: {str(e)}"
            return result
    
    @staticmethod
    def _process_pdf(filepath: str, result: Dict) -> Dict:
        """处理PDF文件 - 优先以二进制字节提交给 Gemini（支持图表/排版），同时尝试提取文本作为补充"""
        try:
            # 始终读取原始字节（Gemini 原生支持 application/pdf 内联 blob）
            with open(filepath, 'rb') as f:
                raw_bytes = f.read()
            result['binary_data'] = raw_bytes
            result['mime_type'] = 'application/pdf'

            # 尝试用 PyPDF2 提取文本，供质量评估使用
            try:
                import PyPDF2
                import warnings
                with warnings.catch_warnings():
                    warnings.simplefilter("ignore")  # 屏蔽 GBK-EUC-H 等编码警告
                    with open(filepath, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        text_parts = []
                        max_pages = min(10, len(reader.pages))
                        for page_num in range(max_pages):
                            page = reader.pages[page_num]
                            text_parts.append(page.extract_text() or "")
                        extracted = '\n'.join(text_parts)

                # 评估文本质量：如果乱码比例高，放弃文本
                def _is_garbled(text: str) -> bool:
                    if not text or len(text) < 20:
                        return True
                    # 统计无法正常显示的替换字符 / Latin1 误读的中文字符
                    garbage_chars = sum(
                        1 for c in text
                        if '\ufffd' == c                         # Unicode 替换字符
                        or (0x00C0 <= ord(c) <= 0x00FF and '\u4e00' <= text[max(0, text.index(c)-2):text.index(c)+3][-1] <= '\u9fff')
                    )
                    # 乱码特征：大量 xad/Ò series Latin 扩展字符出现在中文上下文
                    latin_ext = sum(1 for c in text if 0x00C0 <= ord(c) <= 0x00FF)
                    cjk_chars  = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
                    total = len(text)
                    if total == 0:
                        return True
                    # 若 Latin 扩展字符多于 CJK 且比例超 15%，判断为乱码
                    if latin_ext > cjk_chars and latin_ext / total > 0.15:
                        return True
                    return False

                if extracted.strip() and not _is_garbled(extracted):
                    result['text_content'] = extracted
                    result['metadata']['pages'] = len(reader.pages)
                    result['metadata']['extracted_pages'] = max_pages
                    result['metadata']['text_quality'] = 'good'
                    print(f"[FileProcessor] PDF 文本质量良好，文本+二进制双模式: {result['filename']}")
                else:
                    result['metadata']['text_quality'] = 'garbled'
                    print(f"[FileProcessor] PDF 文本乱码或为空，改用纯二进制模式: {result['filename']}")

            except Exception as e:
                result['metadata']['text_quality'] = 'extract_failed'
                print(f"[FileProcessor] PDF 文本提取失败，使用纯二进制模式: {e}")

            result['metadata']['pages'] = result['metadata'].get('pages', '?')
            result['success'] = True
            print(f"[FileProcessor] PDF 已加载二进制: {result['filename']} ({len(raw_bytes):,} bytes)")
            return result

        except Exception as e:
            result['error'] = f"读取PDF失败: {str(e)}"
            return result
    
    @staticmethod
    def _process_word(filepath: str, result: Dict) -> Dict:
        """处理Word文档"""
        try:
            try:
                from docx import Document
                doc = Document(filepath)
                
                # 提取所有段落文本
                paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
                result['text_content'] = '\n'.join(paragraphs)
                
                # 提取表格内容
                if doc.tables:
                    result['text_content'] += '\n\n[表格内容]:\n'
                    for table in doc.tables:
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            result['text_content'] += ' | '.join(cells) + '\n'
                
                result['metadata']['paragraphs'] = len(paragraphs)
                result['metadata']['tables'] = len(doc.tables)
                result['metadata']['chars'] = len(result['text_content'])
                
                result['success'] = True
                print(f"[FileProcessor] 成功提取Word文档: {result['filename']} ({len(result['text_content'])} 字符)")
                return result
                
            except ImportError:
                result['error'] = "未安装python-docx库，无法读取Word文档。请运行: pip install python-docx"
                return result
                
        except Exception as e:
            result['error'] = f"读取Word文档失败: {str(e)}"
            return result

    @staticmethod
    def _process_powerpoint(filepath: str, result: Dict) -> Dict:
        """处理PowerPoint文档（优先 .pptx）"""
        try:
            from pptx import Presentation

            presentation = Presentation(filepath)
            slide_texts = []

            for idx, slide in enumerate(presentation.slides, start=1):
                parts = []
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text:
                            text = shape.text.strip()
                            if text:
                                parts.append(text)
                    except Exception:
                        continue

                if parts:
                    slide_texts.append(f"[Slide {idx}]\n" + "\n".join(parts))

            result['text_content'] = "\n\n".join(slide_texts)
            result['metadata']['slides'] = len(presentation.slides)
            result['metadata']['extracted_slides'] = len(slide_texts)

            if result['text_content'].strip():
                result['success'] = True
                print(f"[FileProcessor] 成功提取PowerPoint文档: {result['filename']} ({len(result['text_content'])} 字符)")
            else:
                # 无文本时仍标记成功，避免后续把它当成错误文件
                result['success'] = True
                print(f"[FileProcessor] PowerPoint提取完成但无可读文本: {result['filename']}")

            return result

        except Exception as e:
            result['error'] = f"读取PowerPoint文档失败: {str(e)}"
            return result
    
    @staticmethod
    def _process_excel(filepath: str, result: Dict) -> Dict:
        """处理Excel文件"""
        try:
            try:
                import pandas as pd
                
                # 读取Excel文件
                xls = pd.ExcelFile(filepath)
                text_parts = [f"Excel工作簿: {result['filename']}"]
                text_parts.append(f"工作表: {', '.join(xls.sheet_names)}\n")
                
                # 读取前3个工作表
                for sheet_name in xls.sheet_names[:3]:
                    df = pd.read_excel(filepath, sheet_name=sheet_name)
                    text_parts.append(f"\n=== 工作表: {sheet_name} ===")
                    text_parts.append(df.to_string(max_rows=50, max_cols=10))
                
                result['text_content'] = '\n'.join(text_parts)
                result['metadata']['sheets'] = len(xls.sheet_names)
                result['metadata']['extracted_sheets'] = min(3, len(xls.sheet_names))
                
                result['success'] = True
                print(f"[FileProcessor] 成功提取Excel数据: {result['filename']} ({len(xls.sheet_names)} 个工作表)")
                return result
                
            except ImportError:
                result['error'] = "未安装pandas/openpyxl库，无法读取Excel文件。请运行: pip install pandas openpyxl"
                return result
                
        except Exception as e:
            result['error'] = f"读取Excel文件失败: {str(e)}"
            return result
    
    @staticmethod
    def _process_text(filepath: str, result: Dict) -> Dict:
        """处理文本文件 - 尝试多种编码"""
        encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    result['text_content'] = f.read()
                
                result['metadata']['encoding'] = encoding
                result['metadata']['lines'] = result['text_content'].count('\n') + 1
                result['metadata']['chars'] = len(result['text_content'])
                
                result['success'] = True
                print(f"[FileProcessor] 成功读取文本文件: {result['filename']} (编码: {encoding})")
                return result
                
            except (UnicodeDecodeError, LookupError):
                continue
        
        # 所有编码都失败
        result['error'] = f"无法解码文件，尝试过的编码: {', '.join(encodings)}"
        return result
    
    @staticmethod
    def format_result_for_chat(result: Dict, user_message: str = "") -> Tuple[str, Optional[Dict]]:
        """
        格式化处理结果，用于发送给AI聊天
        
        Returns:
            (formatted_message, file_data)
            formatted_message: 包含文件内容的消息
            file_data: 如果是二进制文件（图片/PDF），返回 {'mime_type': str, 'data': bytes}
        """
        if not result['success']:
            error_msg = f"❌ 文件处理失败: {result['error']}"
            return f"{user_message}\n\n{error_msg}", None
        
        # 如果有二进制数据（图片/PDF）
        if result['binary_data']:
            message = user_message
            file_data = {
                'mime_type': result['mime_type'],
                'data': result['binary_data']
            }
            return message, file_data
        
        # 如果有文本内容
        if result['text_content']:
            metadata_str = ""
            if result['metadata']:
                meta_items = [f"{k}: {v}" for k, v in result['metadata'].items()]
                metadata_str = f" ({', '.join(meta_items)})"
            
            formatted = f"{user_message}\n\n"
            formatted += f"📄 文件: {result['filename']}{metadata_str}\n\n"
            formatted += f"=== 文件内容 ===\n{result['text_content']}"
            
            return formatted, None
        
        # 未知情况
        return f"{user_message}\n\n⚠️ 文件已上传但未提取到内容", None


# 便捷函数
def process_uploaded_file(filepath: str, user_message: str = "") -> Tuple[str, Optional[Dict]]:
    """
    处理上传的文件，返回格式化的消息和文件数据
    
    Args:
        filepath: 文件路径
        user_message: 用户的消息
        
    Returns:
        (formatted_message, file_data)
    """
    processor = FileProcessor()
    result = processor.process_file(filepath)
    return processor.format_result_for_chat(result, user_message)
