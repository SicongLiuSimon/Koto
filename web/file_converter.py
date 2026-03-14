#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
file_converter.py — 多格式文件转换引擎

支持转换矩阵:
  文档类:  .docx/.doc ↔ .pdf / .txt / .md   (.doc→.docx 也支持)
  表格类:  .xlsx/.xls ↔ .csv  (.xls→.xlsx 也支持)
  演示类:  .pptx/.ppt → .txt / .pdf
  图片类:  .jpg/.png/.webp/.bmp/.gif 互转

优先级顺序 (文档→PDF):
  1. Microsoft Word/PowerPoint COM（精确还原，仅 Windows）
  2. LibreOffice --headless（跨平台，需安装）
  3. reportlab / python-docx 文本重建（无需额外安装，可能丢失格式）

Skill entry_point: web.file_converter:skill_entry
Agent tool 入口:   FileConverterPlugin → convert_file()
"""

from __future__ import annotations

import os
import re
import shutil
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

# ── 支持的转换矩阵 ────────────────────────────────────────────────────────────
CONVERSION_MATRIX: Dict[str, List[str]] = {
    ".docx":     [".pdf", ".txt", ".md"],
    ".doc":      [".docx", ".pdf", ".txt", ".md"],
    ".pdf":      [".docx", ".txt"],
    ".txt":      [".docx", ".pdf", ".md"],
    ".md":       [".docx", ".pdf", ".html", ".txt"],
    ".markdown": [".docx", ".pdf", ".html", ".txt"],
    ".xlsx":     [".csv", ".txt"],
    ".xls":      [".xlsx", ".csv", ".txt"],
    ".csv":      [".xlsx"],
    ".pptx":     [".txt", ".pdf"],
    ".ppt":      [".txt", ".pdf"],
    # images
    ".jpg":  [".png", ".webp", ".bmp", ".jpeg"],
    ".jpeg": [".png", ".webp", ".bmp", ".jpg"],
    ".png":  [".jpg", ".jpeg", ".webp", ".bmp"],
    ".webp": [".png", ".jpg", ".jpeg"],
    ".bmp":  [".png", ".jpg", ".jpeg"],
    ".gif":  [".png", ".jpg"],
}

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif", ".tiff", ".tif"}

# 用户自然语言 → 扩展名 映射
FORMAT_ALIASES: Dict[str, str] = {
    "word": ".docx",   "doc": ".docx",   "docx": ".docx",
    "pdf": ".pdf",
    "txt": ".txt",     "text": ".txt",   "纯文本": ".txt",
    "md": ".md",       "markdown": ".md",
    "ppt": ".pptx",    "pptx": ".pptx",  "powerpoint": ".pptx", "幻灯片": ".pptx",
    "xls": ".xlsx",    "xlsx": ".xlsx",  "excel": ".xlsx",      "表格": ".xlsx",
    "csv": ".csv",
    "html": ".html",   "htm": ".html",
    "jpg": ".jpg",     "jpeg": ".jpg",
    "png": ".png",
    "webp": ".webp",
    "bmp": ".bmp",
    "gif": ".gif",
}

# 中文格式关键词
CN_FORMAT_PATTERNS: List[Tuple[str, str]] = [
    (r"word文档|docx?文件", ".docx"),
    (r"pdf文件?|PDF", ".pdf"),
    (r"纯文本|txt文件", ".txt"),
    (r"markdown|md文件", ".md"),
    (r"ppt(?:x)?|幻灯片|演示文稿", ".pptx"),
    (r"excel|xlsx?|电子表格", ".xlsx"),
    (r"csv|逗号", ".csv"),
    (r"html?网页", ".html"),
    (r"png图片?", ".png"),
    (r"jpe?g图片?", ".jpg"),
    (r"webp", ".webp"),
]


# ── 公开接口 ──────────────────────────────────────────────────────────────────

def get_supported_conversions() -> Dict[str, List[str]]:
    """返回所有支持的转换对（src_ext → [tgt_ext, ...]）。"""
    return {k: list(v) for k, v in CONVERSION_MATRIX.items()}


def convert(
    source_path: str,
    target_format: str,
    output_path: Optional[str] = None,
    output_dir: Optional[str] = None,
) -> Dict[str, Any]:
    """
    将 source_path 转换为 target_format。

    Args:
        source_path:   源文件绝对路径。
        target_format: 目标格式，支持 "pdf" / ".pdf" / "word" / "docx" 等别名。
        output_path:   指定输出文件路径（可选，优先于 output_dir）。
        output_dir:    输出目录（可选，默认与源文件同目录）。

    Returns::
        {
            "success":     bool,
            "output_path": str,   # 转换后文件的绝对路径
            "from_format": str,   # 源扩展名（不含点）
            "to_format":   str,   # 目标扩展名（不含点）
            "message":     str,   # 人类可读结果
            "warning":     str,   # 降级提示（可为空）
            "error":       str,   # 错误信息（success=False 时）
        }
    """
    source_path = os.path.abspath(source_path)
    if not os.path.exists(source_path):
        return _err(f"文件不存在: {source_path}")

    src_ext = Path(source_path).suffix.lower()

    # 标准化目标格式
    tgt_key = target_format.strip().lower().lstrip(".")
    tgt_ext = FORMAT_ALIASES.get(tgt_key, f".{tgt_key}")
    if not tgt_ext.startswith("."):
        tgt_ext = f".{tgt_ext}"

    if src_ext == tgt_ext:
        return _err(f"源格式与目标格式相同（{src_ext}），无需转换")

    if src_ext not in CONVERSION_MATRIX:
        return _err(
            f"不支持的源格式: {src_ext}\n"
            f"支持的来源格式: {', '.join(sorted(CONVERSION_MATRIX.keys()))}"
        )

    allowed_targets = CONVERSION_MATRIX[src_ext]
    if tgt_ext not in allowed_targets:
        return _err(
            f"不支持 {src_ext} → {tgt_ext} 的转换。\n"
            f"{src_ext} 可转换为: {', '.join(allowed_targets)}"
        )

    # 计算输出路径
    if not output_path:
        out_dir = output_dir or str(Path(source_path).parent)
        os.makedirs(out_dir, exist_ok=True)
        output_path = os.path.join(out_dir, Path(source_path).stem + tgt_ext)

    try:
        out, warning = _dispatch(source_path, src_ext, tgt_ext, output_path)
        return {
            "success":     True,
            "output_path": out,
            "from_format": src_ext.lstrip("."),
            "to_format":   tgt_ext.lstrip("."),
            "message":     f"✅ 转换成功：{Path(source_path).name} → {Path(out).name}",
            "warning":     warning,
            "error":       "",
        }
    except Exception as exc:
        return _err(f"转换失败（{src_ext} → {tgt_ext}）: {exc}")


# ── Skill entry_point ─────────────────────────────────────────────────────────

def skill_entry(user_input: str, context: Dict[str, Any]) -> Dict[str, Any]:
    """
    SkillCapabilityRegistry entry_point，签名: fn(user_input, context) -> Any。

    context 期望字段:
        file_path   : 要转换的文件路径（必需）
        output_dir  : 输出目录（可选）
    """
    file_path = context.get("file_path", "")
    output_dir = context.get("output_dir", "")

    if not file_path:
        return {
            "success": False,
            "message": "❌ 未找到要转换的文件，请先上传文件或在 context 中提供 file_path",
        }

    target_fmt = _extract_target_format(user_input)
    if not target_fmt:
        supported = ", ".join(sorted(set(FORMAT_ALIASES.keys())))
        return {
            "success": False,
            "message": (
                f"❌ 未能从输入中识别目标格式。\n"
                f"请明确说明，例如：「转换成 PDF」「导出为 Excel」。\n"
                f"支持的格式: {supported}"
            ),
        }

    return convert(file_path, target_fmt, output_dir=output_dir or None)


# ── 内部调度 ──────────────────────────────────────────────────────────────────

def _dispatch(src: str, src_ext: str, tgt_ext: str, out: str) -> Tuple[str, str]:
    """根据 (src_ext, tgt_ext) 调用对应的转换函数。"""
    # 图片
    if src_ext in IMAGE_EXTS and tgt_ext in IMAGE_EXTS:
        return _img_to_img(src, out, tgt_ext)

    key = (src_ext, tgt_ext)
    dispatch_table = {
        (".docx", ".pdf"):  _docx_to_pdf,
        (".doc",  ".pdf"):  _docx_to_pdf,
        (".docx", ".txt"):  _docx_to_txt,
        (".doc",  ".txt"):  _docx_to_txt,
        (".docx", ".md"):   _docx_to_md,
        (".doc",  ".md"):   _docx_to_md,
        (".doc",  ".docx"): _doc_to_docx,
        (".pdf",  ".docx"): _pdf_to_docx,
        (".pdf",  ".txt"):  _pdf_to_txt,
        (".txt",  ".docx"): _txt_to_docx,
        (".txt",  ".pdf"):  _txt_to_pdf,
        (".txt",  ".md"):   _text_copy,
        (".md",   ".docx"): _md_to_docx,
        (".md",   ".pdf"):  _md_to_pdf,
        (".md",   ".html"): _md_to_html,
        (".md",   ".txt"):  _text_copy,
        (".markdown", ".docx"): _md_to_docx,
        (".markdown", ".pdf"):  _md_to_pdf,
        (".markdown", ".html"): _md_to_html,
        (".markdown", ".txt"):  _text_copy,
        (".xlsx", ".csv"):  _xlsx_to_csv,
        (".xlsx", ".txt"):  _xlsx_to_txt,
        (".xls",  ".xlsx"): _xls_to_xlsx,
        (".xls",  ".csv"):  _xlsx_to_csv,
        (".xls",  ".txt"):  _xlsx_to_txt,
        (".csv",  ".xlsx"): _csv_to_xlsx,
        (".pptx", ".txt"):  _pptx_to_txt,
        (".ppt",  ".txt"):  _pptx_to_txt,
        (".pptx", ".pdf"):  _pptx_to_pdf,
        (".ppt",  ".pdf"):  _pptx_to_pdf,
    }
    fn = dispatch_table.get(key)
    if fn is None:
        raise NotImplementedError(f"转换路径未实现: {src_ext} → {tgt_ext}")
    return fn(src, out)


# ── 图片 → 图片 ───────────────────────────────────────────────────────────────

def _img_to_img(src: str, out: str, tgt_ext: str) -> Tuple[str, str]:
    from PIL import Image
    fmt_map = {
        ".jpg": "JPEG", ".jpeg": "JPEG", ".png": "PNG",
        ".webp": "WEBP", ".bmp": "BMP", ".gif": "GIF",
    }
    fmt = fmt_map.get(tgt_ext, tgt_ext.lstrip(".").upper())
    with Image.open(src) as img:
        if fmt == "JPEG" and img.mode in ("RGBA", "P", "LA"):
            img = img.convert("RGB")
        img.save(out, fmt)
    return out, ""


# ── DOCX/DOC → PDF ───────────────────────────────────────────────────────────

def _docx_to_pdf(src: str, out: str) -> Tuple[str, str]:
    # 优先: Word COM（Windows, 精确还原）
    try:
        import win32com.client
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(src))
        doc.SaveAs(os.path.abspath(out), FileFormat=17)  # 17 = wdFormatPDF
        doc.Close(False)
        word.Quit()
        return out, ""
    except Exception:
        pass

    # 回退: LibreOffice
    lo = _try_libreoffice(src, str(Path(out).parent), "pdf")
    if lo:
        if lo != out:
            shutil.move(lo, out)
        return out, ""

    # 最后手段: reportlab 文本重建
    from docx import Document as DocxDoc
    doc = DocxDoc(src)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    _text_to_pdf_reportlab(text, out, title=Path(src).stem)
    return (
        out,
        "⚠️ 未检测到 Microsoft Word 或 LibreOffice，已用文本重建 PDF（格式/图片可能丢失）",
    )


# ── DOCX/DOC → TXT ───────────────────────────────────────────────────────────

def _docx_to_txt(src: str, out: str) -> Tuple[str, str]:
    from docx import Document as DocxDoc
    doc = DocxDoc(src)
    lines = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            lines.append(" | ".join(c.text.strip() for c in row.cells))
    Path(out).write_text("\n".join(lines), encoding="utf-8")
    return out, ""


# ── DOCX/DOC → MD ────────────────────────────────────────────────────────────

def _docx_to_md(src: str, out: str) -> Tuple[str, str]:
    from docx import Document as DocxDoc
    doc = DocxDoc(src)
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue
        style = para.style.name
        if style.startswith("Heading 1"):
            lines.append(f"# {text}")
        elif style.startswith("Heading 2"):
            lines.append(f"## {text}")
        elif style.startswith("Heading 3"):
            lines.append(f"### {text}")
        elif "List" in style:
            lines.append(f"- {text}")
        else:
            lines.append(text)
    Path(out).write_text("\n".join(lines), encoding="utf-8")
    return out, "⚠️ 图片、复杂表格样式不包含在 Markdown 输出中"


# ── DOC → DOCX ───────────────────────────────────────────────────────────────

def _doc_to_docx(src: str, out: str) -> Tuple[str, str]:
    try:
        import win32com.client
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(src))
        doc.SaveAs(os.path.abspath(out), FileFormat=16)  # 16 = wdFormatDocumentDefault
        doc.Close(False)
        word.Quit()
        return out, ""
    except Exception:
        pass

    try:
        from docx import Document as DocxDoc
        d = DocxDoc(src)
        d.save(out)
        return out, ""
    except Exception:
        pass

    lo = _try_libreoffice(src, str(Path(out).parent), "docx")
    if lo:
        if lo != out:
            shutil.move(lo, out)
        return out, ""

    raise RuntimeError("无法转换 .doc 文件（需要 Microsoft Word 或 LibreOffice）")


# ── PDF → DOCX ───────────────────────────────────────────────────────────────

def _pdf_to_docx(src: str, out: str) -> Tuple[str, str]:
    # 优先: pdf2docx（版式还原最好）
    try:
        from pdf2docx import Converter
        cv = Converter(src)
        cv.convert(out, start=0, end=None)
        cv.close()
        return out, ""
    except ImportError:
        pass
    except Exception:
        pass

    # 回退: 文本提取 → 重建 DOCX
    from docx import Document as DocxDoc
    text = _extract_pdf_text(src)
    doc = DocxDoc()
    doc.add_heading(Path(src).stem, 0)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line.strip())
    doc.save(out)
    return (
        out,
        "⚠️ PDF→DOCX 仅提取文本，原始版式/图片未保留（如需完美还原可安装 `pip install pdf2docx`）",
    )


# ── PDF → TXT ────────────────────────────────────────────────────────────────

def _pdf_to_txt(src: str, out: str) -> Tuple[str, str]:
    text = _extract_pdf_text(src)
    Path(out).write_text(text, encoding="utf-8")
    return out, ""


def _extract_pdf_text(src: str) -> str:
    for mod_name in ("pypdf", "PyPDF2"):
        try:
            mod = __import__(mod_name)
            PdfReader = getattr(mod, "PdfReader")
            with open(src, "rb") as f:
                reader = PdfReader(f)
                return "\n".join(page.extract_text() or "" for page in reader.pages)
        except ImportError:
            continue
        except Exception:
            break
    try:
        import pdfplumber
        with pdfplumber.open(src) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages)
    except Exception:
        pass
    raise RuntimeError("无法提取 PDF 文本（需要 pypdf 或 pdfplumber）")


# ── TXT → DOCX ───────────────────────────────────────────────────────────────

def _txt_to_docx(src: str, out: str) -> Tuple[str, str]:
    text = Path(src).read_text(encoding="utf-8", errors="replace")
    from docx import Document as DocxDoc
    doc = DocxDoc()
    doc.add_heading(Path(src).stem, 0)
    for line in text.split("\n"):
        m = re.match(r'^(#{1,6})\s+(.*)', line)
        if m:
            doc.add_heading(m.group(2), min(len(m.group(1)), 4))
        elif line.strip():
            doc.add_paragraph(line.rstrip())
        else:
            doc.add_paragraph("")
    doc.save(out)
    return out, ""


# ── TXT → PDF ────────────────────────────────────────────────────────────────

def _txt_to_pdf(src: str, out: str) -> Tuple[str, str]:
    text = Path(src).read_text(encoding="utf-8", errors="replace")
    _text_to_pdf_reportlab(text, out, title=Path(src).stem)
    return out, ""


# ── TXT ↔ MD（纯文本复制）────────────────────────────────────────────────────

def _text_copy(src: str, out: str) -> Tuple[str, str]:
    shutil.copy2(src, out)
    return out, ""


# ── MD → DOCX ────────────────────────────────────────────────────────────────

def _md_to_docx(src: str, out: str) -> Tuple[str, str]:
    text = Path(src).read_text(encoding="utf-8", errors="replace")
    from docx import Document as DocxDoc
    doc = DocxDoc()
    doc.add_heading(Path(src).stem, 0)
    for line in text.split("\n"):
        stripped = line.rstrip()
        if not stripped:
            doc.add_paragraph("")
            continue
        m = re.match(r'^(#{1,6})\s+(.*)', stripped)
        if m:
            doc.add_heading(m.group(2), min(len(m.group(1)), 4))
            continue
        if re.match(r'^[-*+]\s+', stripped):
            doc.add_paragraph(re.sub(r'^[-*+]\s+', '', stripped), style="List Bullet")
            continue
        if re.match(r'^\d+\.\s+', stripped):
            doc.add_paragraph(re.sub(r'^\d+\.\s+', '', stripped), style="List Number")
            continue
        clean = re.sub(r'\*\*?(.*?)\*\*?', r'\1', stripped)
        clean = re.sub(r'`(.*?)`', r'\1', clean)
        doc.add_paragraph(clean)
    doc.save(out)
    return out, "⚠️ 图片链接、HTML 块未包含在 DOCX 输出中"


# ── MD → PDF ─────────────────────────────────────────────────────────────────

def _md_to_pdf(src: str, out: str) -> Tuple[str, str]:
    text = Path(src).read_text(encoding="utf-8", errors="replace")
    _text_to_pdf_reportlab(text, out, title=Path(src).stem, is_markdown=True)
    return out, "⚠️ Markdown 标题样式已尽量保留，图片链接不包含在 PDF 中"


# ── MD → HTML ────────────────────────────────────────────────────────────────

def _md_to_html(src: str, out: str) -> Tuple[str, str]:
    text = Path(src).read_text(encoding="utf-8", errors="replace")
    try:
        import markdown as md_lib
        html_body = md_lib.markdown(text, extensions=["tables", "fenced_code"])
    except ImportError:
        html_body = re.sub(r'^# (.+)$', r'<h1>\1</h1>', text, flags=re.M)
        html_body = re.sub(r'^## (.+)$', r'<h2>\1</h2>', html_body, flags=re.M)
        html_body = re.sub(r'^### (.+)$', r'<h3>\1</h3>', html_body, flags=re.M)
        html_body = html_body.replace("\n", "<br>\n")

    title_safe = re.sub(r'[<>&"\'\\]', '', Path(src).stem)
    html = (
        f'<!DOCTYPE html>\n<html lang="zh">\n<head>\n<meta charset="UTF-8">\n'
        f'<title>{title_safe}</title>\n'
        '<style>body{font-family:sans-serif;max-width:840px;margin:40px auto;'
        'line-height:1.7;color:#333}h1,h2,h3{color:#222}code{background:#f4f4f4;'
        'padding:2px 4px;border-radius:3px}</style>\n'
        f'</head>\n<body>\n{html_body}\n</body>\n</html>\n'
    )
    Path(out).write_text(html, encoding="utf-8")
    return out, ""


# ── XLSX/XLS ↔ CSV ───────────────────────────────────────────────────────────

def _xlsx_to_csv(src: str, out: str) -> Tuple[str, str]:
    import pandas as pd
    xls = pd.ExcelFile(src)
    sheets = xls.sheet_names

    if len(sheets) == 1:
        df = pd.read_excel(src, sheet_name=sheets[0])
        df.to_csv(out, index=False, encoding="utf-8-sig")
        return out, ""

    # 多 sheet：分别导出，同时生成一个合并主文件
    base_dir = Path(out).parent
    stem = Path(out).stem
    extra_paths = []
    for sheet in sheets:
        df = pd.read_excel(src, sheet_name=sheet)
        safe = re.sub(r'[\\/:*?"<>|]', '_', sheet)
        p = str(base_dir / f"{stem}_{safe}.csv")
        df.to_csv(p, index=False, encoding="utf-8-sig")
        extra_paths.append(p)
    shutil.copy(extra_paths[0], out)
    names = ", ".join(Path(p).name for p in extra_paths)
    return out, f"⚠️ 工作簿含 {len(sheets)} 个 Sheet，已分别导出: {names}"


def _xlsx_to_txt(src: str, out: str) -> Tuple[str, str]:
    import pandas as pd
    xls = pd.ExcelFile(src)
    parts = []
    for sheet in xls.sheet_names[:5]:
        df = pd.read_excel(src, sheet_name=sheet)
        parts.append(f"=== Sheet: {sheet} ===\n{df.to_string()}")
    Path(out).write_text("\n\n".join(parts), encoding="utf-8")
    return out, ""


def _xls_to_xlsx(src: str, out: str) -> Tuple[str, str]:
    import pandas as pd
    xls = pd.ExcelFile(src)
    with pd.ExcelWriter(out, engine="openpyxl") as writer:
        for sheet in xls.sheet_names:
            df = pd.read_excel(src, sheet_name=sheet)
            df.to_excel(writer, sheet_name=sheet, index=False)
    return out, ""


def _csv_to_xlsx(src: str, out: str) -> Tuple[str, str]:
    import pandas as pd
    for enc in ("utf-8-sig", "utf-8", "gbk", "gb18030", "latin-1"):
        try:
            df = pd.read_csv(src, encoding=enc)
            break
        except Exception:
            continue
    df.to_excel(out, index=False, engine="openpyxl")
    return out, ""


# ── PPTX → TXT / PDF ─────────────────────────────────────────────────────────

def _pptx_to_txt(src: str, out: str) -> Tuple[str, str]:
    from pptx import Presentation
    prs = Presentation(src)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"=== 第 {i} 页 ===")
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
            except Exception:
                continue
        lines.append("")
    Path(out).write_text("\n".join(lines), encoding="utf-8")
    return out, ""


def _pptx_to_pdf(src: str, out: str) -> Tuple[str, str]:
    # 优先: PowerPoint COM
    try:
        import win32com.client
        ppts = win32com.client.Dispatch("PowerPoint.Application")
        ppts.Visible = False
        ppt = ppts.Presentations.Open(
            os.path.abspath(src), ReadOnly=True, Untitled=False, WithWindow=False
        )
        ppt.SaveAs(os.path.abspath(out), 32)  # 32 = ppSaveAsPDF
        ppt.Close()
        ppts.Quit()
        return out, ""
    except Exception:
        pass

    # 回退: LibreOffice
    lo = _try_libreoffice(src, str(Path(out).parent), "pdf")
    if lo:
        if lo != out:
            shutil.move(lo, out)
        return out, ""

    # 最后手段: 文本提取 → PDF
    txt_tmp = str(Path(out).with_suffix(".tmp.txt"))
    _pptx_to_txt(src, txt_tmp)
    text = Path(txt_tmp).read_text(encoding="utf-8")
    os.remove(txt_tmp)
    _text_to_pdf_reportlab(text, out, title=Path(src).stem)
    return (
        out,
        "⚠️ 未检测到 PowerPoint 或 LibreOffice，已用文本重建 PDF（排版/图片未保留）",
    )


# ── LibreOffice CLI helper ────────────────────────────────────────────────────

def _try_libreoffice(source_path: str, out_dir: str, target_fmt: str = "pdf") -> Optional[str]:
    import subprocess
    candidates = (
        "soffice",
        "libreoffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    )
    for soffice in candidates:
        if not (shutil.which(soffice) or os.path.exists(soffice)):
            continue
        try:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", target_fmt, "--outdir", out_dir, source_path],
                capture_output=True, text=True, timeout=90,
            )
            if result.returncode == 0:
                expected = os.path.join(out_dir, Path(source_path).stem + f".{target_fmt}")
                if os.path.exists(expected):
                    return expected
        except Exception:
            pass
    return None


# ── reportlab 文本 → PDF ──────────────────────────────────────────────────────

def _text_to_pdf_reportlab(
    text: str,
    out: str,
    title: str = "",
    is_markdown: bool = False,
) -> None:
    """用 reportlab 将文本（可含 Markdown 标题）写入 A4 PDF，支持中文字体。"""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    # 注册中文字体（Windows 系统字体）
    cjk_font = "Helvetica"
    for fname, fpath in [
        ("MicrosoftYaHei", r"C:\Windows\Fonts\msyh.ttc"),
        ("SimHei",         r"C:\Windows\Fonts\simhei.ttf"),
        ("SimSun",         r"C:\Windows\Fonts\simsun.ttc"),
    ]:
        if os.path.exists(fpath):
            try:
                pdfmetrics.registerFont(TTFont(fname, fpath))
                cjk_font = fname
                break
            except Exception:
                continue

    styles = getSampleStyleSheet()
    body = ParagraphStyle(
        "CJKBody", parent=styles["Normal"],
        fontName=cjk_font, fontSize=11, leading=17, wordWrap="CJK",
    )
    h1 = ParagraphStyle(
        "CJKH1", parent=styles["Heading1"],
        fontName=cjk_font, fontSize=17, spaceAfter=6,
    )
    h2 = ParagraphStyle(
        "CJKH2", parent=styles["Heading2"],
        fontName=cjk_font, fontSize=13, spaceAfter=4,
    )

    doc = SimpleDocTemplate(
        out, pagesize=A4,
        leftMargin=22 * mm, rightMargin=22 * mm,
        topMargin=22 * mm, bottomMargin=22 * mm,
    )
    story = []
    if title:
        story.append(Paragraph(_safe_rl(title), h1))
        story.append(Spacer(1, 6 * mm))

    for line in text.split("\n"):
        stripped = line.rstrip()
        if not stripped:
            story.append(Spacer(1, 3 * mm))
            continue
        if is_markdown:
            m = re.match(r'^(#{1,6})\s+(.*)', stripped)
            if m:
                style = h1 if len(m.group(1)) == 1 else h2
                story.append(Paragraph(_safe_rl(m.group(2)), style))
                continue
        story.append(Paragraph(_safe_rl(stripped), body))

    doc.build(story)


def _safe_rl(text: str) -> str:
    """转义 reportlab Paragraph 中的 HTML 特殊字符。"""
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


# ── 自然语言格式提取 ──────────────────────────────────────────────────────────

def _extract_target_format(text: str) -> Optional[str]:
    """从用户输入中提取目标格式别名，返回 FORMAT_ALIASES 的键或 None。"""
    tl = text.lower()

    # 先尝试中文关键词模式
    for pat, ext in CN_FORMAT_PATTERNS:
        if re.search(pat, tl, re.IGNORECASE):
            # 找到对应的 alias key
            for k, v in FORMAT_ALIASES.items():
                if v == ext:
                    return k

    # 再按长度倒序匹配英文 alias
    for alias in sorted(FORMAT_ALIASES.keys(), key=len, reverse=True):
        patterns = [
            rf'转(?:换|成|为|出|化)\s*(?:为|成|到|作)?\s*{re.escape(alias)}',
            rf'(?:导出|保存|另存|生成|变成)\s*(?:为|成|到|作|成)?\s*{re.escape(alias)}',
            rf'(?:to|into|as)\s+{re.escape(alias)}\b',
            rf'convert\s+(?:to\s+)?{re.escape(alias)}\b',
            rf'\b{re.escape(alias)}\b',
        ]
        for pat in patterns:
            if re.search(pat, tl, re.IGNORECASE):
                return alias
    return None


# ── 公用错误构造 ──────────────────────────────────────────────────────────────

def _err(msg: str) -> Dict[str, Any]:
    return {
        "success":     False,
        "output_path": "",
        "from_format": "",
        "to_format":   "",
        "message":     f"❌ {msg}",
        "warning":     "",
        "error":       msg,
    }
