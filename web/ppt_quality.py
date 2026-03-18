#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT quality checker.
Scores layout and content density and returns suggestions.
"""

from __future__ import annotations

import os
from typing import Any, Dict, List


class PPTQualityChecker:
    """PPT quality checker."""

    def __init__(self):
        pass

    def evaluate(self, pptx_path: str) -> Dict[str, Any]:
        """Evaluate PPT quality and return score and suggestions."""
        if not os.path.exists(pptx_path):
            return {
                "success": False,
                "error": "PPT file not found",
                "score": 0,
                "issues": ["PPT file not found"],
                "suggestions": ["Verify the output path"],
            }

        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return {
                "success": False,
                "error": "python-pptx is required: pip install python-pptx",
                "score": 0,
                "issues": ["Missing python-pptx dependency"],
                "suggestions": ["Install python-pptx to enable quality checks"],
            }

        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)

        metrics = {
            "slide_count": slide_count,
            "image_slides": 0,
            "missing_title_slides": 0,
            "avg_bullets_per_slide": 0,
            "max_bullet_length": 0,
            "max_text_chars_per_slide": 0,
        }

        bullets_total = 0
        slide_text_char_max = 0
        bullet_length_max = 0

        for slide in prs.slides:
            slide_bullets = 0
            slide_text_chars = 0
            has_title = False
            has_image = False

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_image = True

                if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                    continue

                for p in shape.text_frame.paragraphs:
                    text = (p.text or "").strip()
                    if not text:
                        continue

                    slide_text_chars += len(text)
                    font_size = None
                    if p.runs and p.runs[0].font and p.runs[0].font.size:
                        font_size = p.runs[0].font.size.pt

                    # Simple title detection: large font or top area text
                    if font_size and font_size >= 28:
                        has_title = True
                    elif (
                        shape.top is not None
                        and shape.top < 1000000
                        and len(text) <= 40
                    ):
                        has_title = True

                    # 统计要点
                    if p.level is not None and p.level >= 0:
                        slide_bullets += 1
                        bullet_length_max = max(bullet_length_max, len(text))

            if has_image:
                metrics["image_slides"] += 1
            if not has_title:
                metrics["missing_title_slides"] += 1

            bullets_total += slide_bullets
            slide_text_char_max = max(slide_text_char_max, slide_text_chars)

        metrics["avg_bullets_per_slide"] = bullets_total / max(slide_count, 1)
        metrics["max_bullet_length"] = bullet_length_max
        metrics["max_text_chars_per_slide"] = slide_text_char_max

        score, issues, suggestions = self._score(metrics)

        return {
            "success": True,
            "score": score,
            "metrics": metrics,
            "issues": issues,
            "suggestions": suggestions,
        }

    def _score(self, metrics: Dict[str, Any]) -> (int, List[str], List[str]):
        """Score based on heuristic rules."""
        score = 100
        issues: List[str] = []
        suggestions: List[str] = []

        slide_count = metrics["slide_count"]
        avg_bullets = metrics["avg_bullets_per_slide"]
        missing_titles = metrics["missing_title_slides"]
        image_slides = metrics["image_slides"]
        max_bullet_len = metrics["max_bullet_length"]
        max_chars = metrics["max_text_chars_per_slide"]

        if slide_count < 5:
            score -= 10
            issues.append("Too few slides")
            suggestions.append("Add sections or expand key content")
        if slide_count > 12:
            score -= 10
            issues.append("Too many slides")
            suggestions.append("Merge similar content to reduce repetition")

        if avg_bullets < 2:
            score -= 10
            issues.append("Too few bullets per slide")
            suggestions.append("Aim for 3-5 core bullets per slide")
        if avg_bullets > 6:
            score -= 10
            issues.append("Too many bullets per slide")
            suggestions.append("Split dense content into more slides")

        if missing_titles > 0:
            score -= min(30, missing_titles * 10)
            issues.append("Some slides are missing titles")
            suggestions.append("Ensure every slide has a clear title")

        image_ratio = image_slides / max(slide_count, 1)
        if slide_count >= 6 and image_ratio < 0.1:
            score -= 10
            issues.append("Image coverage is low")
            suggestions.append("Add visuals for key sections")
        if image_ratio > 0.7:
            score -= 5
            issues.append("Image coverage is too high")
            suggestions.append("Keep only essential visuals")

        if max_bullet_len > 120:
            score -= 10
            issues.append("Bullet text is too long")
            suggestions.append("Shorten bullets to concise phrases")
        if max_chars > 700:
            score -= 10
            issues.append("Too much text on a slide")
            suggestions.append("Split content to reduce text density")

        score = max(0, min(100, score))

        if not issues:
            suggestions.append("Quality looks good; no major changes needed")

        return score, issues, suggestions
