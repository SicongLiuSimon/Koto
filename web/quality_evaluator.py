#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
质量评估系统 - 对生成的文档进行自动评分和改进建议
支持 PPT、Word、PDF 等多种文件格式
"""

import json
import logging
import re
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)


@dataclass
class EvaluationScore:
    """评分结果"""

    overall_score: float  # 0-100
    category_scores: Dict[str, float]  # 各类别评分
    issues: List[str]  # 发现的问题列表
    suggestions: List[str]  # 改进建议
    needs_improvement: bool  # 是否需要改进
    improvement_priority: List[str]  # 改进优先级


class PPTEvaluator:
    """PPT质量评估器"""

    def __init__(self):
        self.issues = []
        self.suggestions = []

    def evaluate_pptx_file(self, file_path: str) -> EvaluationScore:
        """评估PPT文件质量"""
        try:
            from pptx import Presentation
        except ImportError:
            return EvaluationScore(
                overall_score=0,
                category_scores={},
                issues=["python-pptx 模块不可用"],
                suggestions=["请安装: pip install python-pptx"],
                needs_improvement=True,
                improvement_priority=[],
            )

        self.issues = []
        self.suggestions = []
        category_scores = {}

        try:
            prs = Presentation(file_path)

            # 1. 评估总体幻灯片数量
            slide_count = len(prs.slides)
            category_scores["slide_count"] = self._score_slide_count(slide_count)

            # 2. 评估每页内容（文本量、图片分布）
            content_score, content_issues = self._evaluate_slide_contents(prs)
            category_scores["content_distribution"] = content_score
            self.issues.extend(content_issues)

            # 3. 评估排版一致性
            layout_score = self._evaluate_layout_consistency(prs)
            category_scores["layout_consistency"] = layout_score

            # 4. 评估图片分布
            image_score, image_issues = self._evaluate_image_distribution(prs)
            category_scores["image_distribution"] = image_score
            self.issues.extend(image_issues)

            # 5. 评估颜色和视觉一致性
            visual_score = self._evaluate_visual_consistency(prs)
            category_scores["visual_consistency"] = visual_score

            # 计算总体评分
            overall_score = (
                sum(category_scores.values()) / len(category_scores)
                if category_scores
                else 0
            )

            # 生成改进建议
            self._generate_suggestions(prs, overall_score)

            # 判断是否需要改进（<75分）
            needs_improvement = overall_score < 75

            # 确定改进优先级
            improvement_priority = self._prioritize_improvements(category_scores)

            return EvaluationScore(
                overall_score=round(overall_score, 1),
                category_scores={k: round(v, 1) for k, v in category_scores.items()},
                issues=self.issues,
                suggestions=self.suggestions,
                needs_improvement=needs_improvement,
                improvement_priority=improvement_priority,
            )

        except Exception as e:
            return EvaluationScore(
                overall_score=0,
                category_scores={},
                issues=[f"文件读取失败: {str(e)}"],
                suggestions=["请检查文件是否损坏"],
                needs_improvement=True,
                improvement_priority=[],
            )

    def _score_slide_count(self, count: int) -> float:
        """评估幻灯片数量（5-20页最优）"""
        if count < 3:
            return 20.0
        elif count < 5:
            return 50.0
        elif count <= 20:
            return 100.0
        elif count <= 30:
            return 80.0
        else:
            self.issues.append(f"幻灯片过多 ({count}页)，可能不适合演示")
            self.suggestions.append("考虑将某些页面合并或删除不必要的内容")
            return 60.0

    def _evaluate_slide_contents(self, prs) -> Tuple[float, List[str]]:
        """评估所有幻灯片的内容分布"""
        issues = []
        problem_slides = []

        for idx, slide in enumerate(prs.slides, 1):
            text_count = 0
            shape_count = 0
            image_count = 0

            for shape in slide.shapes:
                shape_count += 1

                # 计算文本
                if hasattr(shape, "text"):
                    text_count += len(shape.text)

                # 计算图片
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image_count += 1

            # 检查内容过多（单页超过1500字符）
            if text_count > 1500:
                problem_slides.append(f"第{idx}页文本过多 ({text_count}字)")
                issues.append(f"第{idx}页内容超载 ({text_count}字符)")

            # 检查内容过少（封面/过渡页除外，<100字且无图片）
            if text_count < 100 and image_count == 0 and idx > 1:
                issues.append(f"第{idx}页内容过少，信息不足")

            # 检查没有图片的内容页（idx > 1 表示非封面）
            if image_count == 0 and text_count > 300 and idx > 1:
                issues.append(f"第{idx}页缺少图片，可视化不足")

        # 计算评分
        score = 100.0
        score -= len(problem_slides) * 15  # 每个过载页面 -15 分
        score -= (len(issues) - len(problem_slides)) * 5  # 其他问题 -5 分
        score = max(0, min(100, score))

        return score, issues

    def _evaluate_layout_consistency(self, prs) -> float:
        """评估排版一致性（布局、边距、字体）"""
        # 这是一个启发式评估，实际需要深度解析 XML
        try:
            layouts_used = set()
            for slide in prs.slides:
                layouts_used.add(id(slide.slide_layout))

            # 使用的布局种类越少，一致性越高
            layout_variety = len(layouts_used)
            if layout_variety <= 3:
                return 100.0
            elif layout_variety <= 5:
                return 85.0
            elif layout_variety <= 8:
                return 70.0
            else:
                self.suggestions.append("幻灯片布局种类过多，建议统一使用2-3种布局")
                return 50.0
        except Exception as e:
            logger.debug("Failed to evaluate layout variety: %s", e)
            return 75.0

    def _evaluate_image_distribution(self, prs) -> Tuple[float, List[str]]:
        """评估图片分布"""
        issues = []
        total_slides = len(prs.slides)
        slides_with_images = 0
        image_heavy_slides = []

        for idx, slide in enumerate(prs.slides, 1):
            image_count = 0
            total_area = 0

            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    image_count += 1
                    try:
                        total_area += shape.width * shape.height
                    except (AttributeError, TypeError) as e:
                        logger.debug("Failed to calculate image area: %s", e)
                        pass

            if image_count > 0:
                slides_with_images += 1

                # 检查一页上有太多图片（>3张）
                if image_count > 3:
                    image_heavy_slides.append(f"第{idx}页有{image_count}张图片，过多")

        # 计算评分
        image_coverage = slides_with_images / max(1, total_slides - 1)  # 排除封面

        if image_coverage < 0.3:
            self.suggestions.append("图片分布过少，建议至少50%的页面包含图片")
            score = 50.0
        elif image_coverage < 0.5:
            score = 70.0
        elif image_coverage <= 0.9:
            score = 95.0
        else:
            score = 85.0  # 图片太多可能挤占文本空间

        issues.extend(image_heavy_slides)

        return score, issues

    def _evaluate_visual_consistency(self, prs) -> float:
        """评估视觉一致性（颜色、字体等）"""
        # 简化版本：检查主题是否一致
        try:
            font_families = set()
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.name:
                                    font_families.add(run.font.name)

            # 字体种类 1-2 个最优，3-4 个可接受
            font_count = len(font_families)
            if font_count <= 2:
                return 100.0
            elif font_count <= 4:
                return 85.0
            else:
                self.suggestions.append(f"字体种类过多 ({font_count}种)，建议统一为2种")
                return 70.0
        except Exception as e:
            logger.debug("Failed to evaluate visual consistency: %s", e)
            return 80.0

    def _generate_suggestions(self, prs, overall_score):
        """生成改进建议"""
        if overall_score >= 90:
            self.suggestions.append("PPT质量优秀，无需改进")
        else:
            # 检查是否有缺失的最佳实践
            total_slides = len(prs.slides)
            if total_slides < 5:
                self.suggestions.append("考虑扩展内容，至少5-8页能更好地讲述故事")

            # 检查是否有非内容页（封面、目录）
            has_title_slide = any("title" in id(s) for s in prs.slides)
            if not has_title_slide:
                self.suggestions.append("建议添加标题页（封面）")

    def _prioritize_improvements(self, scores: Dict[str, float]) -> List[str]:
        """根据分数确定改进优先级"""
        priority = []
        threshold = 75

        for category, score in sorted(scores.items(), key=lambda x: x[1]):
            if score < threshold:
                if category == "content_distribution":
                    priority.append("调整内容分布 - 避免单页过多内容")
                elif category == "image_distribution":
                    priority.append("增加并均匀分布图片")
                elif category == "layout_consistency":
                    priority.append("统一幻灯片布局")
                elif category == "visual_consistency":
                    priority.append("统一视觉元素（字体、颜色）")
                elif category == "slide_count":
                    priority.append("调整幻灯片数量")

        return priority[:3]  # 最多3个优先级


class DocumentEvaluator:
    """文档质量评估器"""

    def __init__(self):
        self.issues = []
        self.suggestions = []

    def evaluate_document(
        self, content: str, doc_type: str = "docx"
    ) -> EvaluationScore:
        """评估文档内容质量

        Args:
            content: 文档内容（Markdown 或纯文本）
            doc_type: 文档类型 (docx, pdf, md)
        """
        self.issues = []
        self.suggestions = []
        category_scores = {}

        # 1. 评估文档结构
        structure_score = self._evaluate_structure(content)
        category_scores["structure"] = structure_score

        # 2. 评估内容完整性
        completeness_score = self._evaluate_completeness(content)
        category_scores["completeness"] = completeness_score

        # 3. 评估文本长度
        length_score = self._evaluate_length(content)
        category_scores["length"] = length_score

        # 4. 评估格式一致性
        format_score = self._evaluate_format(content)
        category_scores["format"] = format_score

        # 5. 评估语言质量
        language_score = self._evaluate_language(content)
        category_scores["language"] = language_score

        overall_score = (
            sum(category_scores.values()) / len(category_scores)
            if category_scores
            else 0
        )

        self._generate_suggestions(content, overall_score)

        needs_improvement = overall_score < 75
        improvement_priority = self._prioritize_improvements(category_scores)

        return EvaluationScore(
            overall_score=round(overall_score, 1),
            category_scores={k: round(v, 1) for k, v in category_scores.items()},
            issues=self.issues,
            suggestions=self.suggestions,
            needs_improvement=needs_improvement,
            improvement_priority=improvement_priority,
        )

    def _evaluate_structure(self, content: str) -> float:
        """评估文档结构（标题、段落组织）"""
        lines = content.split("\n")
        heading_count = len([l for l in lines if l.startswith("#")])

        # 检查是否有足够的标题
        if heading_count == 0:
            self.issues.append("缺少标题结构，文档不清晰")
            return 40.0
        elif heading_count == 1:
            self.issues.append("仅有一个标题，结构不够清晰")
            return 60.0
        elif heading_count <= 5:
            return 90.0
        else:
            return 85.0  # 标题过多

    def _evaluate_completeness(self, content: str) -> float:
        """评估内容完整性"""
        has_intro = bool(content.startswith("#"))
        has_sections = len([l for l in content.split("\n") if l.startswith("## ")]) > 0
        has_conclusion = "结论" in content or "总结" in content or "结尾" in content

        completeness = sum([has_intro, has_sections, has_conclusion]) / 3 * 100

        if not has_conclusion:
            self.suggestions.append("建议添加结论或总结部分")

        return completeness

    def _evaluate_length(self, content: str) -> float:
        """评估文档长度"""
        char_count = len(content)

        if char_count < 500:
            self.issues.append(f"文档过短 ({char_count}字)，信息不足")
            return 50.0
        elif char_count < 1000:
            return 75.0
        elif char_count <= 5000:
            return 100.0
        elif char_count <= 10000:
            return 85.0
        else:
            self.suggestions.append("文档过长，考虑提炼主要内容或分割成多个文档")
            return 70.0

    def _evaluate_format(self, content: str) -> float:
        """评估格式一致性"""
        # 检查列表、代码块等格式元素
        has_lists = bool(re.search(r"\n[-*]\s", content))
        has_code = bool(re.search(r"```", content))
        has_emphasis = bool(re.search(r"\*\*.*?\*\*", content))

        # 检查不一致的格式
        list_types = len(set(re.findall(r"^\s*[-*]\s", content, re.MULTILINE)))

        score = 85.0
        if list_types > 1:
            self.suggestions.append("列表格式不一致，建议统一使用 - 或 *")
            score -= 15

        if has_code and not has_emphasis:
            score += 10  # 有代码块加分

        return score

    def _evaluate_language(self, content: str) -> float:
        """评估语言质量（简单启发式）"""
        lines = content.split("\n")
        issue_count = 0

        # 检查过长的句子或段落
        for line in lines:
            if len(line) > 150:
                issue_count += 1

        if issue_count > 5:
            self.suggestions.append("存在过长的句子，建议拆分以提高可读性")

        # 检查是否有明显的语法问题（简单检查）
        if re.search(r"[\u4e00-\u9fff]{10,}", content):  # 连续10个汉字
            self.issues.append("存在过长的汉字序列，可能缺少标点")
            return max(0, 100 - issue_count * 5)

        score = 100 - min(20, issue_count * 3)
        return score

    def _generate_suggestions(self, content: str, overall_score: float):
        """生成改进建议"""
        if overall_score >= 90:
            self.suggestions.append("文档质量优秀，无需改进")

        # 检查特定的改进机会
        if "TODO" in content or "FIXME" in content:
            self.suggestions.append("存在未完成的标记 (TODO/FIXME)，建议完成或删除")

        if len(content.split("\n\n")) < 3:
            self.suggestions.append("段落较少，建议增加更多细节内容")

    def _prioritize_improvements(self, scores: Dict[str, float]) -> List[str]:
        """根据分数确定改进优先级"""
        priority = []
        threshold = 75

        for category, score in sorted(scores.items(), key=lambda x: x[1]):
            if score < threshold:
                if category == "structure":
                    priority.append("改进文档结构 - 添加更多层级标题")
                elif category == "completeness":
                    priority.append("完善内容 - 添加缺失的部分")
                elif category == "length":
                    priority.append("调整文档长度 - 扩展或精简内容")
                elif category == "format":
                    priority.append("统一格式 - 保持格式一致性")
                elif category == "language":
                    priority.append("改进语言质量 - 简化表述和拆分长句")

        return priority[:3]


# 导出主要函数
def evaluate_quality(file_type: str, content_or_path: str) -> Dict[str, Any]:
    """统一评估接口

    Args:
        file_type: "pptx" | "docx" | "pdf" | "markdown" | "text"
        content_or_path: 文件路径（对于pptx）或内容字符串（对于其他类型）

    Returns:
        评分结果字典，包含 overall_score, category_scores, issues, suggestions 等
    """
    if file_type.startswith("ppt"):
        evaluator = PPTEvaluator()
        result = evaluator.evaluate_pptx_file(content_or_path)
    else:
        evaluator = DocumentEvaluator()
        result = evaluator.evaluate_document(content_or_path, file_type)

    return {
        "overall_score": result.overall_score,
        "category_scores": result.category_scores,
        "issues": result.issues,
        "suggestions": result.suggestions,
        "needs_improvement": result.needs_improvement,
        "improvement_priority": result.improvement_priority,
    }
