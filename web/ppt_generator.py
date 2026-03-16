#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT生成器 - 高质量演示文稿生成
自动生成结构化、美观的PowerPoint演示文稿
增强版支持：多模型协作、智能配图、搜索增强、排版优化
"""

import os
import io
import re
from typing import List, Dict, Optional, Any, Tuple
from datetime import datetime
from pathlib import Path
import logging

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    logger.warning("Warning: python-pptx not installed. PPT generation will fail.")

from web.ppt_themes import get_theme, PPTTheme
from web.image_generator import ImageGenerator


class SlideContent:
    """幻灯片内容"""
    def __init__(self, title: str, content: List[str], slide_type: str = "bullet", notes: str = ""):
        self.title = title
        self.content = content
        self.slide_type = slide_type
        self.notes = notes


class PPTGenerator:
    """PPT生成器 - 高质量演示文稿"""
    
    # 预设主题配色 (primary, accent, bg, text, subtitle, light_bg)
    THEMES = {
        "business": {
            "primary": (41, 65, 122),       # 深商务蓝
            "accent": (228, 161, 27),       # 金色
            "background": (255, 255, 255),  # 白色
            "text": (51, 51, 51),           # 深灰
            "subtitle": (120, 120, 120),    # 副标题灰
            "light_bg": (240, 244, 250),    # 浅蓝背景
            "bullet_color": (41, 65, 122),  # 要点符号色
        },
        "tech": {
            "primary": (0, 100, 200),       # 科技蓝
            "accent": (0, 200, 150),        # 青绿
            "background": (248, 249, 252),  # 浅灰白
            "text": (40, 40, 50),           # 深灰
            "subtitle": (100, 110, 130),
            "light_bg": (235, 240, 248),
            "bullet_color": (0, 100, 200),
        },
        "creative": {
            "primary": (138, 43, 226),      # 紫色
            "accent": (255, 107, 53),       # 橙色
            "background": (255, 255, 255),
            "text": (50, 50, 50),
            "subtitle": (130, 100, 150),
            "light_bg": (245, 240, 252),
            "bullet_color": (138, 43, 226),
        },
        "minimal": {
            "primary": (30, 30, 30),        # 黑
            "accent": (220, 50, 50),        # 红
            "background": (255, 255, 255),
            "text": (50, 50, 50),
            "subtitle": (130, 130, 130),
            "light_bg": (245, 245, 245),
            "bullet_color": (220, 50, 50),
        },
    }
    
    # CJK 字体优先级
    CJK_FONTS = ['Microsoft YaHei', '微软雅黑', 'PingFang SC', 'Noto Sans CJK SC', 'SimHei', 'Arial']
    LATIN_FONTS = ['Calibri', 'Segoe UI', 'Arial', 'Helvetica']
    
    def __init__(self, theme: str = "business"):
        self.theme = theme.lower()
        
        # Try to load from new Theme System
        try:
            from web.ppt_themes import get_theme
            ppt_theme = get_theme(self.theme)
            if ppt_theme:
                # Map new Theme structure to old Dict structure for backward compatibility
                self.colors = {
                    "primary": ppt_theme.primary_color,
                    "accent": ppt_theme.accent_color,
                    "background": ppt_theme.background_color,
                    "text": ppt_theme.text_color,
                    "subtitle": (100, 100, 100), # Default 
                    "light_bg": ppt_theme.secondary_color,
                    "bullet_color": ppt_theme.primary_color
                }
                # Also store the full theme object for new features
                self.ppt_theme = ppt_theme
            else:
                # Fallback to internal dict
                self.colors = self.THEMES.get(self.theme, self.THEMES["business"])
                self.ppt_theme = None
        except ImportError:
             # Fallback to internal dict if import fails
            self.colors = self.THEMES.get(self.theme, self.THEMES["business"])
            self.ppt_theme = None

    @staticmethod
    def _clean_markdown(text: str, strip_bold: bool = False) -> str:
        """清理文本中残留的 markdown 标记和 AI 对话痕迹，返回纯净的展示文本
        Args:
            strip_bold: 如果 True，同时去除 **bold** 标记。
                        渲染器如果自己处理粗体，传 False 保留 **。
        """
        import re
        if not text:
            return text
        # 去除标题标记 ### ## # 
        text = re.sub(r'^#{1,6}\s+', '', text.strip())
        # 去除加粗 **text** → text（仅当 strip_bold=True）
        if strip_bold:
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        # 去除斜体 *text*（但不影响 **bold**）
        text = re.sub(r'(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)', r'\1', text)
        # 去除行内代码 `text` → text
        text = re.sub(r'`(.+?)`', r'\1', text)
        # 去除行首 bullet 标记 - • 
        text = re.sub(r'^[\s]*[-•]\s+', '', text)
        # 去除数字编号 1. 2. 
        text = re.sub(r'^[\s]*\d+[.、)]\s+', '', text)
        # 去除 ~~strikethrough~~
        text = re.sub(r'~~(.+?)~~', r'\1', text)
        # 去除 markdown 链接 [text](url) → text
        text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
        # 去除连续单引号（AI 模型常见输出）
        text = re.sub(r"'{2,}", "", text)
        # 去除连续星号（***、**** 等）
        text = re.sub(r'\*{3,}', '', text)
        # 去除代码块标记
        text = re.sub(r'```[\s\S]*?```', '', text)
        # 去除 AI 对话痕迹
        ai_patterns = [
            r'^(当然可以[!！。，,\s]*)',
            r'^(好的[!！。，,\s]*)',
            r'^(以下是.*?[：:]\s*)',
            r'^(Sure[!,.\s]*)',
            r'^(Here(?:\'s| is).*?[:：]\s*)',
        ]
        for pat in ai_patterns:
            text = re.sub(pat, '', text, flags=re.IGNORECASE)
        return text.strip()
    

    def generate_from_outline(
        self,
        title: str,
        outline: List[Dict[str, Any]],
        output_path: str,
        subtitle: str = "",
        author: str = "Koto AI",
        progress_callback: Optional[Any] = None,
        enable_ai_images: bool = True  # New Flag
    ) -> Dict[str, Any]:
        """从大纲生成高质量PPT - 支持智能配图"""
        
        # Initialize Image Generator if needed
        image_gen = None
        if enable_ai_images:
            try:
                from web.image_generator import ImageGenerator
                image_gen = ImageGenerator()
                logger.info("[PPT] 🖼️ AI Image Generator Initialized")
            except ImportError:
                logger.warning("[PPT] ⚠️ ImageGenerator not found, skipping AI images")

        _type_name_map = {"detail": "详细页", "overview": "概览页", "highlight": "亮点页",
                          "divider": "过渡页", "comparison": "对比页", "image_full": "图片页"}
        
        prs = Presentation()
        # Set 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 1. Title Slide
        slide_layout = prs.slide_layouts[0] 
        slide = prs.slides.add_slide(slide_layout)
        
        # Apply Theme Colors to Title Slide (Custom Background)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self.colors["primary"])
        

        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = self._clean_markdown(title, strip_bold=True)
        subtitle_shape.text = self._clean_markdown(subtitle or f"{author} | {datetime.now().strftime('%Y-%m-%d')}", strip_bold=True)
        
        # White text on dark primary bg
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.name = self.CJK_FONTS[0]
            
        for paragraph in subtitle_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(200, 200, 200)
            paragraph.font.size = Pt(18)

        total_slides = len(outline)
        
        # 2. Iterate Outline
        for idx, section in enumerate(outline):
            s_title = section.get("title", "No Title")
            s_type = section.get("type", "detail")
            s_points = section.get("points", [])
            s_content = section.get("content", [])
            
            # Progress Callback
            if progress_callback:
                try:
                    progress_callback(idx + 1, total_slides, s_title, _type_name_map.get(s_type, s_type))
                except: pass
            
            # --- AI Image Generation Logic ---
            generated_img_path = None
            if image_gen and (s_type in ["image_full", "highlight", "content_image"] or (s_type == "detail" and idx % 3 == 0)):
                 # Generate image for specific types or every 3rd detail slide
                 img_prompt = section.get("image_prompt")
                 if not img_prompt:
                     # Auto-construct prompt from title and points
                     context_str = " ".join(s_points[:2])
                     img_prompt = f"Professional business illustration for presentation slide about '{s_title}'. Context: {context_str}. Style: {self.theme} style, clean, vector art or photorealistic."
                 
                 # Define distinct filenames
                 safe_name = re.sub(r'[\\/*?:"<>|]', "", s_title)[:20].strip()
                 img_filename = f"slide_{idx}_{safe_name}.png"
                 local_img_path = os.path.join(os.path.dirname(output_path), "images", img_filename)
                 os.makedirs(os.path.dirname(local_img_path), exist_ok=True)
                 
                 # Call API (Blocking for now, but per-slide ensures progress updates)
                 logger.info(f"[PPT] 🎨 Generating image for slide {idx+1}: {s_title}")
                 if image_gen.generate_image(img_prompt, local_img_path, aspect_ratio="16:9"):
                     generated_img_path = local_img_path
            
            # Create Slide based on type
            if s_type == "divider":
                self._create_divider_slide(prs, s_title, section.get("description", ""))
            
            elif s_type == "overview":
                self._create_overview_slide(prs, s_title, section.get("subsections", []))
            
            elif s_type == "comparison":
                self._create_comparison_slide(prs, s_title, section.get("left", {}), section.get("right", {}))
                
            elif s_type == "highlight":
                self._create_highlight_slide(prs, s_title, s_points, bg_image=generated_img_path)
            
            elif s_type == "image_full" and generated_img_path:
                 # Full screen image slide
                 slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
                 slide.shapes.add_picture(generated_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                 # Add overlay text

                 textbox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(10), Inches(1.5))
                 tf = textbox.text_frame
                 p = tf.add_paragraph()
                 p.text = self._clean_markdown(s_title, strip_bold=True)
                 p.font.size = Pt(36)
                 p.font.color.rgb = RGBColor(255, 255, 255)
                 p.font.bold = True
                 # Add Shadow
                 # (Shadow API in python-pptx is complex, skipping for simplicity)
                 
            else:
                # Default Detail Layout
                # If we have an generated image, use "Picture with Caption" layout logic manually
                if generated_img_path:
                    self._create_picture_slide(prs, s_title, s_points, generated_img_path)
                else:
                    self._create_detail_slide(prs, s_title, s_points)

        # Save
        prs.save(output_path)
        return {"output_path": output_path, "slide_count": len(prs.slides)}



    def _create_picture_slide(self, prs, title, points, image_path):
        """Creates a slide with image on right, text on left"""
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
        
        # 1. Title
        title_shape = slide.shapes.title
        title_shape.text = self._clean_markdown(title, strip_bold=True)
        
        # Style Title using same logic as Detail Slide
        for paragraph in title_shape.text_frame.paragraphs:
             paragraph.font.name = self.CJK_FONTS[0]
             # title color usually comes from theme master
        
        # 2. Image (Right Half, taking up ~45% width)
        left = Inches(7.5)
        top = Inches(2.0)
        height = Inches(4.5)
        
        try:
            slide.shapes.add_picture(image_path, left, top, height=height)
        except Exception:
            pass # Image load fail
            
        # 3. Content (Left Half)
        body_shape = slide.placeholders[1]
        body_shape.width = Inches(6.5) # Limit to left side
        tf = body_shape.text_frame
        tf.clear()
        
        for point in points:
            p = tf.add_paragraph()
            p.text = self._clean_markdown(str(point), strip_bold=True)


    


    def _create_detail_slide(self, prs, title, points):
        """Standard detail slide with bullet points"""
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = self._clean_markdown(title, strip_bold=True)
        
        # Color title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = self.CJK_FONTS[0]
            if self.theme and self.ppt_theme:
                 paragraph.font.color.rgb = RGBColor(*self.ppt_theme.primary_color)
            else:
                 paragraph.font.color.rgb = RGBColor(*self.colors.get("primary", (0,0,0)))

        # Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for point in points:
            p = tf.add_paragraph()
            p.text = self._clean_markdown(str(point), strip_bold=True)
            p.font.size = Pt(24)
            p.space_before = Pt(14)
            p.level = 0
            # Set bullet color if possible (hard with python-pptx, skipping custom bullet)


    def _create_overview_slide(self, prs, title, subsections):
        """Overview slide with multiple sections/cards"""
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
        
        title_shape = slide.shapes.title
        title_shape.text = self._clean_markdown(title, strip_bold=True)
        
        # Create 2-3 columns depending on subsection count
        if not subsections:
            return
            
        count = len(subsections)
        width = (prs.slide_width - Inches(2)) / count
        for i, sub in enumerate(subsections):
            left = Inches(1) + (width * i)
            top = Inches(2.5)
            height = Inches(4)
            
            # Draw a box/shape
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width - Inches(0.2), height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*self.colors.get("light_bg", (240,240,240)))
            shape.line.color.rgb = RGBColor(*self.colors.get("primary", (0,0,0)))
            
            # Add text inside
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = self._clean_markdown(sub.get("title", f"Section {i+1}"), strip_bold=True)
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(*self.colors.get("primary", (0,0,0)))
            
            # Content
            details = sub.get("points", [])
            for d in details[:3]:
                p = tf.add_paragraph()
                p.text = f"• {self._clean_markdown(str(d), strip_bold=True)}"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(*self.colors.get("text", (50,50,50)))


    def _create_highlight_slide(self, prs, title, points, bg_image=None):
        """Highlight slide with big numbers or key takeaway"""
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
        
        if bg_image:
             try:
                 slide.shapes.add_picture(bg_image, 0, 0, width=prs.slide_width, height=prs.slide_height)
             except: pass
        else:
             # Use solid color background
             background = slide.background
             fill = background.fill
             fill.solid()
             fill.fore_color.rgb = RGBColor(*self.colors.get("primary", (0,0,0)))

        # Add Title (White usually)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(1.5))
        tf = tb.text_frame
        p = tf.add_paragraph()
        p.text = self._clean_markdown(title, strip_bold=True)
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Points
        top_offset = Inches(3)
        for point in points:
            tb = slide.shapes.add_textbox(Inches(2), top_offset, Inches(9), Inches(1))
            tf = tb.text_frame
            p = tf.add_paragraph()
            p.text = self._clean_markdown(str(point), strip_bold=True)
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            top_offset += Inches(1.2)


    def _create_comparison_slide(self, prs, title, left, right):
        """Comparison side by side"""
        slide = prs.slides.add_slide(prs.slide_layouts[1]) 
        slide.shapes.title.text = self._clean_markdown(title, strip_bold=True)
        
        # Left Box
        left_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5.5), Inches(5))
        tf_l = left_box.text_frame
        tf_l.text = self._clean_markdown(left.get("title", "Left Side"), strip_bold=True)
        tf_l.paragraphs[0].font.bold = True
        for p_text in left.get("points", []):
            p = tf_l.add_paragraph()
            p.text = f"• {self._clean_markdown(str(p_text), strip_bold=True)}"
            
        # Right Box
        right_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5.5), Inches(5))
        tf_r = right_box.text_frame
        tf_r.text = self._clean_markdown(right.get("title", "Right Side"), strip_bold=True)
        tf_r.paragraphs[0].font.bold = True
        for p_text in right.get("points", []):
            p = tf_r.add_paragraph()
            p.text = f"• {self._clean_markdown(str(p_text), strip_bold=True)}"
            

    def _create_divider_slide(self, prs, title, subtitle):
        """Section divider"""
        self._create_highlight_slide(prs, title, [subtitle])

        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)   # 16:9 宽屏
            prs.slide_height = Inches(7.5)
            
            total_slides = len(outline)
            
            # 1. 封面页
            _progress(0, total_slides, title, "封面")
            self._add_title_slide(prs, title, subtitle, author)
            
            # 2. 目录页 (≥3 页内容才需要，过滤掉过渡页)
            content_slides = [s for s in outline if s.get("type", "detail") != "divider"]
            if len(content_slides) >= 3:
                self._add_agenda_slide(prs, content_slides)
            
            # 3. 内容页 — 根据 slide type 分发到不同渲染器
            divider_count = 0
            for idx, section in enumerate(outline):
                slide_type = section.get("type", "detail")
                section_title = section.get("title", "")
                _progress(idx + 1, total_slides, section_title, slide_type)
                
                if slide_type == "divider":
                    divider_count += 1
                    self._add_section_divider_slide(prs, section, divider_count)
                elif slide_type == "overview":
                    self._add_overview_slide(prs, section, idx + 1, len(outline))
                elif slide_type == "highlight":
                    self._add_highlight_slide(prs, section, idx + 1, len(outline))
                elif slide_type == "comparison":
                    self._add_comparison_slide(prs, section, idx + 1, len(outline))
                else:  # "detail" 或未指定 → 默认详细内容页
                    self._add_content_slide(prs, section, idx + 1, len(outline))
            
            # 4. 结束页
            _progress(total_slides, total_slides, "结束页", "ending")
            self._add_ending_slide(prs, title)
            
            # 5. 给所有非封面和结束页添加页码
            total = len(prs.slides)
            for i, slide in enumerate(prs.slides):
                if 0 < i < total - 1:
                    self._add_slide_number(prs, slide, i, total - 2)
            
            os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
            prs.save(output_path)
            
            return {
                "success": True,
                "output_path": output_path,
                "slide_count": len(prs.slides),
                "message": f"成功生成 {len(prs.slides)} 页PPT"
            }
        except ImportError:
            return {"success": False, "error": "需要安装 python-pptx: pip install python-pptx"}
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    # ─── 辅助方法 ───────────────────────────────────
    
    def _rgb(self, key):
        from pptx.dml.color import RGBColor
        return RGBColor(*self.colors[key])
    
    def _set_font(self, run, size, bold=False, italic=False, color_key="text"):
        from pptx.util import Pt
        from pptx.oxml.ns import qn
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = self._rgb(color_key)
        run.font.name = self.LATIN_FONTS[0]
        rpr = run._r.get_or_add_rPr()
        ea = rpr.find(qn('a:ea'))
        if ea is None:
            import lxml.etree as etree
            ea = etree.SubElement(rpr, qn('a:ea'))
        ea.set('typeface', self.CJK_FONTS[0])
    
    def _add_rect(self, slide, left, top, width, height, color_key="primary", alpha=None):
        from pptx.util import Emu
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(color_key)
        shape.line.fill.background()
        if alpha is not None:
            from pptx.oxml.ns import qn
            solid = shape.fill._fill
            srgb = solid.find('.//' + qn('a:srgbClr'))
            if srgb is not None:
                import lxml.etree as etree
                alpha_el = etree.SubElement(srgb, qn('a:alpha'))
                alpha_el.set('val', str(int(alpha * 1000)))
        return shape
    
    def _add_line(self, slide, x1, y1, x2, y2, color_key="accent", width_pt=3):
        from pptx.util import Pt
        connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT
        connector.line.color.rgb = self._rgb(color_key)
        connector.line.width = Pt(width_pt)
        return connector
    
    def _add_slide_number(self, prs, slide, current, total):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        w = prs.slide_width
        box = slide.shapes.add_textbox(w - Inches(1.5), prs.slide_height - Inches(0.5), Inches(1.2), Inches(0.35))
        tf = box.text_frame
        tf.text = f"{current}/{total}"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        self._set_font(p.runs[0], 10, color_key="subtitle")
    
    # ─── 封面页 ─────────────────────────────────────
    
    def _add_title_slide(self, prs, title, subtitle, author):
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W = prs.slide_width
        H = prs.slide_height
        
        # 左侧色带 (宽 0.4 英寸)
        self._add_rect(slide, 0, 0, Inches(0.45), H, "primary")
        
        # 底部装饰条
        self._add_rect(slide, 0, H - Inches(0.9), W, Inches(0.9), "primary")
        
        # 装饰性 accent 线
        self._add_line(slide, Inches(2), Inches(3.8), Inches(6.5), Inches(3.8), "accent", 4)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(9), Inches(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        self._set_font(p.runs[0], 44, bold=True, color_key="primary")
        
        # 副标题
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.alignment = PP_ALIGN.LEFT
            p2.space_before = Pt(12)
            self._set_font(p2.runs[0], 22, color_key="subtitle")
        
        # 作者 + 日期 (底部色带上)
        footer_text = f"{author}  |  {datetime.now().strftime('%Y年%m月%d日')}"
        footer_box = slide.shapes.add_textbox(Inches(2), H - Inches(0.75), Inches(8), Inches(0.5))
        fp = footer_box.text_frame.paragraphs[0]
        fp.text = footer_text
        fp.alignment = PP_ALIGN.LEFT
        from pptx.dml.color import RGBColor
        run = fp.runs[0]
        self._set_font(run, 14, color_key="text")
        run.font.color.rgb = RGBColor(255, 255, 255)
    
    # ─── 目录页 ─────────────────────────────────────
    
    def _add_agenda_slide(self, prs, outline):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W = prs.slide_width
        
        # 顶部色带
        self._add_rect(slide, 0, 0, W, Inches(1.3), "primary")
        
        # 标题 (白色，在色带上)
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(8), Inches(0.8))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = "目 录"
        from pptx.dml.color import RGBColor
        run = tp.runs[0]
        self._set_font(run, 36, bold=True, color_key="text")
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        # 目录项 — 两列布局
        items = outline
        half = (len(items) + 1) // 2
        col_width = Inches(5.5)
        
        for col in range(2):
            start_idx = col * half
            end_idx = min(start_idx + half, len(items))
            x = Inches(1) + col * col_width
            y_start = Inches(2)
            
            for i in range(start_idx, end_idx):
                y = y_start + (i - start_idx) * Inches(0.75)
                
                # 序号圆圈
                num_box = slide.shapes.add_textbox(x, y, Inches(0.5), Inches(0.5))
                np = num_box.text_frame.paragraphs[0]
                np.text = f"{i + 1:02d}"
                np.alignment = PP_ALIGN.CENTER
                self._set_font(np.runs[0], 16, bold=True, color_key="accent")
                
                # 标题文字
                item_box = slide.shapes.add_textbox(x + Inches(0.7), y + Inches(0.05), Inches(4.5), Inches(0.45))
                ip = item_box.text_frame.paragraphs[0]
                ip.text = items[i].get('title', '')
                self._set_font(ip.runs[0], 18, color_key="text")
    
    # ─── 内容页 ─────────────────────────────────────
    
    def _add_content_slide(self, prs, section, page_num, total_pages):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W = prs.slide_width
        H = prs.slide_height
        
        # 左侧细色带 (装饰)
        self._add_rect(slide, 0, 0, Inches(0.15), H, "primary")
        
        # 顶部区域：章节标题
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        section_title = section.get('title', '')
        p.text = section_title
        p.alignment = PP_ALIGN.LEFT
        self._set_font(p.runs[0], 30, bold=True, color_key="primary")
        
        # 标题下方 accent 分隔线
        self._add_line(slide, Inches(0.6), Inches(1.35), Inches(4), Inches(1.35), "accent", 3)
        
        # 内容区域
        points = section.get('points', []) or section.get('content', [])
        image_path = section.get('image', None)
        has_image = image_path and os.path.exists(str(image_path))
        
        if has_image:
            content_width = Inches(6)
        else:
            content_width = Inches(11.5)
        
        # 根据要点数量调整字号
        point_count = len(points)
        if point_count <= 3:
            font_size = 24  # Size up for fewer points
            spacing = Pt(20)
        elif point_count <= 5:
            font_size = 20
            spacing = Pt(14)
        else:
            font_size = 18
            spacing = Pt(10)
        
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.7), content_width, Inches(5.3)
        )
        ctf = content_box.text_frame
        ctf.word_wrap = True
        ctf.auto_size = None  # We manually manage font size
        
        for i, point in enumerate(points):
            if i == 0:
                p = ctf.paragraphs[0]
            else:
                p = ctf.add_paragraph()
            
            # Use native bullet level (for indent) but manual char for safety
            p.level = 0
            
            # Add bullet manually to ensure visibility
            bullet_run = p.add_run()
            bullet_run.text = "•  "
            self._set_font(bullet_run, font_size, color_key="accent")

            # 先清理 markdown 残留
            point = self._clean_markdown(point)
            
            # 解析粗体...
            import re as _re
            parts = _re.split(r'(\*\*.*?\*\*)', point)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = p.add_run()
                    run.text = part[2:-2]
                    self._set_font(run, font_size, bold=True, color_key="text")
                elif part:
                    run = p.add_run()
                    # 允许较长的内容，python-pptx 会自动换行
                    run.text = part[:200] + ('...' if len(part) > 200 else '')
                    self._set_font(run, font_size, color_key="text")
            
            p.space_before = spacing
            p.space_after = Pt(4)
        
        # 图片 (右侧) - Fix: Preserve Aspect Ratio
        if has_image:
            try:
                # Add picture without specifying dimensions first to get natural size
                pic = slide.shapes.add_picture(str(image_path), Inches(7.5), Inches(1.8))
                
                # Define constraints
                max_width = Inches(5.0)
                max_height = Inches(4.5)
                
                # Calculate scaling factor to fit within box while preserving aspect ratio
                width_ratio = max_width / pic.width
                height_ratio = max_height / pic.height
                scale = min(width_ratio, height_ratio)
                
                pic.width = int(pic.width * scale)
                pic.height = int(pic.height * scale)
                
                # Optional: Center image in the reserved area? 
                # Currently top-aligned at Inches(1.8), Left-aligned at Inches(7.5)
                
            except Exception as e:
                logger.info(f"[PPT] 添加图片失败: {e}")
    
    # ─── 概览页（多主题并列） ─────────────────────────
    
    def _add_overview_slide(self, prs, section, page_num, total_pages):
        """概览页 - 2~4 个小主题卡片并列展示（适合简要 / 速览内容）"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W = prs.slide_width
        H = prs.slide_height
        
        # 左侧装饰条（与内容页一致）
        self._add_rect(slide, 0, 0, Inches(0.15), H, "primary")
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(0.7))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = section.get("title", "概览")
        tp.alignment = PP_ALIGN.LEFT
        self._set_font(tp.runs[0], 28, bold=True, color_key="primary")
        self._add_line(slide, Inches(0.6), Inches(1.15), Inches(4), Inches(1.15), "accent", 3)
        
        # 获取子主题
        subsections = section.get("subsections", [])
        if not subsections:
            # 没有子主题 → 尝试从 points 自动分组为子主题
            points = section.get("points", section.get("content", []))
            if points:
                # 智能分组: 每 2-3 个 point 归为一组，取第 1 个的前几个字作为小标题
                import re as _re
                group_size = 3 if len(points) <= 9 else 2
                auto_subs = []
                for g_start in range(0, len(points), group_size):
                    group_pts = points[g_start:g_start + group_size]
                    # 用第一个 point 的核心词做小标题
                    first = self._clean_markdown(group_pts[0])
                    # 取前20个字符，如果有 "：" 或 ":" 取冒号前面
                    if "：" in first:
                        label = first.split("：")[0].strip()
                    elif ":" in first:
                        label = first.split(":")[0].strip()
                    elif "—" in first:
                        label = first.split("—")[0].strip()
                    elif "，" in first:
                        label = first.split("，")[0].strip()
                    else:
                        label = first[:20] + ('...' if len(first) > 20 else '')
                    auto_subs.append({
                        "subtitle": label,
                        "label": label,
                        "points": [self._clean_markdown(p) for p in group_pts]
                    })
                subsections = auto_subs
            else:
                return  # 无内容
        
        n = len(subsections)
        cols = min(n, 3) if n <= 3 else 2
        rows = 1 if n <= 3 else 2
        
        usable_w = 12.0   # 可用宽度 (英寸)
        usable_h = 5.2    # 可用高度
        col_w = usable_w / cols
        row_h = usable_h / rows
        x0, y0 = 0.7, 1.5
        
        for i, sub in enumerate(subsections[:cols * rows]):
            ci = i % cols
            ri = i // cols
            x = x0 + ci * col_w
            y = y0 + ri * row_h
            
            # 卡片背景
            margin = 0.12
            self._add_rect(
                slide, Inches(x + margin), Inches(y + margin),
                Inches(col_w - 2 * margin), Inches(row_h - 2 * margin), "light_bg"
            )
            
            # 子标题
            sub_box = slide.shapes.add_textbox(
                Inches(x + 0.35), Inches(y + 0.25),
                Inches(col_w - 0.7), Inches(0.45)
            )
            sp = sub_box.text_frame.paragraphs[0]
            sp.text = self._clean_markdown(sub.get("subtitle", sub.get("label", "")), strip_bold=True)
            sp.alignment = PP_ALIGN.LEFT
            self._set_font(sp.runs[0], 18, bold=True, color_key="accent")
            
            # 子标题下小分隔线
            self._add_line(
                slide, Inches(x + 0.35), Inches(y + 0.78),
                Inches(x + 1.6), Inches(y + 0.78), "accent", 2
            )
            
            # 子要点
            pts_box = slide.shapes.add_textbox(
                Inches(x + 0.35), Inches(y + 0.9),
                Inches(col_w - 0.7), Inches(row_h - 1.3)
            )
            ptf = pts_box.text_frame
            ptf.word_wrap = True
            ptf.auto_size = None
            for j, pt in enumerate(sub.get("points", [])[:6]):
                pt = self._clean_markdown(pt, strip_bold=True)
                if j == 0:
                    p = ptf.paragraphs[0]
                else:
                    p = ptf.add_paragraph()
                
                # Use native bullet
                p.level = 0
                
                bullet = p.add_run()
                bullet.text = "• "
                self._set_font(bullet, 12, color_key="accent")
                
                run = p.add_run()
                display = pt[:100] + ('...' if len(pt) > 100 else '')
                run.text = display
                self._set_font(run, 13, color_key="text")
                p.space_before = Pt(3)
    
    def _render_bullet_area(self, slide, points, left, top, width, height):
        """在指定区域渲染 bullet 列表（内部辅助方法）"""
        from pptx.util import Pt
        import re as _re
        
        count = len(points)
        font_size = 20 if count <= 3 else (17 if count <= 5 else 15)
        spacing = Pt(16) if count <= 3 else (Pt(10) if count <= 5 else Pt(6))
        
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        
        for i, point in enumerate(points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            
            # Use native bullet
            p.level = 0
            
            # Manual bullet
            bullet = p.add_run()
            bullet.text = "• "
            self._set_font(bullet, 12, color_key="accent")
            
            # 清理 markdown 残留
            point = self._clean_markdown(point)
            
            parts = _re.split(r'(\*\*.*?\*\*)', point)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = p.add_run()
                    run.text = part[2:-2]
                    self._set_font(run, font_size, bold=True, color_key="text")
                elif part:
                    run = p.add_run()
                    # 允许较长的内容，python-pptx 会自动换行
                    run.text = part[:200] + ('...' if len(part) > 200 else '')
                    self._set_font(run, font_size, color_key="text")
            p.space_before = spacing
            p.space_after = Pt(4)
    
    # ─── 亮点数据页 ─────────────────────────────────
    
    def _add_highlight_slide(self, prs, section, page_num, total_pages):
        """亮点数据页 - 大字展示关键数字 / 成果（最多 3 张数据卡片）"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W = prs.slide_width
        H = prs.slide_height
        
        # 左侧装饰条
        self._add_rect(slide, 0, 0, Inches(0.15), H, "primary")
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(0.7))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = section.get("title", "")
        tp.alignment = PP_ALIGN.LEFT
        self._set_font(tp.runs[0], 28, bold=True, color_key="primary")
        self._add_line(slide, Inches(0.6), Inches(1.15), Inches(4), Inches(1.15), "accent", 3)
        
        # 从 points 中解析 "数值 | 说明" 格式
        highlights = []
        for pt in section.get("points", section.get("content", [])):
            pt = self._clean_markdown(pt, strip_bold=True)
            if "|" in pt:
                parts = pt.split("|", 1)
                highlights.append({"value": parts[0].strip(), "label": parts[1].strip()})
            elif "：" in pt:
                parts = pt.split("：", 1)
                highlights.append({"value": parts[0].strip(), "label": parts[1].strip()})
            else:
                highlights.append({"value": "", "label": pt.strip()})
        
        if not highlights:
            # 无数据 → 回退到普通内容页
            self._add_content_slide(prs, section, page_num, total_pages)
            return
        
        n = min(len(highlights), 4)
        card_w = 3.2
        gap = 0.4
        total_w = n * card_w + (n - 1) * gap
        x_start = (13.333 - total_w) / 2
        card_h = 3.8
        y_card = 2.0
        
        for i, hl in enumerate(highlights[:n]):
            x = x_start + i * (card_w + gap)
            
            # 卡片背景
            self._add_rect(slide, Inches(x), Inches(y_card), Inches(card_w), Inches(card_h), "light_bg")
            
            # 顶部 accent 色条
            self._add_rect(slide, Inches(x), Inches(y_card), Inches(card_w), Inches(0.08), "accent")
            
            # 大数值
            if hl["value"]:
                val_box = slide.shapes.add_textbox(
                    Inches(x + 0.2), Inches(y_card + 0.5),
                    Inches(card_w - 0.4), Inches(1.6)
                )
                vtf = val_box.text_frame
                vtf.word_wrap = True
                vp = vtf.paragraphs[0]
                vp.text = hl["value"]
                vp.alignment = PP_ALIGN.CENTER
                self._set_font(vp.runs[0], 42, bold=True, color_key="accent")
                
                # 说明文字 (standard)
                lbl_box = slide.shapes.add_textbox(
                    Inches(x + 0.25), Inches(y_card + 2.2),
                    Inches(card_w - 0.5), Inches(1.3)
                )
                label_size = 15
                label_y_offset = 0 # Included in box position
            else:
                # No value -> Center the label and make it bigger
                lbl_box = slide.shapes.add_textbox(
                    Inches(x + 0.25), Inches(y_card + 1.2),
                    Inches(card_w - 0.5), Inches(2.3)
                )
                label_size = 24
            
            ltf = lbl_box.text_frame
            ltf.word_wrap = True
            lp = ltf.paragraphs[0]
            lp.text = hl["label"]
            lp.alignment = PP_ALIGN.CENTER
            self._set_font(lp.runs[0], label_size, color_key="text" if hl["value"] else "primary")
    
    # ─── 章节过渡页 ─────────────────────────────────
    
    def _add_section_divider_slide(self, prs, section, part_number=1):
        """章节过渡页 - 左文右色块，引入新的大章节"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W = prs.slide_width
        H = prs.slide_height
        
        # 右侧色块（占约 45%）
        self._add_rect(slide, Inches(7.5), 0, Inches(5.833), H, "primary")
        
        # 左侧区域: 章节编号
        num_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(5), Inches(1.2))
        np_ = num_box.text_frame.paragraphs[0]
        np_.text = f"PART {part_number:02d}"
        np_.alignment = PP_ALIGN.LEFT
        self._set_font(np_.runs[0], 20, bold=False, color_key="accent")
        
        # accent 分隔线
        self._add_line(slide, Inches(1.2), Inches(3.2), Inches(6), Inches(3.2), "accent", 4)
        
        # 章节标题
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(5.5), Inches(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        tp = tf.paragraphs[0]
        tp.text = section.get("title", "")
        tp.alignment = PP_ALIGN.LEFT
        self._set_font(tp.runs[0], 36, bold=True, color_key="primary")
        
        # 描述文字（如果有）
        desc = section.get("description", "")
        if desc:
            dp = tf.add_paragraph()
            dp.text = desc
            dp.space_before = Pt(16)
            dp.alignment = PP_ALIGN.LEFT
            self._set_font(dp.runs[0], 16, color_key="subtitle")
    
    # ─── 对比页 ─────────────────────────────────────
    
    def _add_comparison_slide(self, prs, section, page_num, total_pages):
        """对比页 - 左右两栏对比展示"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W = prs.slide_width
        H = prs.slide_height
        
        # 左侧装饰条
        self._add_rect(slide, 0, 0, Inches(0.15), H, "primary")
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(0.7))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = section.get("title", "对比")
        tp.alignment = PP_ALIGN.LEFT
        self._set_font(tp.runs[0], 28, bold=True, color_key="primary")
        self._add_line(slide, Inches(0.6), Inches(1.15), Inches(4), Inches(1.15), "accent", 3)
        
        # 获取左右数据
        left_data = section.get("left", {})
        right_data = section.get("right", {})
        
        # 如果没有 left/right，尝试从 subsections 取前两个
        if not left_data and not right_data:
            subs = section.get("subsections", [])
            if len(subs) >= 2:
                left_data = subs[0]
                right_data = subs[1]
            elif len(subs) == 1:
                left_data = subs[0]
                right_data = {"label": "", "points": []}
            else:
                # 回退到普通内容页
                points = section.get("points", section.get("content", []))
                if points:
                    self._render_bullet_area(slide, points, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.3))
                return
        
        col_w = 5.5
        left_x = 0.8
        right_x = 7.0
        y_start = 1.7
        
        # 中间分隔线
        center_x = 6.6
        self._add_line(slide, Inches(center_x), Inches(y_start),
                       Inches(center_x), Inches(6.8), "light_bg", 2)
        
        for col_data, x in [(left_data, left_x), (right_data, right_x)]:
            if not col_data:
                continue
            
            # 列标题
            header_box = slide.shapes.add_textbox(
                Inches(x), Inches(y_start), Inches(col_w), Inches(0.55)
            )
            hp = header_box.text_frame.paragraphs[0]
            hp.text = self._clean_markdown(col_data.get("label", col_data.get("subtitle", "")), strip_bold=True)
            hp.alignment = PP_ALIGN.LEFT
            self._set_font(hp.runs[0], 22, bold=True, color_key="accent")
            
            # 列要点
            pts_box = slide.shapes.add_textbox(
                Inches(x), Inches(y_start + 0.7),
                Inches(col_w - 0.3), Inches(4.5)
            )
            ptf = pts_box.text_frame
            ptf.word_wrap = True
            ptf.auto_size = None
            for j, pt in enumerate(col_data.get("points", [])[:6]):
                pt = self._clean_markdown(pt, strip_bold=True)
                if j == 0:
                    p = ptf.paragraphs[0]
                else:
                    p = ptf.add_paragraph()
                
                # Use native bullet
                p.level = 0
                
                bullet = p.add_run()
                bullet.text = "• "
                self._set_font(bullet, 12, color_key="accent")
                
                run = p.add_run()
                run.text = pt[:150] + ('...' if len(pt) > 150 else '')
                self._set_font(run, 16, color_key="text")
                p.space_before = Pt(10)
                p.space_after = Pt(4)
    
    # ─── 结束页 ─────────────────────────────────────
    
    def _add_ending_slide(self, prs, title):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W = prs.slide_width
        H = prs.slide_height
        
        # 全屏背景色
        self._add_rect(slide, 0, 0, W, H, "primary")
        
        # 装饰线
        self._add_line(slide, Inches(4), Inches(3.3), Inches(9.3), Inches(3.3), "accent", 4)
        
        # 谢谢观看
        end_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(1.2))
        tf = end_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "谢谢观看"
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        self._set_font(run, 52, bold=True, color_key="text")
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Thank You
        p2 = tf.add_paragraph()
        p2.text = "THANK YOU"
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
        run2 = p2.runs[0]
        self._set_font(run2, 24, color_key="subtitle")
        run2.font.color.rgb = RGBColor(200, 200, 220)
        
        # 底部信息
        info_box = slide.shapes.add_textbox(Inches(1), H - Inches(1.5), Inches(11.3), Inches(0.6))
        ip = info_box.text_frame.paragraphs[0]
        ip.text = f"Generated by Koto AI  •  {datetime.now().strftime('%Y-%m-%d')}"
        ip.alignment = PP_ALIGN.CENTER
        run3 = ip.runs[0]
        self._set_font(run3, 12, color_key="subtitle")
        run3.font.color.rgb = RGBColor(180, 180, 200)
    
    def generate_from_text(
        self,
        content: str,
        output_path: str,
        title: str = "演示文稿",
        ai_model: Optional[Any] = None
    ) -> Dict[str, Any]:
        """
        从长文本生成PPT（使用AI提取结构）
        
        Args:
            content: 输入文本内容
            output_path: 输出路径
            title: 标题
            ai_model: AI模型（用于提取结构）
        
        Returns:
            生成结果
        """
        # 简单的启发式方法提取结构
        lines = content.strip().split('\n')
        outline = []
        current_section = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 识别标题（以#开头或全大写）
            if line.startswith('#') or (len(line) < 50 and line.isupper()):
                if current_section:
                    outline.append(current_section)
                current_section = {
                    'title': line.lstrip('#').strip(),
                    'points': []
                }
            elif current_section:
                # 识别要点（以-或数字开头）
                if line.startswith('-') or line.startswith('•') or (len(line) > 0 and line[0].isdigit()):
                    current_section['points'].append(line.lstrip('-•0123456789. ').strip())
                else:
                    # 普通段落，按句子分割
                    if len(current_section['points']) < 5:  # 限制每页要点数
                        current_section['points'].append(line[:100])
        
        if current_section:
            outline.append(current_section)
        
        # 如果提取失败，创建默认结构
        if not outline:
            outline = [{
                'title': '内容',
                'points': [line[:100] for line in lines[:5] if line.strip()]
            }]
        
        return self.generate_from_outline(title, outline, output_path)
    
    def add_image_to_slide(self, prs, slide_index: int, image_path: str, position: str = "right") -> bool:
        """向指定幻灯片添加图片"""
        try:
            from pptx.util import Inches
            
            if slide_index >= len(prs.slides):
                return False
            
            slide = prs.slides[slide_index]
            shapes = slide.shapes
            
            # 根据位置确定坐标
            if position == "right":
                left, top = Inches(5.5), Inches(2)
                width, height = Inches(4), Inches(4.5)
            elif position == "center":
                left, top = Inches(3), Inches(2.5)
                width, height = Inches(4), Inches(4)
            elif position == "full":
                left, top = Inches(1), Inches(1.5)
                width, height = Inches(8), Inches(5.5)
            else:
                left, top = Inches(1), Inches(2)
                width, height = Inches(3.5), Inches(4)
            
            shapes.add_picture(image_path, left, top, width=width, height=height)
            return True
        except Exception as e:
            logger.info(f"[PPT] 添加图片到幻灯片 {slide_index} 失败: {e}")
            return False


class EnhancedPPTGenerator:
    """增强版PPT生成器 - 支持多模型协作"""
    
    def __init__(self, theme: str = "business"):
        self.generator = PPTGenerator(theme=theme)
        self.theme = theme
    
    async def generate_with_multimodal(
        self,
        title: str,
        user_request: str,
        output_path: str,
        search_results: Optional[List[Dict]] = None,
        images: Optional[List[str]] = None,
        ai_client = None,
        quality_feedback: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        多模型协作生成PPT
        
        工作流:
        1. 基于搜索结果生成结构化大纲
        2. 为每个章节生成详细内容
        3. 匹配图片到合适的幻灯片
        4. 优化排版和设计
        
        Args:
            title: PPT标题
            user_request: 用户原始请求
            output_path: 输出路径
            search_results: 搜索结果列表
            images: 生成的图片路径列表
            ai_client: AI客户端（用于内容生成）
        """
        from google.genai import types
        
        # 1. 生成智能大纲
        outline_prompt = self._build_outline_prompt(
            title,
            user_request,
            search_results,
            quality_feedback
        )
        
        if ai_client:
            try:
                response = ai_client.models.generate_content(
                    model="gemini-3-flash-preview",
                    contents=outline_prompt,
                    config=types.GenerateContentConfig(
                        system_instruction="你是专业的PPT内容策划师，擅长结构化内容组织。",
                        temperature=0.7,
                        max_output_tokens=3000,
                    )
                )
                outline_text = response.text
            except Exception as e:
                logger.info(f"[PPT] 大纲生成失败: {e}")
                outline_text = self._generate_fallback_outline(title, user_request)
        else:
            outline_text = self._generate_fallback_outline(title, user_request)
        
        # 2. 解析大纲
        outline = self._parse_enhanced_outline(outline_text)
        
        # 3. 匹配图片到幻灯片
        if images:
            outline = self._match_images_to_slides(outline, images)
        
        # 4. 生成PPT
        result = self.generator.generate_from_outline(
            title=title,
            outline=outline,
            output_path=output_path,
            subtitle=self._extract_subtitle(user_request)
        )
        
        return result
    
    def _build_outline_prompt(
        self,
        title: str,
        user_request: str,
        search_results: Optional[List[Dict]],
        quality_feedback: Optional[Dict[str, Any]]
    ) -> str:
        """构建大纲生成提示词"""
        prompt = f"""请为演示文稿 "{title}" 生成详细的结构化大纲。

用户需求: {user_request}

"""
        
        if search_results:
            prompt += "\n参考信息:\n"
            for i, result in enumerate(search_results[:5], 1):
                title_text = result.get('title', '')
                snippet = result.get('snippet', '')
                prompt += f"{i}. {title_text[:100]}\n   {snippet[:200]}\n\n"

        if quality_feedback:
            issues = quality_feedback.get("issues", [])
            suggestions = quality_feedback.get("suggestions", [])
            prompt += "\n质量反馈（请重点修正以下问题）:\n"
            for issue in issues[:6]:
                prompt += f"- 问题: {issue}\n"
            for suggestion in suggestions[:6]:
                prompt += f"- 改进: {suggestion}\n"
        
        prompt += """\n请生成包含以下内容的大纲（Markdown格式）：

# 演示文稿标题

## 章节1：引言/背景
- 要点1：核心问题或背景
- 要点2：为什么重要
- 要点3：本次演示的目标

## 章节2：主要内容（可以是多个章节）
- 要点1：关键信息
- 要点2：详细说明
- 要点3：数据或案例
- 要点4：深入分析

## 章节3：总结/展望
- 要点1：核心观点总结
- 要点2：行动建议
- 要点3：未来展望

要求：
1. 结构清晰，逻辑连贯
2. 每个章节3-5个要点
3. 内容充实，避免空洞
4. 适合口头演讲
5. 总共4-7个章节
"""
        return prompt
    
    def _generate_fallback_outline(self, title: str, user_request: str) -> str:
        """生成备用大纲"""
        return f"""# {title}

## 背景介绍
- 主题概述
- 重要性说明
- 目标受众

## 核心内容
- 关键要点1
- 关键要点2
- 关键要点3
- 详细分析

## 总结与展望
- 核心观点总结
- 实践建议
- 未来方向
"""
    
    def _parse_enhanced_outline(self, md_text: str) -> List[Dict]:
        """解析增强型大纲"""
        lines = md_text.split('\n')
        outline = []
        current_section = None
        
        for line in lines:
            line = line.rstrip()
            if line.startswith('## '):
                if current_section:
                    outline.append(current_section)
                current_section = {
                    'title': line[3:].strip(),
                    'points': [],
                    'content': [],
                    'image': None
                }
            elif line.startswith('- ') and current_section:
                current_section['points'].append(line[2:].strip())
                current_section['content'].append(line[2:].strip())
            elif line.startswith('• ') and current_section:
                current_section['points'].append(line[2:].strip())
                current_section['content'].append(line[2:].strip())
        
        if current_section:
            outline.append(current_section)
        
        return outline
    
    def _match_images_to_slides(self, outline: List[Dict], images: List[str]) -> List[Dict]:
        """智能匹配图片到幻灯片"""
        # 简单策略：为前N个有内容的章节分配图片
        image_idx = 0
        for section in outline:
            if image_idx < len(images) and len(section.get('points', [])) > 0:
                section['image'] = images[image_idx]
                image_idx += 1
                if image_idx >= len(images):
                    break
        return outline
    
    def _extract_subtitle(self, user_request: str) -> str:
        """从用户请求提取副标题"""
        # 简单提取：如果有年份或关键主题
        import re
        year_match = re.search(r'(20\d{2})年?', user_request)
        if year_match:
            return f"{year_match.group(1)}年度报告"
        return ""


# ================= 测试 =================

if __name__ == "__main__":
    generator = PPTGenerator(theme="business")
    
    # 示例大纲
    outline = [
        {
            "title": "项目背景",
            "points": [
                "市场需求日益增长",
                "技术发展为实现提供可能",
                "团队具备丰富经验",
                "时机已经成熟"
            ]
        },
        {
            "title": "解决方案",
            "points": [
                "采用先进的AI技术",
                "构建用户友好的界面",
                "确保系统稳定可靠",
                "提供7x24小时支持"
            ]
        },
        {
            "title": "预期成果",
            "points": [
                "提升工作效率50%",
                "降低运营成本30%",
                "改善用户体验",
                "增强市场竞争力"
            ]
        }
    ]
    
    result = generator.generate_from_outline(
        title="AI助手项目方案",
        subtitle="让工作更智能",
        outline=outline,
        output_path="workspace/documents/presentation.pptx"
    )
    
    logger.info(result)
