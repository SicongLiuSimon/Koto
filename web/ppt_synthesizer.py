#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT合成引擎 - 使用蓝图生成高质量PPT
包含高级排版、美化、图像集成等功能
"""

import logging
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)


class PPTSynthesizer:
    """PPT合成器 - 从蓝图生成最终PPT"""

    def __init__(self, theme: str = "business"):
        self.theme = theme
        self.pptx = None
        self.slide_count = 0

    async def synthesize_from_blueprint(
        self,
        blueprint,  # PPTBlueprint
        output_path: str,
        apply_beauty_rules: bool = True,
        image_paths: Optional[Dict[int, List[str]]] = None,
        progress_callback=None,
    ) -> Dict[str, Any]:
        """
        从蓝图合成PPT

        Args:
            blueprint: PPTBlueprint对象
            output_path: 输出文件路径
            apply_beauty_rules: 是否应用美化规则
            image_paths: 幻灯片索引 -> 图像路径列表
            progress_callback: 进度回调 (msg, progress)

        Returns:
            {
                "success": bool,
                "output_path": str,
                "slide_count": int,
                "details": {...}
            }
        """

        def _report(msg, p=None):
            if progress_callback:
                try:
                    progress_callback(msg, p)
                except:
                    pass

        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt

            _report("初始化 PowerPoint 引擎...", 50)

            # 创建演示文稿
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # 设置主题配色
            theme_colors = self._get_theme_colors(blueprint.theme)

            total_slides = len(blueprint.slides)

            # 遍历所有幻灯片
            for i, slide_blueprint in enumerate(blueprint.slides):
                current_slide_num = i + 1
                progress_pct = 50 + int(
                    (current_slide_num / total_slides) * 45
                )  # 50% -> 95%
                _report(
                    f"正在渲染幻灯片 {current_slide_num}/{total_slides}: {slide_blueprint.title} ({slide_blueprint.slide_type.value})",
                    progress_pct,
                )

                slide_layout = self._select_slide_layout(prs, slide_blueprint)
                slide = prs.slides.add_slide(slide_layout)

                # 填充内容
                await self._populate_slide(
                    slide, slide_blueprint, theme_colors, image_paths=image_paths
                )

                # 应用美化规则
                if apply_beauty_rules:
                    self._apply_beauty_rules(slide, slide_blueprint, theme_colors)

                self.slide_count += 1

            _report(f"幻灯片渲染完成，正在保存文件...", 98)

            # 生成文件
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            prs.save(output_path)

            return {
                "success": True,
                "output_path": output_path,
                "slide_count": self.slide_count,
                "file_size": os.path.getsize(output_path) / 1024,
                "blueprint_info": (
                    blueprint.to_dict() if hasattr(blueprint, "to_dict") else {}
                ),
            }

        except Exception as e:
            import traceback

            traceback.print_exc()
            return {
                "success": False,
                "error": str(e),
                "traceback": traceback.format_exc(),
            }

    def _select_slide_layout(self, prs, slide_blueprint) -> Any:
        """根据幻灯片类型选择布局"""

        # 使用空白布局为所有幻灯片，以便自定义
        blank_layout = prs.slide_layouts[6]  # 空白布局
        return blank_layout

    async def _populate_slide(
        self,
        slide,
        slide_blueprint,
        theme_colors: Dict[str, Tuple],
        image_paths: Optional[Dict[int, List[str]]] = None,
    ):
        """填充幻灯片内容"""

        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        layout_config = slide_blueprint.layout_config

        # 根据幻灯片类型填充内容
        if slide_blueprint.slide_type.value == "title":
            await self._populate_title_slide(slide, slide_blueprint, theme_colors)

        elif slide_blueprint.slide_type.value == "section":
            await self._populate_section_slide(slide, slide_blueprint, theme_colors)

        elif slide_blueprint.slide_type.value == "content":
            await self._populate_content_slide(slide, slide_blueprint, theme_colors)

        elif slide_blueprint.slide_type.value == "content_image":
            await self._populate_content_image_slide(
                slide, slide_blueprint, theme_colors, image_paths
            )

        elif slide_blueprint.slide_type.value == "summary":
            await self._populate_summary_slide(slide, slide_blueprint, theme_colors)

        elif slide_blueprint.slide_type.value == "comparison":
            await self._populate_comparison_slide(slide, slide_blueprint, theme_colors)

        elif slide_blueprint.slide_type.value == "data":
            # [Phase 3] Data/Chart Slide
            await self._populate_data_slide(slide, slide_blueprint, theme_colors)

        else:
            await self._populate_content_slide(slide, slide_blueprint, theme_colors)

    async def _populate_comparison_slide(self, slide, slide_blueprint, theme_colors):
        """[Phase 2] 填充对比页 (Comparison Layout)"""
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        self._setup_background(slide, theme_colors)
        self._add_header(slide, slide_blueprint.title, theme_colors)

        # Left Panel (Arguments/Pros/Side A)
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.5)
        )
        lf = left_box.text_frame
        lf.word_wrap = True

        # Right Panel (Counter-arguments/Cons/Side B)
        right_box = slide.shapes.add_textbox(
            Inches(5.3), Inches(1.5), Inches(4.2), Inches(5.5)
        )
        rf = right_box.text_frame
        rf.word_wrap = True

        content = slide_blueprint.content
        mid = len(content) // 2
        left_points = content[:mid]
        right_points = content[mid:]

        for p_list, frame in [(left_points, lf), (right_points, rf)]:
            for i, txt in enumerate(p_list):
                if i > 0:
                    frame.add_paragraph()
                p = frame.paragraphs[i]
                p.text = f"• {txt}"
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(*theme_colors["text"])

        # Divider Line
        line = slide.shapes.add_shape(1, Inches(5), Inches(1.5), Inches(0), Inches(5))
        line.line.color.rgb = RGBColor(*theme_colors["accent"])
        line.line.width = Pt(2)

    async def _populate_data_slide(self, slide, slide_blueprint, theme_colors):
        """[Phase 3] 填充数据页 (Big Number / Chart Placeholder)"""
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        self._setup_background(slide, theme_colors)
        self._add_header(slide, slide_blueprint.title, theme_colors)

        # Big Number Focus
        # Assuming the first point contains the main stat (e.g., "75% Increase")
        main_stat = "DATA"
        description = "Key Metric"

        if slide_blueprint.content:
            # Try to find a number or short phrase
            main_stat = slide_blueprint.content[0]
            if len(slide_blueprint.content) > 1:
                description = slide_blueprint.content[1]

        # Big Central Number
        stat_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(2)
        )
        sf = stat_box.text_frame
        p = sf.paragraphs[0]
        p.text = main_stat
        p.font.size = Pt(80)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(*theme_colors["accent"])

        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(2), Inches(4.5), Inches(6), Inches(2)
        )
        df = desc_box.text_frame
        p2 = df.paragraphs[0]
        p2.text = description
        p2.font.size = Pt(24)
        p2.alignment = PP_ALIGN.CENTER
        p2.font.color.rgb = RGBColor(*theme_colors["text"])

    def _setup_background(self, slide, theme_colors):
        """Helper for background setup"""
        from pptx.dml.color import RGBColor

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme_colors["background"])

    def _add_header(self, slide, title, theme_colors):
        """Helper for header setup"""
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
        )
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(*theme_colors["primary"])
        # Divider
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
        line.line.color.rgb = RGBColor(*theme_colors["accent"])
        line.line.width = Pt(2)

    async def _populate_title_slide(
        self, slide, slide_blueprint, theme_colors: Dict[str, Tuple]
    ):
        """填充标题页"""

        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        # 添加背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme_colors["primary"])

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_p = title_frame.paragraphs[0]
        title_p.text = slide_blueprint.title
        title_p.font.size = Pt(60)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.alignment = PP_ALIGN.CENTER

        # 添加副标题/subtitle
        if slide_blueprint.content:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.8), Inches(9), Inches(1.5)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.text = (
                slide_blueprint.content[0]
                if isinstance(slide_blueprint.content, list)
                else str(slide_blueprint.content)
            )
            subtitle_p.font.size = Pt(28)
            subtitle_p.font.color.rgb = RGBColor(*theme_colors["accent"])
            subtitle_p.alignment = PP_ALIGN.CENTER

    async def _populate_section_slide(
        self, slide, slide_blueprint, theme_colors: Dict[str, Tuple]
    ):
        """填充章节页"""

        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        # 背景色 - 稍浅
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(
            int(theme_colors["primary"][0] * 0.8),
            int(theme_colors["primary"][1] * 0.8),
            int(theme_colors["primary"][2] * 0.8),
        )

        # 添加章节标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(9), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_p = title_frame.paragraphs[0]
        title_p.text = slide_blueprint.title
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.alignment = PP_ALIGN.CENTER

    async def _populate_content_slide(
        self, slide, slide_blueprint, theme_colors: Dict[str, Tuple]
    ):
        """填充内容页"""

        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme_colors["background"])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_blueprint.title
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(*theme_colors["primary"])

        # 添加分割线
        line = slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.3), Inches(9), Inches(0)  # 直线
        )
        line.line.color.rgb = RGBColor(*theme_colors["accent"])
        line.line.width = Pt(2)

        # 内容点
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8), Inches(8.5), Inches(5.2)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, point in enumerate(slide_blueprint.content):
            if i > 0:
                text_frame.add_paragraph()

            p = text_frame.paragraphs[i]
            p.text = f"• {point}"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(*theme_colors["text"])
            p.level = 0
            p.space_before = Pt(8)
            p.space_after = Pt(8)

    async def _populate_content_image_slide(
        self,
        slide,
        slide_blueprint,
        theme_colors: Dict[str, Tuple],
        image_paths: Optional[Dict[int, List[str]]] = None,
    ):
        """填充文字+图片页"""

        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme_colors["background"])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_blueprint.title
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(*theme_colors["primary"])

        # 左侧内容
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, point in enumerate(slide_blueprint.content):
            if i > 0:
                text_frame.add_paragraph()

            p = text_frame.paragraphs[i]
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(*theme_colors["text"])
            p.space_before = Pt(6)
            p.space_after = Pt(6)

        # 右侧图片
        image_list = (
            image_paths.get(slide_blueprint.slide_index, []) if image_paths else []
        )
        if image_list and os.path.exists(image_list[0]):
            try:
                slide.shapes.add_picture(
                    image_list[0], Inches(5.3), Inches(1.5), width=Inches(4)
                )
            except Exception as e:
                logger.info(f"[PPTSynthesizer] 添加图片失败: {e}")

    async def _populate_summary_slide(
        self, slide, slide_blueprint, theme_colors: Dict[str, Tuple]
    ):
        """填充总结页"""

        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        # 背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme_colors["primary"])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_p = title_frame.paragraphs[0]
        title_p.text = slide_blueprint.title
        title_p.font.size = Pt(64)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(*theme_colors["accent"])
        title_p.alignment = PP_ALIGN.CENTER

    def _apply_beauty_rules(
        self, slide, slide_blueprint, theme_colors: Dict[str, Tuple]
    ):
        """应用美化规则"""

        from pptx.dml.color import RGBColor

        # 为所有形状添加阴影和边框（如果配置了）
        layout_config = slide_blueprint.layout_config

        if layout_config.get("shadow_enabled"):
            # 这里可以添加阴影效果（需要通过XML）
            pass

        if layout_config.get("border_accent"):
            # 添加强调边框（可选）
            pass

    def _get_theme_colors(self, theme: str) -> Dict[str, Tuple]:
        """获取主题配色"""

        themes = {
            "business": {
                "primary": (54, 96, 146),  # 商务蓝
                "accent": (255, 192, 0),  # 金色
                "background": (255, 255, 255),  # 白色
                "text": (0, 0, 0),  # 黑色
            },
            "tech": {
                "primary": (0, 120, 215),  # 科技蓝
                "accent": (16, 124, 16),  # 绿色
                "background": (245, 245, 245),  # 浅灰
                "text": (31, 31, 31),  # 深灰
            },
            "creative": {
                "primary": (156, 39, 176),  # 紫色
                "accent": (255, 87, 34),  # 橙色
                "background": (255, 255, 255),  # 白色
                "text": (33, 33, 33),  # 深灰
            },
        }

        return themes.get(theme, themes["business"])


class PPTBeautyOptimizer:
    """PPT美化优化器 - 应用高级美化规则"""

    @staticmethod
    def optimize_slide_aesthetics(slide, slide_blueprint, theme_colors: Dict):
        """优化幻灯片美观度"""

        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        # 1. 文本对齐和间距优化
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    # 优化行距
                    if slide_blueprint.density.value == "dense":
                        paragraph.line_spacing = 1.3
                    elif slide_blueprint.density.value == "light":
                        paragraph.line_spacing = 1.6
                    else:
                        paragraph.line_spacing = 1.5

                    # 优化字体
                    for run in paragraph.runs:
                        run.font.name = "微软雅黑"
                        if not run.font.bold:
                            run.font.bold = False

        # 2. 形状美化（如果有边框配置）
        layout_config = slide_blueprint.layout_config
        bullet_style = layout_config.get("bullet_style", "circle")

        # 3. 添加装饰元素（可选）
        pass

    @staticmethod
    def add_visual_hierarchy(slide, theme_colors: Dict):
        """添加视觉层次"""

        # 根据幻灯片类型添加不同的视觉强调
        pass

    @staticmethod
    def optimize_image_placement(slide, images: List[str], layout_type: str):
        """优化图片放置"""

        if not images or len(images) == 0:
            return

        from pptx.util import Inches

        # 根据布局类型放置图片
        if layout_type == "balanced":
            # 右侧放置单张图片
            if os.path.exists(images[0]):
                slide.shapes.add_picture(
                    images[0], Inches(5.3), Inches(1.5), width=Inches(4)
                )

        elif layout_type == "image_dominant":
            # 全屏放置
            if os.path.exists(images[0]):
                slide.shapes.add_picture(
                    images[0], Inches(0.5), Inches(1.2), width=Inches(9)
                )


class PPTQualityEnsurance:
    """PPT质量保证 - 检查最终PPT质量"""

    @staticmethod
    async def verify_blueprint_quality(blueprint) -> Dict[str, Any]:
        """验证蓝图质量"""

        checks = {
            "slide_count_valid": 5 <= len(blueprint.slides) <= 20,
            "all_slides_titled": all(s.title for s in blueprint.slides),
            "mixed_slide_types": len(set(s.slide_type for s in blueprint.slides)) > 2,
            "layout_planned": all(s.layout_config for s in blueprint.slides),
            "content_density_varied": len(set(s.density for s in blueprint.slides)) > 1,
            "image_prompts_generated": sum(
                len(s.image_prompts) for s in blueprint.slides
            )
            > 0,
        }

        quality_score = sum(checks.values()) / len(checks) * 100

        return {
            "quality_score": quality_score,
            "checks": checks,
            "recommendations": PPTQualityEnsurance._generate_recommendations(checks),
        }

    @staticmethod
    def _generate_recommendations(checks: Dict[str, bool]) -> List[str]:
        """生成改进建议"""

        recommendations = []

        if not checks.get("slide_count_valid"):
            recommendations.append("调整幻灯片数量至5-20张之间")

        if not checks.get("all_slides_titled"):
            recommendations.append("确保所有幻灯片都有标题")

        if not checks.get("mixed_slide_types"):
            recommendations.append("增加幻灯片类型多样性")

        if not checks.get("layout_planned"):
            recommendations.append("为所有幻灯片规划布局")

        if not checks.get("content_density_varied"):
            recommendations.append("调整内容密度以避免单调")

        if not checks.get("image_prompts_generated"):
            recommendations.append("添加更多的图像")

        return recommendations
