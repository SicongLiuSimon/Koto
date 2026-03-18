#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
File Quality Checker — 文件质量自检系统
对 Koto 生成的文件（PPT/Word/Excel）进行质量评估和内容清洗。
- 检测并清除 Markdown 残留 (**, *, `, #, - 等)
- 检测 AI 模型常见无意义内容 ("当然可以", "以下是...")
- 评估内容丰富度 / 格式质量
- 提供自动修复和改进建议
- strip_markdown_for_export: 转换 Word/Excel 时彻底去除 Markdown 符号（持久功能）
"""

import os
import re
from typing import Any, Dict, List, Optional, Tuple

# ────── 导出清洗：转成 Word/Excel 时强制去除的 Markdown 符号 ──────
# 这些规则会在 check_and_fix_for_export() 中永久生效
_EXPORT_MD_STRIP_PATTERNS = [
    # 加粗/斜体/粗斜体 — 保留内部文字
    (r"\*\*\*(.+?)\*\*\*", r"\1"),  # ***text***
    (r"\*\*(.+?)\*\*", r"\1"),  # **text**
    (r"__(.+?)__", r"\1"),  # __text__
    (r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"\1"),  # *text* (单星号斜体)
    (r"(?<!_)_(?!_)(.+?)(?<!_)_(?!_)", r"\1"),  # _text_
    # 删除线
    (r"~~(.+?)~~", r"\1"),
    # 行内代码
    (r"`{3}[\s\S]*?`{3}", ""),  # 代码块整体删除（Excel/Word 中无意义）
    (r"``(.+?)``", r"\1"),
    (r"`([^`\n]+)`", r"\1"),
    # 链接 [text](url) → text
    (r"\[([^\]]+)\]\([^\)]*\)", r"\1"),
    # 图片 ![alt](url) → 空
    (r"!\[[^\]]*\]\([^\)]*\)", ""),
    # 标题 # ## ### — 仅保留文字 (行首)
    (r"(?m)^#{1,6}\s+", ""),
    # 引用块 > text → text
    (r"(?m)^\s*>\s?", ""),
    # 水平线
    (r"(?m)^\s*[-_*]{3,}\s*$", ""),
    # 多余的 * 或 _ 残留（配对失败的孤立符号）
    (r"\*{2,}", ""),
    (r"(?<!\w)_(?!\w)", ""),
    # 表格分隔符行 |---|---| → 空行
    (r"(?m)^\|?[\s:|-]+\|[\s:|-|]+\|?\s*$", ""),
]

# 表格单元格内的 Markdown 也要处理（去除后保留纯文字用于 Excel 单元格值）
_CELL_STRIP_PATTERNS = [
    (r"\*\*\*(.+?)\*\*\*", r"\1"),
    (r"\*\*(.+?)\*\*", r"\1"),
    (r"__(.+?)__", r"\1"),
    (r"\*(.+?)\*", r"\1"),
    (r"_(.+?)_", r"\1"),
    (r"`([^`]+)`", r"\1"),
    (r"\[([^\]]+)\]\([^\)]*\)", r"\1"),
]


def strip_markdown_for_export(text: str, target_format: str = "word") -> str:
    """
    将 Markdown 文本转换为适合导出到 Word/Excel 的纯格式文本。
    ─────────────────────────────────────────────────────────
    对于 Word: 保留段落结构 (# → 去掉 # 符号, - → 去掉 -), 文字完整保留
    对于 Excel: 去除所有 Markdown, 保留纯文字, 适合写入单元格
    ─────────────────────────────────────────────────────────
    这是一个永久性持久功能：每次导出文件时都会自动调用。
    """
    if not text:
        return text

    result = text

    # 先处理多行代码块 (完整删除代码围栏, 保留内容)
    result = re.sub(r"```[a-zA-Z0-9_-]*\n([\s\S]*?)```", r"\1", result)
    result = result.replace("```", "")

    # 应用导出清洗规则
    for pattern, replacement in _EXPORT_MD_STRIP_PATTERNS:
        result = re.sub(pattern, replacement, result, flags=re.MULTILINE | re.DOTALL)

    # Excel 模式: 进一步清理表格行的 | 分隔符 (只保留单元格文字)
    if target_format == "excel":
        # 处理每个 | 分隔的行: | 内容A | 内容B | → 内容A\t内容B
        def _table_row_to_plain(m):
            cells = [c.strip() for c in m.group(0).strip("|").split("|")]
            # 每个单元格内也清洗残留 Markdown
            for pat, rep in _CELL_STRIP_PATTERNS:
                cells = [re.sub(pat, rep, c) for c in cells]
            return "\t".join(cells)

        result = re.sub(r"(?m)^\|.+\|$", _table_row_to_plain, result)

    # 清理多余空白行（不超过 2 个连续空行）
    result = re.sub(r"\n{3,}", "\n\n", result)

    # 清理每行首尾空白
    result = "\n".join(line.rstrip() for line in result.split("\n"))

    return result.strip()


def strip_markdown_from_cell(cell_value: str) -> str:
    """单个 Excel 单元格值的 Markdown 清洗（用于写入 openpyxl 之前）"""
    if not cell_value or not isinstance(cell_value, str):
        return cell_value
    result = cell_value
    for pattern, replacement in _CELL_STRIP_PATTERNS:
        result = re.sub(pattern, replacement, result)
    # 清理多余的 * # ` 残留
    result = re.sub(r"[*`#]{1,}", "", result)
    return result.strip()


def detect_markdown_in_export(text: str) -> List[str]:
    """
    检测文本中存在的 Markdown 符号，返回问题描述列表。
    用于检查层（check layer）报告问题。
    """
    issues = []
    checks = [
        (r"\*\*[^*]+\*\*", "存在 **加粗** 标记"),
        (r"(?<!\*)\*[^*\n]+\*(?!\*)", "存在 *斜体* 标记"),
        (r"`[^`\n]+`", "存在 `行内代码` 标记"),
        (r"(?m)^#{1,6}\s", "存在 # 标题符号"),
        (r"(?m)^\s*[-*+]\s", "存在 - 列表符号"),
        (r"(?m)^\s*\d+\.\s", "存在 1. 编号列表"),
        (r"~~[^~]+~~", "存在 ~~删除线~~ 标记"),
        (r"\[[^\]]+\]\([^\)]+\)", "存在 Markdown 链接"),
        (r"(?m)^>\s", "存在 > 引用块"),
    ]
    for pattern, desc in checks:
        if re.search(pattern, text):
            issues.append(desc)
    return issues


# ────── 内容清洗规则 ──────

# AI 模型常见无意义前缀/后缀 (需要移除)
AI_NOISE_PATTERNS = [
    # 中文
    r"^(当然可以[!！。，,\s]*)",
    r"^(好的[!！。，,\s]*)",
    r"^(没问题[!！。，,\s]*)",
    r"^(以下是.*?[：:]\s*)",
    r"^(下面是.*?[：:]\s*)",
    r"^(这是.*?[：:]\s*)",
    r"^(我来为你.*?[：:。]\s*)",
    r"^(我为你.*?[：:。]\s*)",
    r"^(根据你的要求.*?[：:。]\s*)",
    r"^(根据您的要求.*?[：:。]\s*)",
    r"^(请看以下.*?[：:]\s*)",
    r"^(为您生成.*?[：:。]\s*)",
    r"^(如你所愿.*?[：:。]\s*)",
    r"^(让我来.*?[：:。]\s*)",
    # "以下是" 无冒号时也移除
    r"^(以下是[^。：:]*?(?:的|了)?[。\s]*)",
    # 英文
    r"^(Sure[!,.\s]*)",
    r"^(Of course[!,.\s]*)",
    r"^(Here(?:\'s| is| are).*?[:：]\s*)",
    r"^(Let me.*?[:：。.]\s*)",
    r"^(I\'ll.*?[:：。.]\s*)",
    r"^(Below is.*?[:：]\s*)",
    r"^(The following.*?[:：]\s*)",
]

# AI 模型对话痕迹（不应出现在文件正文中）
AI_DIALOGUE_PATTERNS = [
    r"(如果你需要.*?[，,].*?可以告诉我)",
    r"(如果您需要.*?[，,].*?请告诉我)",
    r"(希望这对你有帮助)",
    r"(希望这能帮到你)",
    r"(需要我.*?吗[？?])",
    r"(还有什么.*?可以帮.*?[？?])",
    r"(如有.*?问题.*?随时.*?)",
    r"(Is there anything else.*?\?)",
    r"(Let me know if.*)",
    r"(Feel free to.*)",
    r"(Happy to help.*)",
]

# Markdown 残留标记
MARKDOWN_RESIDUE_PATTERNS = [
    (r"\*\*(.+?)\*\*", r"\1"),  # **bold** → bold
    (r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"\1"),  # *italic* → italic
    (r"``(.+?)``", r"\1"),  # ``code`` → code (double backtick)
    (r"`(.+?)`", r"\1"),  # `code` → code
    (r"^#{1,6}\s+", ""),  # ### heading → heading
    (r"^\s*[-•]\s+", ""),  # - bullet → bullet
    (r"^\s*\d+[.、)]\s+", ""),  # 1. numbered → numbered
    (r"~~(.+?)~~", r"\1"),  # ~~strikethrough~~ → strikethrough
    (r"\[([^\]]+)\]\([^\)]+\)", r"\1"),  # [link](url) → link
    (r"(?:^|\s)>\s+", " "),  # > quote → quote
    (r"```[\s\S]*?```", ""),  # ```code block```
    (r"\*{3,}", ""),  # ***
    (r"_{3,}", ""),  # ___
    (r"—{2,}", "—"),  # ——— → —
]


class ContentSanitizer:
    """对输入文本进行清洗，移除 AI 残留 / Markdown 标记"""

    @staticmethod
    def sanitize_text(text: str, context: str = "general") -> str:
        """
        清洗单段文本。
        context: "ppt_point" | "ppt_title" | "doc_paragraph" | "general"
        """
        if not text:
            return text

        original = text
        text = text.strip()
        was_modified = False

        # 1. 移除 AI 无意义开头
        for pat in AI_NOISE_PATTERNS:
            new_text = re.sub(pat, "", text, flags=re.IGNORECASE).strip()
            if new_text != text:
                was_modified = True
                text = new_text

        # 2. 移除 AI 对话痕迹
        for pat in AI_DIALOGUE_PATTERNS:
            new_text = re.sub(pat, "", text, flags=re.IGNORECASE).strip()
            if new_text != text:
                was_modified = True
                text = new_text

        # 3. 移除 Markdown 残留
        if context in ("ppt_point", "ppt_title", "general"):
            for pattern, replacement in MARKDOWN_RESIDUE_PATTERNS:
                new_text = re.sub(
                    pattern, replacement, text, flags=re.MULTILINE
                ).strip()
                if new_text != text:
                    was_modified = True
                    text = new_text

        # 4. 清理多余空白和符号
        text = re.sub(r"\s{2,}", " ", text)  # 多余空格
        text = re.sub(r"[\']{2,}", "", text)  # 多余单引号
        text = text.strip('\'"""「」【】')

        # 如果清洗后为空但确实匹配了 AI 噪声模式，说明整段都是 AI 废话，返回空
        # 只有在没有任何模式匹配时，才回退到原始文本（保护正常内容）
        if not text and was_modified:
            return ""
        return text if text else original

    @staticmethod
    def sanitize_ppt_outline(outline: list) -> Tuple[list, List[str]]:
        """
        清洗 PPT 大纲数据结构。
        返回 (清洗后的大纲, 修复日志列表)
        """
        fixes = []
        for slide in outline:
            # 清洗标题
            old_title = slide.get("title", "")
            new_title = ContentSanitizer.sanitize_text(old_title, "ppt_title")
            if new_title != old_title:
                slide["title"] = new_title
                fixes.append(f"标题清洗: '{old_title[:30]}' → '{new_title[:30]}'")

            # 清洗要点
            for key in ("points", "content"):
                pts = slide.get(key, [])
                if not isinstance(pts, list):
                    continue
                cleaned_pts = []
                for pt in pts:
                    if not isinstance(pt, str):
                        cleaned_pts.append(pt)
                        continue
                    cleaned = ContentSanitizer.sanitize_text(pt, "ppt_point")
                    if cleaned != pt:
                        fixes.append(f"要点清洗: '{pt[:30]}' → '{cleaned[:30]}'")
                    if cleaned:  # 过滤空要点
                        cleaned_pts.append(cleaned)
                slide[key] = cleaned_pts

            # 清洗子主题
            for sub in slide.get("subsections", []):
                old_sub = sub.get("subtitle", "")
                new_sub = ContentSanitizer.sanitize_text(old_sub, "ppt_title")
                if new_sub != old_sub:
                    sub["subtitle"] = new_sub
                    sub["label"] = new_sub
                    fixes.append(f"子标题清洗: '{old_sub[:30]}' → '{new_sub[:30]}'")
                sub_pts = sub.get("points", [])
                cleaned_sub_pts = []
                for pt in sub_pts:
                    if isinstance(pt, str):
                        cleaned = ContentSanitizer.sanitize_text(pt, "ppt_point")
                        if cleaned:
                            cleaned_sub_pts.append(cleaned)
                    else:
                        cleaned_sub_pts.append(pt)
                sub["points"] = cleaned_sub_pts

            # 清洗 left/right（对比页）
            for side in ("left", "right"):
                side_data = slide.get(side)
                if side_data and isinstance(side_data, dict):
                    for k in ("subtitle", "label"):
                        if k in side_data:
                            side_data[k] = ContentSanitizer.sanitize_text(
                                side_data[k], "ppt_title"
                            )
                    if "points" in side_data:
                        side_data["points"] = [
                            ContentSanitizer.sanitize_text(p, "ppt_point")
                            for p in side_data["points"]
                            if isinstance(p, str)
                            and ContentSanitizer.sanitize_text(p, "ppt_point")
                        ]

        return outline, fixes

    @staticmethod
    def sanitize_document_text(text: str) -> Tuple[str, List[str]]:
        """
        清洗文档文本（Word/PDF），保留正确的 Markdown 格式用于 docx 转换。
        仅移除 AI 对话痕迹和无意义前缀。
        """
        fixes = []
        lines = text.split("\n")
        cleaned_lines = []

        for line in lines:
            original = line
            stripped = line.strip()

            # 移除 AI 无意义开头（仅对第一段有效）
            if not cleaned_lines and stripped:
                for pat in AI_NOISE_PATTERNS:
                    new = re.sub(pat, "", stripped, flags=re.IGNORECASE).strip()
                    if new != stripped:
                        fixes.append(f"移除AI前缀: '{stripped[:40]}'")
                        stripped = new
                        line = stripped

            # 移除 AI 对话痕迹（任何位置）
            for pat in AI_DIALOGUE_PATTERNS:
                new = re.sub(pat, "", line, flags=re.IGNORECASE).strip()
                if new != line:
                    fixes.append(f"移除对话痕迹: '{line.strip()[:40]}'")
                    line = new

            # 移除多余单引号聚集
            line = re.sub(r"'{2,}", "", line)

            cleaned_lines.append(line)

        return "\n".join(cleaned_lines), fixes


class FileQualityEvaluator:
    """
    文件质量评估器 — 对生成的文件进行综合打分。
    """

    QUALITY_THRESHOLDS = {
        "ppt": {
            "min_slides": 4,
            "max_slides": 25,
            "min_points_per_slide": 2,
            "max_points_per_slide": 8,
            "min_point_length": 8,
            "max_point_length": 200,
            "min_title_length": 2,
            "required_title": True,
        },
        "docx": {
            "min_chars": 200,
            "min_paragraphs": 3,
            "min_headings": 1,
        },
        "xlsx": {
            "min_rows": 2,
            "min_cols": 2,
        },
    }

    @classmethod
    def evaluate_ppt_outline(
        cls, outline: list, user_request: str = ""
    ) -> Dict[str, Any]:
        """
        评估 PPT 大纲质量（在渲染前调用）。
        返回 { score, issues, suggestions, pass }
        """
        score = 100
        issues = []
        suggestions = []
        thresholds = cls.QUALITY_THRESHOLDS["ppt"]

        # 1. 幻灯片数量
        slide_count = len(outline)
        if slide_count < thresholds["min_slides"]:
            score -= 15
            issues.append(f"幻灯片数量过少 ({slide_count} 页)")
            suggestions.append("建议增加更多章节或展开内容")
        elif slide_count > thresholds["max_slides"]:
            score -= 10
            issues.append(f"幻灯片数量过多 ({slide_count} 页)")
            suggestions.append("建议合并相似内容，精简页面")

        # 2. 检查每页质量
        empty_slides = 0
        thin_slides = 0
        markdown_contaminated = 0
        total_points = 0
        titles_with_issues = 0

        for i, slide in enumerate(outline):
            stype = slide.get("type", "detail")
            title = slide.get("title", "")
            points = slide.get("points", []) or slide.get("content", [])
            subs = slide.get("subsections", [])

            # 标题检查
            if not title or len(title.strip()) < thresholds["min_title_length"]:
                titles_with_issues += 1
                issues.append(f"第 {i+1} 页标题缺失或过短")

            # 内容检查（过渡页除外）
            if stype == "divider":
                continue

            if stype in ("overview", "comparison") and subs:
                sub_pts = sum(len(sub.get("points", [])) for sub in subs)
                if sub_pts == 0:
                    empty_slides += 1
                elif sub_pts < len(subs) * 2:
                    thin_slides += 1
                total_points += sub_pts
            else:
                if not points:
                    empty_slides += 1
                elif len(points) < thresholds["min_points_per_slide"]:
                    thin_slides += 1
                total_points += len(points)

                # 每个要点长度检查
                for pt in points:
                    if isinstance(pt, str):
                        if len(pt.strip()) < thresholds["min_point_length"]:
                            thin_slides += 1
                            break
                        # Markdown 残留检查
                        if re.search(r"\*\*|`|^#{1,3}\s", pt):
                            markdown_contaminated += 1
                            break

        if empty_slides > 0:
            score -= empty_slides * 15
            issues.append(f"{empty_slides} 页内容为空")
            suggestions.append("需要为空页面补充内容")

        if thin_slides > 0:
            penalty = min(30, thin_slides * 8)
            score -= penalty
            issues.append(f"{thin_slides} 页内容单薄（要点过少或过短）")
            suggestions.append("建议每页至少 3-4 个有实质内容的要点")

        if markdown_contaminated > 0:
            score -= markdown_contaminated * 5
            issues.append(f"{markdown_contaminated} 页存在 Markdown 标记残留")
            suggestions.append("需要清洗 Markdown 标记 (**, `, # 等)")

        if titles_with_issues > 0:
            score -= titles_with_issues * 5
            suggestions.append("确保每页都有明确的标题")

        # 3. 内容丰富度
        content_slides = [s for s in outline if s.get("type", "detail") != "divider"]
        avg_points = total_points / max(len(content_slides), 1)
        if avg_points < 2.5:
            score -= 10
            issues.append(f"平均每页要点数偏低 ({avg_points:.1f})")

        # 4. 类型多样性
        types_used = set(s.get("type", "detail") for s in outline)
        if len(types_used) <= 1 and slide_count >= 6:
            score -= 5
            suggestions.append("建议使用多种幻灯片类型增加视觉多样性")

        # 5. 用户需求匹配度（简单关键词检查）
        if user_request:
            all_text = " ".join(
                (
                    s.get("title", "")
                    + " "
                    + " ".join(
                        s.get("points", [])
                        if isinstance(s.get("points", []), list)
                        else []
                    )
                )
                for s in outline
            )
            # 提取用户请求中的关键词
            keywords = re.findall(r"[\u4e00-\u9fff]{2,}|[a-zA-Z]{3,}", user_request)
            matched = sum(1 for kw in keywords if kw.lower() in all_text.lower())
            match_ratio = matched / max(len(keywords), 1)
            if match_ratio < 0.3:
                score -= 10
                issues.append("内容与用户需求的关键词匹配度偏低")

        score = max(0, min(100, score))
        return {
            "score": score,
            "pass": score >= 60,
            "issues": issues,
            "suggestions": suggestions,
            "metrics": {
                "slide_count": slide_count,
                "total_points": total_points,
                "avg_points": round(avg_points, 1),
                "empty_slides": empty_slides,
                "thin_slides": thin_slides,
                "markdown_contaminated": markdown_contaminated,
                "types_used": list(types_used),
            },
        }

    @classmethod
    def evaluate_document_text(
        cls, text: str, user_request: str = ""
    ) -> Dict[str, Any]:
        """评估文档文本质量（Word/PDF 内容）"""
        score = 100
        issues = []
        suggestions = []

        if not text or not text.strip():
            return {
                "score": 0,
                "pass": False,
                "issues": ["文档内容为空"],
                "suggestions": ["重新生成文档"],
            }

        char_count = len(text.strip())
        lines = [l for l in text.strip().split("\n") if l.strip()]
        paragraphs = [l for l in lines if not l.strip().startswith("#")]
        headings = [l for l in lines if re.match(r"^#{1,6}\s+", l.strip())]

        # 字数
        if char_count < 200:
            score -= 20
            issues.append(f"文档内容过短 ({char_count} 字)")
            suggestions.append("建议丰富内容，至少 500 字")
        elif char_count < 500:
            score -= 10
            issues.append(f"文档内容偏短 ({char_count} 字)")

        # 段落数
        if len(paragraphs) < 3:
            score -= 15
            issues.append("段落过少")
            suggestions.append("增加更多段落使文档结构完整")

        # 标题
        if not headings:
            score -= 10
            issues.append("缺少标题结构")
            suggestions.append("添加 Markdown 标题 (# / ## / ###) 组织结构")

        # Markdown 残留（在文档正文段落中出现）
        md_count = 0
        for line in paragraphs[:30]:
            if re.search(r"\*\*[^*]+\*\*", line):
                # 文档中的 **加粗** 是正常的 Markdown，会被 save_docx 正确转换
                pass
            if re.search(r"'{2,}|`{3,}", line):
                md_count += 1

        # AI 对话痕迹
        ai_noise_count = 0
        for line in lines[:10]:
            for pat in AI_NOISE_PATTERNS + AI_DIALOGUE_PATTERNS:
                if re.search(pat, line, re.IGNORECASE):
                    ai_noise_count += 1
                    break

        if ai_noise_count > 0:
            score -= ai_noise_count * 5
            issues.append(f"检测到 {ai_noise_count} 处 AI 对话痕迹")
            suggestions.append("清除 AI 生成文本中的对话性前缀/后缀")

        score = max(0, min(100, score))
        return {
            "score": score,
            "pass": score >= 60,
            "issues": issues,
            "suggestions": suggestions,
            "metrics": {
                "char_count": char_count,
                "paragraphs": len(paragraphs),
                "headings": len(headings),
                "ai_noise": ai_noise_count,
            },
        }

    @classmethod
    def evaluate_pptx_file(cls, pptx_path: str) -> Dict[str, Any]:
        """
        评估已生成的 PPTX 文件（后置检查）。
        检查实际渲染结果中是否有 Markdown 残留。
        """
        if not os.path.exists(pptx_path):
            return {
                "score": 0,
                "pass": False,
                "issues": ["文件不存在"],
                "suggestions": [],
            }

        try:
            from pptx import Presentation
        except ImportError:
            return {
                "score": 50,
                "pass": True,
                "issues": ["无法检查: python-pptx 未安装"],
                "suggestions": [],
            }

        prs = Presentation(pptx_path)
        score = 100
        issues = []
        suggestions = []

        slide_count = len(prs.slides)
        if slide_count < 3:
            score -= 15
            issues.append(f"PPT 只有 {slide_count} 页")

        md_residue_count = 0
        empty_text_slides = 0
        long_text_slides = 0

        for si, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text or ""
                        slide_text += text + " "

                        # 检测 Markdown 残留
                        if re.search(r"\*{2,}[^*]+\*{2,}", text):
                            md_residue_count += 1
                        if re.search(r"`[^`]+`", text):
                            md_residue_count += 1
                        if re.search(r"^#{1,3}\s", text):
                            md_residue_count += 1
                        # 检测连续单引号
                        if re.search(r"'{2,}", text):
                            md_residue_count += 1

            # 跳过封面和结束页
            if si == 0 or si == slide_count - 1:
                continue

            if len(slide_text.strip()) < 10:
                empty_text_slides += 1
            elif len(slide_text.strip()) > 800:
                long_text_slides += 1

        if md_residue_count > 0:
            score -= min(30, md_residue_count * 5)
            issues.append(f"检测到 {md_residue_count} 处 Markdown 标记残留")
            suggestions.append("需要清除 **, `, # 等 Markdown 标记")

        if empty_text_slides > 0:
            score -= empty_text_slides * 10
            issues.append(f"{empty_text_slides} 页内容为空")

        if long_text_slides > 0:
            score -= long_text_slides * 5
            issues.append(f"{long_text_slides} 页文字过多")
            suggestions.append("精简文字，避免单页内容过密")

        score = max(0, min(100, score))
        return {
            "score": score,
            "pass": score >= 60,
            "issues": issues,
            "suggestions": suggestions,
            "metrics": {
                "slide_count": slide_count,
                "md_residue": md_residue_count,
                "empty_slides": empty_text_slides,
                "long_text_slides": long_text_slides,
            },
        }


class FileQualityGate:
    """
    文件质量门控 — 在文件生成流水线中拦截质量不达标的输出。
    集成 ContentSanitizer + FileQualityEvaluator，提供一键式 API。
    """

    @staticmethod
    def check_and_fix_ppt_outline(
        outline: list, user_request: str = "", progress_callback=None
    ) -> Dict[str, Any]:
        """
        PPT 大纲质量门控：清洗 → 评估 → 返回结果。
        如果质量分 < 60，建议重新生成。

        返回:
        {
            'outline': 清洗后的大纲,
            'quality': 评估结果,
            'fixes': 修复记录,
            'action': 'proceed' | 'warn' | 'regenerate'
        }
        """
        # Step 1: 清洗
        sanitized, fixes = ContentSanitizer.sanitize_ppt_outline(outline)

        if fixes and progress_callback:
            progress_callback(
                f"🧹 已清洗 {len(fixes)} 处内容问题", "移除 Markdown 残留和 AI 对话痕迹"
            )

        # Step 2: 评估
        quality = FileQualityEvaluator.evaluate_ppt_outline(sanitized, user_request)

        if progress_callback:
            emoji = "✅" if quality["pass"] else "⚠️"
            progress_callback(
                f"{emoji} 内容质量评分: {quality['score']}/100",
                "; ".join(quality["issues"][:3]) if quality["issues"] else "质量良好",
            )

        # Step 3: 决策
        if quality["score"] >= 75:
            action = "proceed"
        elif quality["score"] >= 60:
            action = "warn"
        else:
            action = "regenerate"

        return {
            "outline": sanitized,
            "quality": quality,
            "fixes": fixes,
            "action": action,
        }

    @staticmethod
    def check_and_fix_document(
        text: str, user_request: str = "", progress_callback=None
    ) -> Dict[str, Any]:
        """
        文档内容质量门控：清洗 → 评估 → 返回结果。
        """
        # Step 1: 清洗
        sanitized, fixes = ContentSanitizer.sanitize_document_text(text)

        if fixes and progress_callback:
            progress_callback(f"🧹 已清洗 {len(fixes)} 处内容问题", "")

        # Step 2: 评估
        quality = FileQualityEvaluator.evaluate_document_text(sanitized, user_request)

        if progress_callback:
            emoji = "✅" if quality["pass"] else "⚠️"
            progress_callback(
                f"{emoji} 文档质量评分: {quality['score']}/100",
                "; ".join(quality["issues"][:3]) if quality["issues"] else "质量良好",
            )

        if quality["score"] >= 75:
            action = "proceed"
        elif quality["score"] >= 60:
            action = "warn"
        else:
            action = "regenerate"

        return {"text": sanitized, "quality": quality, "fixes": fixes, "action": action}

    @staticmethod
    def post_check_pptx(pptx_path: str, progress_callback=None) -> Dict[str, Any]:
        """
        PPT 文件后置检查（渲染后）。
        """
        quality = FileQualityEvaluator.evaluate_pptx_file(pptx_path)

        if progress_callback:
            emoji = "✅" if quality["pass"] else "⚠️"
            progress_callback(
                f"{emoji} PPT 文件检查: {quality['score']}/100",
                (
                    "; ".join(quality["issues"][:3])
                    if quality["issues"]
                    else "文件质量良好"
                ),
            )

        return quality

    @staticmethod
    def check_and_fix_for_export(
        text: str,
        target_format: str = "word",
        user_request: str = "",
        progress_callback=None,
    ) -> Dict[str, Any]:
        """
        【导出清洗门控】在转换为 Word/Excel/PDF 等格式之前执行的完整检查+清洗。
        ─────────────────────────────────────────────────────────────────────
        此方法是一个「永久性持久功能」：
          - 每次生成文件时都会自动调用
          - 检测 Markdown 符号并报告
          - 自动移除不适合该格式的 Markdown 符号
          - 对于 Word: 保留标题/列表结构（只去掉 ## # 符号，保留其对应的文字层级）
                       注意: save_docx 已经能正确解析 ## → Heading，因此不对 word 做符号去除
                       只清洗真正会变成"乱码"的符号（**在标题里、孤立的 ```）
          - 对于 Excel: 彻底去除所有 Markdown 符号，单元格只保留纯文字
        ─────────────────────────────────────────────────────────────────────
        返回:
        {
          'text':    清洗后的文本,
          'issues':  检测到的 Markdown 问题列表,
          'fixes':   AI 噪声清洗记录,
          'stripped_md': 是否进行了 Markdown 去除 (bool),
          'quality': 评估结果字典
        }
        """
        if not text:
            return {
                "text": text,
                "issues": [],
                "fixes": [],
                "stripped_md": False,
                "quality": {
                    "score": 0,
                    "pass": False,
                    "issues": ["内容为空"],
                    "suggestions": [],
                },
            }

        # Step 1: AI 对话痕迹清洗（所有格式都需要）
        sanitized, fixes = ContentSanitizer.sanitize_document_text(text)

        # Step 2: 检测 Markdown（用于报告）
        md_issues = detect_markdown_in_export(sanitized)

        if md_issues and progress_callback:
            progress_callback(
                f"🔍 检测到 {len(md_issues)} 类 Markdown 符号", "; ".join(md_issues[:3])
            )

        # Step 3: 根据目标格式决定是否需要去除 Markdown
        stripped_md = False
        if target_format == "excel":
            # Excel: 彻底去除所有 Markdown
            sanitized = strip_markdown_for_export(sanitized, target_format="excel")
            stripped_md = bool(md_issues)
            if stripped_md and progress_callback:
                progress_callback(
                    "🧹 Excel 导出：已去除所有 Markdown 符号", "单元格只保留纯文字"
                )

        elif target_format in ("word", "pdf"):
            # Word/PDF: save_docx 已能处理 ## / **，只去除孤立符号
            # 清理"代码围栏"、多余 ``` 等不会被 save_docx 处理的残留
            sanitized = re.sub(r"```[a-zA-Z0-9_-]*\n([\s\S]*?)```", r"\1", sanitized)
            sanitized = sanitized.replace("```", "")
            # 清理完全孤立的 ** （没有配对成粗体的）
            sanitized = re.sub(r"\*\*\s*\*\*", "", sanitized)
            # 清理首段 AI 噪声连带的 # 开头（如果 # 后跟非空格）
            sanitized = re.sub(r"(?m)^(#+)(?!\s)", r"\1 ", sanitized)
            if md_issues and progress_callback:
                md_ok = [i for i in md_issues if "加粗" not in i and "标题" not in i]
                if md_ok:
                    progress_callback(
                        "🔍 Word 导出：Markdown 格式将由排版引擎处理", "已清理孤立符号"
                    )
                stripped_md = True

        # Step 4: 质量评估
        quality = FileQualityEvaluator.evaluate_document_text(sanitized, user_request)
        if fixes:
            quality["issues"] = [f"已清洗 {len(fixes)} 处 AI 噪声"] + quality.get(
                "issues", []
            )

        if progress_callback:
            emoji = "✅" if quality["pass"] else "⚠️"
            progress_callback(
                f"{emoji} 导出质量检查: {quality['score']}/100",
                "; ".join(quality["issues"][:2]) if quality["issues"] else "质量良好",
            )

        return {
            "text": sanitized,
            "issues": md_issues,
            "fixes": fixes,
            "stripped_md": stripped_md,
            "quality": quality,
        }
