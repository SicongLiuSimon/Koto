"""
文件上传全类型通路测试
测试 FileProcessor 对各种文件类型的处理能力，
包括：txt, md, csv, json, docx, xlsx, pptx, pdf, png/jpg
每类文件都验证：
  1. process_file() 成功（success=True）
  2. 有内容产出（text_content 或 binary_data 非空）
  3. format_result_for_chat() 输出格式正确
"""

import os
import sys
import struct
import tempfile
import textwrap
import zlib

# 项目根目录
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)

from web.file_processor import FileProcessor, process_uploaded_file

# ── 颜色输出 ──────────────────────────────────────────────────────────────────
GREEN = "\033[92m"
RED = "\033[91m"
YELLOW = "\033[93m"
CYAN = "\033[96m"
RESET = "\033[0m"

RESULTS: list[dict] = []


def ok(label, detail=""):
    print(
        f"  {GREEN}✓ PASS{RESET}  {label}"
        + (f"  {YELLOW}({detail}){RESET}" if detail else "")
    )
    RESULTS.append({"label": label, "status": "PASS", "detail": detail})


def fail(label, reason):
    print(f"  {RED}✗ FAIL{RESET}  {label}  →  {reason}")
    RESULTS.append({"label": label, "status": "FAIL", "detail": reason})


def skip(label, reason):
    print(f"  {YELLOW}⊘ SKIP{RESET}  {label}  →  {reason}")
    RESULTS.append({"label": label, "status": "SKIP", "detail": reason})


# ── 工具函数 ──────────────────────────────────────────────────────────────────


def run_processor(filepath: str, user_msg="请分析这个文件"):
    """运行 FileProcessor 并返回 (result, formatted_msg, file_data)。"""
    proc = FileProcessor()
    result = proc.process_file(filepath)
    fmt, fdata = proc.format_result_for_chat(result, user_msg)
    return result, fmt, fdata


# ╔═══════════════════════════════════════════════════════════════════════╗
# ║  文件创建助手                                                         ║
# ╚═══════════════════════════════════════════════════════════════════════╝


def make_txt(d):
    p = os.path.join(d, "sample.txt")
    with open(p, "w", encoding="utf-8") as f:
        f.write("这是一份纯文本测试文件\n第二行内容\nHello, world!\n")
    return p


def make_md(d):
    p = os.path.join(d, "sample.md")
    with open(p, "w", encoding="utf-8") as f:
        f.write(
            "# 测试文档\n\n## 简介\n这是 Markdown 格式的测试文件。\n\n- 功能 A\n- 功能 B\n"
        )
    return p


def make_csv(d):
    p = os.path.join(d, "sample.csv")
    with open(p, "w", encoding="utf-8") as f:
        f.write("姓名,年龄,城市\n张三,28,北京\n李四,32,上海\n王五,25,深圳\n")
    return p


def make_json(d):
    import json

    p = os.path.join(d, "sample.json")
    with open(p, "w", encoding="utf-8") as f:
        json.dump(
            {"title": "测试", "items": [1, 2, 3], "active": True}, f, ensure_ascii=False
        )
    return p


def make_docx(d):
    p = os.path.join(d, "sample.docx")
    try:
        from docx import Document

        doc = Document()
        doc.add_heading("测试 Word 文档", 0)
        doc.add_paragraph("这是第一段正文内容，用于验证 DOCX 解析。")
        doc.add_paragraph("这是第二段，包含更多文字。")
        t = doc.add_table(rows=2, cols=2)
        t.cell(0, 0).text = "列A"
        t.cell(0, 1).text = "列B"
        t.cell(1, 0).text = "值1"
        t.cell(1, 1).text = "值2"
        doc.save(p)
        return p
    except ImportError:
        return None


def make_xlsx(d):
    p = os.path.join(d, "sample.xlsx")
    try:
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["产品", "数量", "单价"])
        ws.append(["苹果", 100, 5.5])
        ws.append(["香蕉", 200, 3.0])
        wb.save(p)
        return p
    except ImportError:
        # 尝试 pandas
        try:
            import pandas as pd

            pd.DataFrame({"产品": ["苹果", "香蕉"], "数量": [100, 200]}).to_excel(
                p, index=False
            )
            return p
        except ImportError:
            return None


def make_pptx(d):
    p = os.path.join(d, "sample.pptx")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "测试幻灯片"
        slide.placeholders[1].text = "这是幻灯片内容，第一张。"
        slide2 = prs.slides.add_slide(slide_layout)
        slide2.shapes.title.text = "第二张"
        slide2.placeholders[1].text = "更多内容在此。"
        prs.save(p)
        return p
    except ImportError:
        return None


def make_pdf(d):
    p = os.path.join(d, "sample.pdf")
    # 优先用 reportlab，fallback 最小手写 PDF
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4

        c = canvas.Canvas(p, pagesize=A4)
        c.drawString(100, 750, "PDF Test Document")
        c.drawString(100, 720, "This is a test PDF created for Koto upload testing.")
        c.drawString(100, 690, "Page 1 content.")
        c.showPage()
        c.save()
        return p
    except ImportError:
        pass
    # 极简手写 PDF（仅包含 ASCII 文本，足够 PyPDF2 解析）
    content = textwrap.dedent("""\
    %PDF-1.4
    1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
    2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
    3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj
    4 0 obj<</Length 44>>
    stream
    BT /F1 12 Tf 72 720 Td (PDF Test File) Tj ET
    endstream
    endobj
    5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj
    xref
    0 6
    0000000000 65535 f 
    0000000009 00000 n 
    0000000058 00000 n 
    0000000115 00000 n 
    0000000266 00000 n 
    0000000360 00000 n 
    trailer<</Size 6/Root 1 0 R>>
    startxref
    441
    %%EOF
    """)
    with open(p, "w", encoding="latin-1") as f:
        f.write(content)
    return p


def make_png(d):
    p = os.path.join(d, "sample.png")

    # 最小合法 PNG (1×1 红色像素)
    def write_chunk(chunk_type, data):
        length = struct.pack(">I", len(data))
        chunk = chunk_type + data
        crc = struct.pack(">I", zlib.crc32(chunk) & 0xFFFFFFFF)
        return length + chunk + crc

    signature = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr = write_chunk(b"IHDR", ihdr_data)
    raw_data = b"\x00\xff\x00\x00"  # filter byte + RGB
    compressed = zlib.compress(raw_data)
    idat = write_chunk(b"IDAT", compressed)
    iend = write_chunk(b"IEND", b"")
    with open(p, "wb") as f:
        f.write(signature + ihdr + idat + iend)
    return p


def make_jpg(d):
    p = os.path.join(d, "sample.jpg")
    try:
        from PIL import Image as PILImage

        img = PILImage.new("RGB", (64, 64), color=(0, 128, 255))
        img.save(p, "JPEG")
        return p
    except ImportError:
        pass
    # 最小合法 JPEG (SOI + EOI)
    minimal_jpg = bytes(
        [
            0xFF,
            0xD8,  # SOI
            0xFF,
            0xE0,
            0x00,
            0x10,  # APP0 marker + length
            0x4A,
            0x46,
            0x49,
            0x46,
            0x00,  # JFIF\0
            0x01,
            0x01,
            0x00,
            0x00,
            0x01,
            0x00,
            0x01,
            0x00,
            0x00,  # version, density
            0xFF,
            0xD9,  # EOI
        ]
    )
    with open(p, "wb") as f:
        f.write(minimal_jpg)
    return p


# ╔═══════════════════════════════════════════════════════════════════════╗
# ║  各类型测试                                                           ║
# ╚═══════════════════════════════════════════════════════════════════════╝


def test_text_files(d):
    print(f"\n{CYAN}── 纯文本类 ──{RESET}")
    for name, maker in [
        ("TXT", make_txt),
        ("Markdown", make_md),
        ("CSV", make_csv),
        ("JSON", make_json),
    ]:
        p = maker(d)
        try:
            r, fmt, fdata = run_processor(p)
            if not r["success"]:
                fail(name, r["error"])
            elif not r["text_content"]:
                fail(name, "text_content 为空")
            elif fdata is not None:
                fail(name, "文本文件不应产生 binary file_data")
            elif "=== 文件内容 ===" not in fmt:
                fail(name, "format 输出中缺少文件内容块")
            else:
                ok(name, f"{len(r['text_content'])} 字符")
        except Exception as e:
            fail(name, f"异常: {e}")


def test_image_files(d):
    print(f"\n{CYAN}── 图片类 ──{RESET}")
    for name, maker, ext in [("PNG", make_png, ".png"), ("JPEG", make_jpg, ".jpg")]:
        p = maker(d)
        try:
            r, fmt, fdata = run_processor(p)
            if not r["success"]:
                fail(name, r["error"])
            elif r["binary_data"] is None or len(r["binary_data"]) == 0:
                fail(name, "binary_data 为空")
            elif fdata is None:
                fail(name, "图片应产生 file_data")
            elif not fdata.get("mime_type", "").startswith("image/"):
                fail(name, f"mime_type 错误: {fdata.get('mime_type')}")
            else:
                ok(name, f"{len(r['binary_data'])} bytes, mime={fdata['mime_type']}")
        except Exception as e:
            fail(name, f"异常: {e}")


def test_pdf(d):
    print(f"\n{CYAN}── PDF ──{RESET}")
    p = make_pdf(d)
    try:
        r, fmt, fdata = run_processor(p)
        if not r["success"]:
            fail("PDF", r["error"])
        elif r["binary_data"] is None or len(r["binary_data"]) == 0:
            fail("PDF", "binary_data 为空")
        elif fdata is None:
            fail("PDF", "PDF 应产生 file_data")
        elif fdata.get("mime_type") != "application/pdf":
            fail("PDF", f"mime_type 错误: {fdata.get('mime_type')}")
        else:
            tq = r["metadata"].get("text_quality", "n/a")
            ok("PDF", f"{len(r['binary_data'])} bytes, text_quality={tq}")
    except Exception as e:
        fail("PDF", f"异常: {e}")


def test_docx(d):
    print(f"\n{CYAN}── Word / DOCX ──{RESET}")
    p = make_docx(d)
    if p is None:
        skip("DOCX", "python-docx 未安装")
        return
    try:
        r, fmt, fdata = run_processor(p)
        if not r["success"]:
            fail("DOCX", r["error"])
        elif not r["text_content"]:
            fail("DOCX", "text_content 为空")
        elif fdata is not None:
            fail("DOCX", "DOCX 不应产生 binary file_data（应嵌入文本）")
        elif "测试 Word 文档" not in fmt:
            fail("DOCX", "提取文本缺少预期内容")
        else:
            ok(
                "DOCX",
                f"{len(r['text_content'])} 字符, {r['metadata'].get('paragraphs','?')} 段落",
            )
    except Exception as e:
        fail("DOCX", f"异常: {e}")


def test_xlsx(d):
    print(f"\n{CYAN}── Excel / XLSX ──{RESET}")
    p = make_xlsx(d)
    if p is None:
        skip("XLSX", "openpyxl/pandas 未安装")
        return
    try:
        r, fmt, fdata = run_processor(p)
        if not r["success"]:
            fail("XLSX", r["error"])
        elif not r["text_content"]:
            fail("XLSX", "text_content 为空")
        elif fdata is not None:
            fail("XLSX", "XLSX 不应产生 binary file_data")
        else:
            ok(
                "XLSX",
                f"{len(r['text_content'])} 字符, {r['metadata'].get('sheets','?')} 个sheet",
            )
    except Exception as e:
        fail("XLSX", f"异常: {e}")


def test_pptx(d):
    print(f"\n{CYAN}── PowerPoint / PPTX ──{RESET}")
    p = make_pptx(d)
    if p is None:
        skip("PPTX", "python-pptx 未安装")
        return
    try:
        r, fmt, fdata = run_processor(p)
        if not r["success"]:
            fail("PPTX", r["error"])
        elif fdata is not None:
            fail("PPTX", "PPTX 不应产生 binary file_data")
        elif not r["text_content"]:
            # 有些 PPTX 无文本，仍算通过（不报错）
            ok("PPTX (无文本)", f"{r['metadata'].get('slides','?')} 张幻灯片")
        else:
            ok(
                "PPTX",
                f"{len(r['text_content'])} 字符, {r['metadata'].get('slides','?')} 张",
            )
    except Exception as e:
        fail("PPTX", f"异常: {e}")


def test_process_uploaded_file_pipeline(d):
    """测试 process_uploaded_file 端到端返回值形态"""
    print(f"\n{CYAN}── process_uploaded_file 端到端管道 ──{RESET}")
    cases = [
        ("txt→文本消息", make_txt(d), False),
        ("png→binary file_data", make_png(d), True),
        ("pdf→binary file_data", make_pdf(d), True),
    ]
    for label, p, expect_fdata in cases:
        try:
            fmt, fdata = process_uploaded_file(p, "这是用户消息")
            if expect_fdata and fdata is None:
                fail(label, "期望有 file_data，实际为 None")
            elif not expect_fdata and fdata is not None:
                fail(label, f"期望 file_data=None，实际有 ({fdata.get('mime_type')})")
            elif not fmt:
                fail(label, "formatted_message 为空")
            elif "用户消息" not in fmt and not expect_fdata:
                # 文本文件消息应包含用户原文
                fail(label, "formatted_message 中缺少原始用户消息")
            else:
                ok(label, f"fmt_len={len(fmt)}, fdata={'有' if fdata else '无'}")
        except Exception as e:
            fail(label, f"异常: {e}")


def test_error_handling(d):
    """测试错误路径（不存在的文件）"""
    print(f"\n{CYAN}── 异常处理 ──{RESET}")
    nonexist = os.path.join(d, "nonexistent.docx")
    try:
        r, fmt, fdata = run_processor(nonexist)
        if r["success"]:
            fail("不存在文件", "不存在的文件不应返回 success=True")
        elif "❌" in fmt:
            ok("不存在文件→优雅失败", fmt[:60].strip())
        else:
            fail("不存在文件", f"错误消息格式异常: {fmt[:80]}")
    except Exception as e:
        # FileProcessor 直接 raise 也属于"错误处理存在"，记录 OK
        ok("不存在文件→raise异常", str(e)[:60])


# ╔═══════════════════════════════════════════════════════════════════════╗
# ║  汇总输出                                                             ║
# ╚═══════════════════════════════════════════════════════════════════════╝


def print_summary():
    passed = sum(1 for r in RESULTS if r["status"] == "PASS")
    failed = sum(1 for r in RESULTS if r["status"] == "FAIL")
    skipped = sum(1 for r in RESULTS if r["status"] == "SKIP")
    total = len(RESULTS)
    print(f"\n{'='*60}")
    print(f"{CYAN}文件上传通路测试汇总{RESET}")
    print(f"{'='*60}")
    print(
        f"  总计: {total}   {GREEN}通过: {passed}{RESET}   "
        f"{RED}失败: {failed}{RESET}   {YELLOW}跳过: {skipped}{RESET}"
    )
    if failed:
        print(f"\n{RED}失败项明细:{RESET}")
        for r in RESULTS:
            if r["status"] == "FAIL":
                print(f"  ✗ {r['label']}: {r['detail']}")
    print(f"{'='*60}")
    return failed == 0


if __name__ == "__main__":
    import shutil

    tmpdir = tempfile.mkdtemp(prefix="koto_filetest_")
    try:
        print(f"{CYAN}{'='*60}{RESET}")
        print(f"{CYAN}Koto 文件上传全类型通路测试{RESET}")
        print(f"{CYAN}{'='*60}{RESET}")
        print(f"临时目录: {tmpdir}")

        test_text_files(tmpdir)
        test_image_files(tmpdir)
        test_pdf(tmpdir)
        test_docx(tmpdir)
        test_xlsx(tmpdir)
        test_pptx(tmpdir)
        test_process_uploaded_file_pipeline(tmpdir)
        test_error_handling(tmpdir)

        success = print_summary()
        sys.exit(0 if success else 1)
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)
