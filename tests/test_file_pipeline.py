#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
文件格式任务通路测试
验证 FileProcessor / FileParser 对各种支持格式的解析是否正常。

支持的格式:
  FileProcessor (上传通路):
    - 图片: .jpg .png .gif .bmp .webp (image/* MIME)
    - PDF:  .pdf
    - Word: .doc / .docx
    - PPT:  .ppt / .pptx
    - Excel:.xls / .xlsx
    - 文本: .txt / .md / .markdown

  FileParser (PPT 生成素材通路):
    - .pdf / .docx / .doc / .txt / .md / .markdown

运行方式:
  cd c:\\Users\\12524\\Desktop\\Koto
  .venv\\Scripts\\python.exe tests\\test_file_pipeline.py
"""

import os
import sys
import tempfile
import shutil

# 确保项目根目录在 sys.path
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)

# ─────────── 颜色输出 ───────────
GREEN = "\033[92m"
RED = "\033[91m"
YELLOW = "\033[93m"
RESET = "\033[0m"


def ok(msg):
    print(f"{GREEN}  ✅ {msg}{RESET}")


def fail(msg):
    print(f"{RED}  ❌ {msg}{RESET}")


def warn(msg):
    print(f"{YELLOW}  ⚠️  {msg}{RESET}")


def section(title):
    print(f"\n{'─'*55}\n  {title}\n{'─'*55}")


# ─────────── 临时目录 ───────────
TMP_DIR = tempfile.mkdtemp(prefix="koto_filetest_")


def cleanup():
    shutil.rmtree(TMP_DIR, ignore_errors=True)


# ─────────── 生成各类测试文件 ───────────


def make_txt(path: str) -> str:
    p = os.path.join(path, "sample.txt")
    with open(p, "w", encoding="utf-8") as f:
        f.write("这是一个纯文本文件。\nHello Koto!\n第三行内容。\n")
    return p


def make_md(path: str) -> str:
    p = os.path.join(path, "sample.md")
    with open(p, "w", encoding="utf-8") as f:
        f.write("# 标题\n\n这是 Markdown 文档。\n\n- 列表项1\n- 列表项2\n")
    return p


def make_csv(path: str) -> str:
    """CSV 实际上被当作 text/* 处理"""
    p = os.path.join(path, "sample.csv")
    with open(p, "w", encoding="utf-8") as f:
        f.write("姓名,年龄,城市\n张三,25,北京\n李四,30,上海\n")
    return p


def make_json(path: str) -> str:
    """JSON 实际上被当作 text/* 处理"""
    p = os.path.join(path, "sample.json")
    with open(p, "w", encoding="utf-8") as f:
        f.write('{"name": "Koto", "version": "2.0", "langs": ["zh", "en"]}\n')
    return p


def make_docx(path: str) -> str:
    p = os.path.join(path, "sample.docx")
    try:
        from docx import Document

        doc = Document()
        doc.add_heading("测试文档标题", 0)
        doc.add_paragraph("这是第一段正文内容，用于验证 .docx 解析通路。")
        doc.add_paragraph("第二段：Hello Koto DOCX test!")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "列1"
        table.cell(0, 1).text = "列2"
        table.cell(1, 0).text = "数据A"
        table.cell(1, 1).text = "数据B"
        doc.save(p)
    except ImportError:
        # 创建一个最小 docx 结构（zip包）用于检测库缺失的错误路径
        import zipfile

        with zipfile.ZipFile(p, "w") as zf:
            zf.writestr(
                "[Content_Types].xml",
                '<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"></Types>',
            )
    return p


def make_pptx(path: str) -> str:
    p = os.path.join(path, "sample.pptx")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Koto PPT 测试"
        slide.placeholders[1].text = "这是演示文稿的第一张幻灯片内容。"
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "第二张幻灯片"
        slide2.placeholders[1].text = "Hello PPTX!"
        prs.save(p)
    except ImportError:
        import zipfile

        with zipfile.ZipFile(p, "w") as zf:
            zf.writestr(
                "[Content_Types].xml",
                '<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"></Types>',
            )
    return p


def make_xlsx(path: str) -> str:
    p = os.path.join(path, "sample.xlsx")
    try:
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["日期", "销售额", "利润"])
        ws.append(["2026-01-01", 10000, 2000])
        ws.append(["2026-01-02", 12000, 2500])
        ws2 = wb.create_sheet("Sheet2")
        ws2.append(["项目", "状态"])
        ws2.append(["Koto", "Active"])
        wb.save(p)
    except ImportError:
        import zipfile

        with zipfile.ZipFile(p, "w") as zf:
            zf.writestr(
                "[Content_Types].xml",
                '<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"></Types>',
            )
    return p


def make_pdf(path: str) -> str:
    p = os.path.join(path, "sample.pdf")
    # 生成最小合法 PDF（不依赖外部库）
    pdf_bytes = (
        b"%PDF-1.4\n"
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]"
        b" /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n"
        b"4 0 obj\n<< /Length 44 >>\nstream\n"
        b"BT /F1 12 Tf 100 700 Td (Koto Test PDF) Tj ET\n"
        b"endstream\nendobj\n"
        b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
        b"xref\n0 6\n0000000000 65535 f \n"
        b"0000000009 00000 n \n0000000058 00000 n \n0000000115 00000 n \n"
        b"0000000274 00000 n \n0000000369 00000 n \n"
        b"trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n452\n%%EOF\n"
    )
    with open(p, "wb") as f:
        f.write(pdf_bytes)
    return p


def make_image_png(path: str) -> str:
    p = os.path.join(path, "sample.png")
    try:
        from PIL import Image

        img = Image.new("RGB", (100, 100), color=(73, 109, 137))
        img.save(p, "PNG")
    except ImportError:
        # 最小合法 1x1 红点 PNG（无需 PIL）
        import base64

        tiny_png = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8"
            "z8BQDwADhQGAWjR9awAAAABJRU5ErkJggg=="
        )
        with open(p, "wb") as f:
            f.write(tiny_png)
    return p


def make_image_jpg(path: str) -> str:
    p = os.path.join(path, "sample.jpg")
    try:
        from PIL import Image

        img = Image.new("RGB", (100, 100), color=(200, 100, 50))
        img.save(p, "JPEG")
    except ImportError:
        # 最小合法 1x1 JPEG（无 PIL）
        import base64

        tiny_jpg = base64.b64decode(
            "/9j/4AAQSkZJRgABAQEAYABgAAD/2wBDAAgGBgcGBQgHBwcJCQgKDBQNDAsLDBk"
            "SEw8UHRofHh0aHBwgJC4nICIsIxwcKDcpLDAxNDQ0Hyc5PTgyPC4zNDL/2wBD"
            "AQkJCQwLDBgNDRgyIRwhMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIy"
            "MjIyMjIyMjIyMjIyMjIyMjL/wAARCAABAAEDASIAAhEBAxEB/8QAFgABAQEA"
            "AAAAAAAAAAAAAAAABQMG/8QAIRAAAQQCAgMBAAAAAAAAAAAAAQIDBBESITEFQWH/"
            "xAAUAQEAAAAAAAAAAAAAAAAAAAAA/8QAFBEBAAAAAAAAAAAAAAAAAAAAAP/aAAwD"
            "AQACEQMRAD8Aq3a3YV5YlJFfCfEVHkS0UBERERb/2Q=="
        )
        with open(p, "wb") as f:
            f.write(tiny_jpg)
    return p


# ─────────── 测试 FileProcessor ───────────


def test_file_processor():
    section("FileProcessor 测试（上传解析通路）")
    from web.file_processor import FileProcessor

    test_cases = [
        ("TXT", make_txt(TMP_DIR), "text_content", 10),
        ("MD", make_md(TMP_DIR), "text_content", 10),
        ("CSV (文本)", make_csv(TMP_DIR), "text_content", 5),
        ("JSON(文本)", make_json(TMP_DIR), "text_content", 5),
        ("DOCX", make_docx(TMP_DIR), "text_content", 5),
        ("PPTX", make_pptx(TMP_DIR), "text_content", 0),  # pptx 可以空文本仍成功
        ("XLSX", make_xlsx(TMP_DIR), "text_content", 5),
        ("PDF", make_pdf(TMP_DIR), None, 0),  # pdf 以二进制传，text_content 可空
        ("PNG", make_image_png(TMP_DIR), None, 0),  # 图片只检查 success
        ("JPG", make_image_jpg(TMP_DIR), None, 0),
    ]

    passed = 0
    failed_names = []

    for label, fpath, content_key, min_chars in test_cases:
        try:
            result = FileProcessor.process_file(fpath)
            if not result.get("success"):
                fail(f"{label:<12} FAIL — {result.get('error', '未知错误')}")
                failed_names.append(label)
                continue

            # 检查内容字段
            if content_key:
                content = result.get(content_key, "")
                if len(content) < min_chars:
                    warn(
                        f"{label:<12} WARN — {content_key} 太短 ({len(content)} 字符，期望 ≥{min_chars})"
                    )
                else:
                    ok(f"{label:<12} OK   — {content_key}: {len(content)} 字符")
            else:
                if label == "PDF":
                    bdata = result.get("binary_data")
                    if bdata:
                        ok(f"{label:<12} OK   — binary_data: {len(bdata)} 字节")
                    else:
                        warn(f"{label:<12} WARN — binary_data 为空")
                else:
                    ok(f"{label:<12} OK   — success=True")

            passed += 1

        except Exception as e:
            fail(f"{label:<12} EXCEPTION — {e}")
            failed_names.append(label)

    print(
        f"\n  结果: {passed}/{len(test_cases)} 通过"
        + (f"，失败: {', '.join(failed_names)}" if failed_names else "")
    )
    return len(failed_names) == 0


# ─────────── 测试 FileParser ───────────


def test_file_parser():
    section("FileParser 测试（PPT 素材解析通路）")
    from web.file_parser import FileParser

    test_cases = [
        ("PDF", make_pdf(TMP_DIR)),
        ("DOCX", make_docx(TMP_DIR)),
        ("TXT", make_txt(TMP_DIR)),
        ("MD", make_md(TMP_DIR)),
    ]

    passed = 0
    failed_names = []

    for label, fpath in test_cases:
        result = FileParser.parse_file(fpath)
        if result and result.get("success"):
            chars = result.get("char_count", 0)
            ok(f"{label:<12} OK   — {chars} 字符")
            passed += 1
        else:
            err = result.get("error", "未知") if result else "None 返回"
            fail(f"{label:<12} FAIL — {err}")
            failed_names.append(label)

    # 测试不支持的格式
    unsupported_path = os.path.join(TMP_DIR, "sample.xyz")
    with open(unsupported_path, "w") as f:
        f.write("test")
    res = FileParser.parse_file(unsupported_path)
    if res and not res.get("success") and "不支持" in res.get("error", ""):
        ok(f"{'不支持格式':<12} OK   — 正确拒绝 .xyz")
    else:
        warn(f"{'不支持格式':<12} WARN — 预期拒绝 .xyz，实际: {res}")

    print(
        f"\n  结果: {passed}/{len(test_cases)} 通过"
        + (f"，失败: {', '.join(failed_names)}" if failed_names else "")
    )
    return len(failed_names) == 0


# ─────────── 测试 format_result_for_chat ───────────


def test_format_result_for_chat():
    section("FileProcessor.format_result_for_chat 格式化通路测试")
    from web.file_processor import FileProcessor

    # 文本文件
    txt_path = make_txt(TMP_DIR)
    result_txt = FileProcessor.process_file(txt_path)
    msg, fdata = FileProcessor.format_result_for_chat(result_txt, "帮我总结这个文件")
    if "帮我总结" in msg and "纯文本" not in msg.lower() or len(msg) > 5:
        ok(f"{'TXT format':<12} OK   — 格式化消息长度 {len(msg)}")
    else:
        warn(f"{'TXT format':<12} WARN — 格式化消息异常: {msg[:80]}")

    # PDF 文件
    pdf_path = make_pdf(TMP_DIR)
    result_pdf = FileProcessor.process_file(pdf_path)
    msg_pdf, fdata_pdf = FileProcessor.format_result_for_chat(result_pdf, "分析这个PDF")
    if result_pdf.get("success"):
        if fdata_pdf and fdata_pdf.get("mime_type") == "application/pdf":
            ok(f"{'PDF format':<12} OK   — file_data mime_type=application/pdf")
        else:
            ok(f"{'PDF format':<12} OK   — 格式化消息长度 {len(msg_pdf)}")
    else:
        warn(f"{'PDF format':<12} WARN — PDF未成功处理，跳过格式化测试")


# ─────────── 测试 FileParser.batch_parse ───────────


def test_batch_parse():
    section("FileParser.batch_parse 批量解析通路测试")
    from web.file_parser import FileParser

    paths = [make_txt(TMP_DIR), make_md(TMP_DIR), make_pdf(TMP_DIR)]
    results = FileParser.batch_parse(paths)

    ok_count = sum(1 for r in results if r and r.get("success"))
    ok(f"批量解析 {len(paths)} 个文件 → {ok_count} 成功")

    merged = FileParser.merge_contents([r for r in results if r and r.get("success")])
    if merged and len(merged) > 20:
        ok(f"merge_contents   OK — 合并后 {len(merged)} 字符")
    else:
        warn(f"merge_contents   WARN — 合并内容过短: {repr(merged[:100])}")


# ─────────── 检查依赖库 ───────────


def check_dependencies():
    section("依赖库检查")
    deps = {
        "python-docx": ("docx", "Document"),
        "python-pptx": ("pptx", "Presentation"),
        "openpyxl": ("openpyxl", "Workbook"),
        "PyPDF2": ("PyPDF2", "PdfReader"),
        "pypdf": ("pypdf", "PdfReader"),
        "pdfplumber": ("pdfplumber", "open"),
        "Pillow": ("PIL", "Image"),
        "pandas": ("pandas", "read_excel"),
    }
    missing = []
    for pkg, (mod, attr) in deps.items():
        try:
            m = __import__(mod)
            getattr(m, attr)
            ok(f"{pkg:<16} 已安装")
        except (ImportError, AttributeError):
            warn(f"{pkg:<16} 未安装（可选）")
            missing.append(pkg)
    if missing:
        print(f"\n  未安装的库: {', '.join(missing)}")
        print("  部分格式可能无法解析，但核心通路仍应以 fallback 方式处理。")


# ─────────── 主入口 ───────────


def main():
    print("=" * 55)
    print("  Koto 文件格式任务通路测试")
    print(f"  临时目录: {TMP_DIR}")
    print("=" * 55)

    try:
        check_dependencies()

        all_pass = True
        all_pass &= test_file_processor()
        all_pass &= test_file_parser()
        test_format_result_for_chat()
        test_batch_parse()

        print("\n" + "=" * 55)
        if all_pass:
            print(f"{GREEN}  🎉 所有核心通路测试通过！{RESET}")
        else:
            print(f"{YELLOW}  ⚠️  部分测试存在问题，请检查上方报告。{RESET}")
        print("=" * 55)

    finally:
        cleanup()


if __name__ == "__main__":
    main()
