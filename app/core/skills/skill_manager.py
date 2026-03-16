# -*- coding: utf-8 -*-
"""
🎯 Koto Skills Manager（v2 — 原子化 Schema 升级）

可插拔的 Prompt 技能系统。
每个 Skill 现在由 SkillDefinition（原子化标准 Schema）描述，
支持 MCP Tool 导出、IO 变量约束、输出验收规格。

向后兼容：旧版 dict 格式的 BUILTIN_SKILLS 会自动升级为 SkillDefinition，
原有的 inject_into_prompt / list_skills / set_enabled 等 API 不变。

新增 API:
    SkillManager.get_definition(skill_id)   → SkillDefinition
    SkillManager.list_mcp_tools()           → List[MCP Tool dict]
    SkillManager.register_custom(skill_def) → 注册自定义 Skill
    SkillManager.load_custom_skills_dir()   → 从 config/skills/ 目录加载 JSON Skill 文件

用法:
    from app.core.skills.skill_manager import SkillManager

    enhanced = SkillManager.inject_into_prompt(base_instruction, task_type="CHAT")
    skills   = SkillManager.list_skills()
    SkillManager.set_enabled("concise_mode", True)

    # 新：导出 MCP 工具列表
    mcp_tools = SkillManager.list_mcp_tools()
"""

import json
import logging
import os
from pathlib import Path
from typing import Dict, List, Optional

from app.core.skills.skill_schema import OutputSpec, SkillDefinition

logger = logging.getLogger(__name__)

# ── 内置技能定义 ────────────────────────────────────────────────────────────────
# 每个技能字段说明：
#   id           : 唯一标识符（snake_case）
#   name         : 中文名称
#   icon         : emoji 图标
#   category     : 分类 behavior | style | domain
#   description  : 简短描述（显示在 UI）
#   task_types   : 生效的任务类型列表；空列表 = 所有任务类型
#   prompt       : 注入到 system_instruction 的文本片段
#   enabled      : 默认启用状态
# ───────────────────────────────────────────────────────────────────────────────
BUILTIN_SKILLS: List[Dict] = [
    # ── 行为类 ────────────────────────────────────────────────────────────────
    {
        "id": "step_by_step",
        "name": "步骤化输出",
        "icon": "🪜",
        "category": "behavior",
        "skill_nature": "model_hint",
        "description": "将回答拆解为带编号的清晰步骤，自动适配操作/排查/决策等不同场景",
        "intent_description": "用户需要操作指南、安装配置步骤、故障排查流程、学习路径或任何有先后顺序的内容",
        "task_types": [],
        "prompt": (
            "\n\n## 🪜 行为要求：步骤化输出"
            "\n\n**第一步：判断是否适合步骤化（必须检查）**"
            "\n- ✅ 适用：操作流程、安装配置、故障排查、学习路径、多阶段决策"
            "\n- ❌ 不适用：单一事实、词汇定义、是/否判断、简单数值计算、情感回应"
            "\n  → 不适用时，直接给出结论，不要强行加编号。"
            "\n\n**操作 / 教程类问题 — 标准格式**"
            "\n```"
            "\n📋 【前提条件】（无则省略）"
            "\n- 所需工具 / 权限 / 基础知识"
            "\n"
            "\n步骤 1：[简洁动词短语，说明目标]"
            "\n  ▸ 操作：（每步限 1-3 个动作；内容复杂时用子步骤 1.1 / 1.2 展开）"
            "\n  ▸ 预期结果：（完成操作后应看到或得到什么）"
            "\n  ⚠️ 注意：（仅在有易错点或风险时写，否则省略）"
            "\n"
            "\n步骤 2：…"
            "\n（依此格式直到流程完成）"
            "\n"
            "\n✅ 完成验证：[如何确认全部步骤执行成功的具体指标]"
            "\n🔁 常见失败处理：[最高频的失败场景及对应排查方向]"
            "\n```"
            "\n\n**故障排查 / 调试类问题 — 四步诊断格式**"
            "\n（当用户问题含「报错/异常/失败/不正常/bug/无法」等关键词时优先用此格式）"
            "\n```"
            "\n步骤 1：🔍 定位问题"
            "\n  ▸ 从错误信息/现象出发，排除明显原因，缩小范围"
            "\n步骤 2：🧪 验证假设"
            "\n  ▸ 给出能快速确认根因的最小测试命令 / 操作"
            "\n步骤 3：🔧 应用修复"
            "\n  ▸ 针对根本原因的最小改动方案（附代码块如适用）"
            "\n步骤 4：✅ 验证修复"
            "\n  ▸ 执行后应看到的成功信号"
            "\n```"
            "\n\n**格式规则**"
            "\n- 步骤编号用阿拉伯数字（1、2、3），禁止「首先/其次/然后/最后」等过渡词"
            "\n- 步骤描述用主动语态动词开头（「点击 X」而非「X 需要被点击」）"
            "\n- 代码独占代码块，不嵌在普通段落行内"
            "\n- 理想步骤数量 3-6 步；超过 7 步时将流程拆为「阶段」再细分"
            "\n- 全部完成后不要再写「总结一下以上步骤…」重复内容"
        ),
        "enabled": False,
    },
    {
        "id": "strict_mode",
        "name": "严谨模式",
        "icon": "🔬",
        "category": "behavior",
        "skill_nature": "model_hint",
        "description": "要求更严格的推理：引用来源、标明不确定性、避免模糊结论",
        "intent_description": "用户需要可靠的事实或分析、担心信息过时、处理有争议的话题、或明确要求严谨/有据可依的回答时",
        "task_types": ["CHAT", "RESEARCH"],
        "prompt": (
            "\n\n## 🔬 行为要求：严谨模式\n"
            "- 区分「已知事实」和「推测/估计」，后者用「可能/据报道」等标注\n"
            "- 涉及时间敏感内容时，提醒信息可能已过时\n"
            "- 避免过度自信：不确定时明确说「我不确定」\n"
            "- 给出多角度分析时，注明前提假设是什么"
        ),
        "enabled": False,
    },
    {
        "id": "concise_mode",
        "name": "精简模式",
        "icon": "⚡",
        "category": "behavior",
        "skill_nature": "model_hint",
        "description": "强制简短回答：只给核心结论和最关键步骤，适合快速查询场景",
        "intent_description": "用户需要简短、直接的回答，例如快速查询、简单事实确认",
        "task_types": ["CHAT"],
        "conflict_with": ["step_by_step", "teaching_mode", "proactive_suggestions"],
        "priority": 90,  # 最高优先级，覆盖其他行为
        "prompt": (
            "\n\n## ⚡ 行为要求：精简模式（最高优先级）\n"
            "- 回答控制在 150 字以内\n"
            "- 只给出「结论 + 最小可行步骤」，省略背景和解释\n"
            "- 禁止使用「首先…其次…最后…」等冗长过渡语\n"
            "- 若用户明确要求详细解释，则忽略此限制"
        ),
        "enabled": False,
    },
    {
        "id": "teaching_mode",
        "name": "教学模式",
        "icon": "🎓",
        "category": "behavior",
        "skill_nature": "model_hint",
        "description": "苏格拉底式引导 + 类比 + 逐层展开，帮助用户真正理解而非死记答案",
        "intent_description": "用户说「教我」「帮我理解」「解释一下」「我是新手」或想从零开始系统学习某个概念时",
        "task_types": ["CHAT", "RESEARCH"],
        "prompt": (
            "\n\n## 🎓 行为要求：教学模式\n"
            "- 用「类比」解释抽象概念（「就像…」「可以把它理解为…」）\n"
            "- 从基础前提出发，逐层深入，不跳过中间推导\n"
            "- 在适当节点加入「💡 为什么这样做？」的思考引导\n"
            "- 结尾附「📝 关键要点」小结，帮助用户归纳记忆"
        ),
        "enabled": False,
    },
    {
        "id": "proactive_suggestions",
        "name": "主动建议",
        "icon": "💡",
        "category": "behavior",
        "skill_nature": "model_hint",
        "description": "回答后主动给出相关的延伸建议、潜在问题或下一步操作",
        "intent_description": "用户在探索某个主题、刚完成一项任务、或需要进一步指引时，Koto 主动推荐下一步行动",
        "task_types": [],
        "prompt": (
            "\n\n## 💡 行为要求：主动建议\n"
            "- 完成主要回答后，用「💡 延伸建议」段落提出 1-3 个相关建议\n"
            "- 如果检测到用户可能遇到的陷阱，用「⚠️ 注意」提前提示\n"
            "- 适时建议下一步操作：「您可能还想要…」"
        ),
        "enabled": False,
    },
    # ── 风格类 ────────────────────────────────────────────────────────────────
    {
        "id": "professional_tone",
        "name": "专业语气",
        "icon": "👔",
        "category": "style",
        "skill_nature": "model_hint",
        "description": "使用正式的商务/学术语气，适合需要对外输出的报告、邮件、文档",
        "intent_description": "用户要写对外报告、商务邮件、官方声明、论文摘要或任何需要正式书面语气的内容时",
        "task_types": ["CHAT", "FILE_GEN", "RESEARCH"],
        "prompt": (
            "\n\n## 👔 风格要求：专业语气\n"
            "- 使用正式书面语，避免口语化表达和网络用语\n"
            "- 称谓使用「您」；指代用户需求时用「贵方/您的需求」\n"
            "- 数字和单位规范书写（如「3 个」而非「三个」）\n"
            "- 结构清晰：每段只表达一个主要观点"
        ),
        "enabled": False,
    },
    {
        "id": "creative_writing",
        "name": "创意写作",
        "icon": "✨",
        "category": "style",
        "skill_nature": "model_hint",
        "description": "使用生动、富有感染力的语言，适合文案、故事、创意内容生成",
        "task_types": ["CHAT", "FILE_GEN"],
        "prompt": (
            "\n\n## ✨ 风格要求：创意写作\n"
            "- 使用具体、感官性的描写代替抽象表述（「金黄色的光」而非「美丽的光」）\n"
            "- 善用比喻、拟人、排比等修辞手法增加文字感染力\n"
            "- 句子长短交错，避免单调节奏\n"
            "- 鼓励独特角度和出人意料的表达"
        ),
        "enabled": False,
    },
    {
        "id": "emoji_assist",
        "name": "Emoji 辅助",
        "icon": "🎨",
        "category": "style",
        "skill_nature": "model_hint",
        "description": "适当使用 emoji 增强视觉层次感和可读性",
        "task_types": [],
        "prompt": (
            "\n\n## 🎨 风格要求：Emoji 辅助\n"
            "- 各段标题/分类用 emoji 标注以增强视觉区分\n"
            "- 重要警告用 ⚠️，成功结果用 ✅，关键步骤用 👉\n"
            "- 全文 emoji 密度适中（每段 1-2 个），不要每句话都用"
        ),
        "enabled": False,
    },
    {
        "id": "bilingual",
        "name": "双语标注",
        "icon": "🌐",
        "category": "style",
        "skill_nature": "model_hint",
        "description": "对专业术语和概念提供中英双语标注，适合学习技术内容",
        "task_types": ["CHAT", "RESEARCH", "CODER"],
        "prompt": (
            "\n\n## 🌐 风格要求：双语标注\n"
            "- 首次出现的专业术语在中文后括注英文原文，如「机器学习 (Machine Learning)」\n"
            "- 代码相关的概念保留英文名称\n"
            "- 不需要对每个普通词汇都标注，只针对技术和专有名词"
        ),
        "enabled": False,
    },
    # ── 领域类 ────────────────────────────────────────────────────────────────
    {
        "id": "code_best_practices",
        "name": "代码最佳实践",
        "icon": "🛡️",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "编写代码时强调可读性、注释、异常处理和测试，提升代码质量",
        "task_types": ["CODER", "CHAT"],
        "prompt": (
            "\n\n## 🛡️ 领域要求：代码最佳实践\n"
            "- 每个函数/类必须有中文文档字符串（docstring）\n"
            "- 关键逻辑加行内注释，说明「为什么」而不仅是「做什么」\n"
            "- 对外部调用、文件IO、网络请求一律使用 try/except 包裹\n"
            "- 变量和函数命名使用 snake_case，类名使用 PascalCase\n"
            "- 输出代码后，额外说明「可能的改进」或「注意事项」"
        ),
        "enabled": False,
    },
    {
        "id": "security_aware",
        "name": "安全意识",
        "icon": "🔒",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "在代码和系统操作建议中主动指出安全风险和防护措施",
        "task_types": ["CODER", "SYSTEM", "AGENT", "CHAT"],
        "prompt": (
            "\n\n## 🔒 领域要求：安全意识\n"
            "- 涉及文件删除/格式化/权限修改等高风险操作时，主动添加确认提示\n"
            "- 代码中如有硬编码密码/密钥，必须提示替换为环境变量\n"
            "- SQL 拼接、文件路径拼接处提醒注入风险\n"
            "- 网络请求建议始终校验 HTTPS 和证书"
        ),
        "enabled": False,
    },
    {
        "id": "data_analysis",
        "name": "数据分析视角",
        "icon": "📊",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "回答时多从数据和量化角度出发，提供可验证的建议",
        "task_types": ["CHAT", "RESEARCH"],
        "prompt": (
            "\n\n## 📊 领域要求：数据分析视角\n"
            "- 尽量用数据和指标支撑观点，而不仅是定性描述\n"
            "- 建议可测量的方法：「通过监控 X 指标来判断」\n"
            "- 对比分析时使用表格格式呈现\n"
            "- 推荐 Python/Excel 等工具时，给出具体的实现思路\n\n"
            "## 🔧 Koto 实际可用的数据分析工具（如实告知用户）\n"
            "当用户询问 Koto 能做什么数据分析时，只介绍以下已实现的能力，不夸大：\n"
            "1. **加载数据文件**（load_data）：读取 CSV / Excel / JSON，返回行列数和预览\n"
            "2. **统计摘要**（describe_data）：均值、中位数、标准差、缺失值统计、分类列频率\n"
            "3. **数据查询**（query_data）：在 DataFrame 上运行 pandas 表达式过滤/聚合\n"
            "4. **探索性问题建议**（suggest_questions）：根据数据结构自动生成分析方向\n"
            "5. **趋势分析**（analyze_trends）：时序数据趋势、环比增长率、峰谷检测、异常识别\n"
            "6. **数据保存**（save_data）：将结果导出为 CSV / Excel / JSON\n"
            "7. **图表图片解读**（analyze_chart_image）：上传图表截图，Gemini 多模态视觉解析\n"
            "8. **学习资源推荐**（learning_guide 技能）：通过 web_search 搜索最新学习路径\n\n"
            "**真正尚未具备的能力（不要在回复中声称拥有）：**\n"
            "- 直接从网页 URL 读取在线图表（需用户先下载到本地）\n"
            "- 自动绘制可视化图表并展示（可生成绘图代码，但不能直接渲染图形）"
        ),
        "enabled": True,
    },
    # ── 领域增强 ──────────────────────────────────────────────────────────────
    {
        "id": "writing_assistant",
        "name": "写作助手",
        "icon": "✍️",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "提升文字表达质量：润色语言、优化段落结构、保持风格一致",
        "intent_description": "用户要写文章、报告、邮件、文案或需要修改润色文字时",
        "task_types": ["CHAT", "FILE_GEN"],
        "priority": 60,
        "prompt": (
            "\n\n## ✍️ 领域要求：写作助手\n"
            "- 每次输出前默默检查：逻辑是否连贯、表达是否准确、有无冗余\n"
            "- 段落开头避免「首先/其次/另外」等程式化连接词\n"
            "- 对同一意思给出至少一种更简练或更有力的替代表达\n"
            "- 指出可能的歧义或不清晰之处，建议具体改法\n"
            "- 保持用户原文的语气和风格，不做风格切换"
        ),
        "enabled": False,
    },
    {
        "id": "debate_mode",
        "name": "批判思维",
        "icon": "⚖️",
        "category": "behavior",
        "skill_nature": "model_hint",
        "description": "主动挑战假设和结论，从多角度分析，识别逻辑漏洞和可能的反驳",
        "intent_description": "用户在思考复杂决策、评估方案优劣、或需要检验观点可靠性时",
        "task_types": ["CHAT", "RESEARCH"],
        "conflict_with": ["proactive_suggestions"],
        "priority": 65,
        "prompt": (
            "\n\n## ⚖️ 行为要求：批判思维\n"
            "- 分析问题时，主动列出「支持」和「反对」两侧论据\n"
            "- 识别并明确指出输入中存在的假设（显性和隐性）\n"
            "- 对结论标注置信度：确定/可能/存疑\n"
            "- 指出「还需要哪些信息才能做出更可靠的判断」\n"
            "- 避免立场倾斜：即使用户倾向某一侧，也要呈现完整对立面"
        ),
        "enabled": False,
    },
    {
        "id": "research_depth",
        "name": "深度研究",
        "icon": "🔭",
        "category": "behavior",
        "skill_nature": "model_hint",
        "description": "提供有深度的研究性回答：多维度分析、原因链条、历史背景和未来影响",
        "intent_description": "用户在做研究、深入了解某个话题、或准备报告/论文内容时",
        "task_types": ["RESEARCH", "CHAT"],
        "conflict_with": ["concise_mode"],
        "priority": 55,
        "prompt": (
            "\n\n## 🔭 行为要求：深度研究\n"
            "- 先给出「结论摘要」，再分层展开论据\n"
            "- 追溯原因链：「A 导致 B，因为 C，背后是 D」\n"
            "- 主动给出历史背景或类似案例以增加参考维度\n"
            "- 区分「已有广泛共识」和「存在争议」的内容\n"
            "- 给出延伸阅读方向：「若想深入了解，可查阅…」"
        ),
        "enabled": False,
    },
    {
        "id": "task_planner",
        "name": "任务规划",
        "icon": "🗓️",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "将复杂目标分解为具体可执行的任务列表，评估优先级和依赖关系",
        "intent_description": "用户在规划项目、拆解复杂任务、制定行动计划或整理待办事项时",
        "task_types": ["CHAT", "FILE_GEN"],
        "priority": 60,
        "prompt": (
            "\n\n## 🗓️ 行为要求：任务规划\n"
            "- 将用户的目标分解为 3-7 个具体、可执行的子任务\n"
            "- 每个子任务标注：预计耗时 / 优先级（高/中/低）/ 前置依赖\n"
            "- 识别关键路径：哪些任务必须先完成\n"
            "- 用 Markdown 复选框格式（- [ ] 任务名）便于直接使用\n"
            "- 最后提醒可能的风险点和应对建议"
        ),
        "enabled": False,
    },
    # ── 专项文档批注技能 ──────────────────────────────────────────────────────
    {
        "id": "annotate_academic",
        "name": "学术批注",
        "icon": "📚",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "对学术论文进行批注润色：修正语言、增强逻辑连贯性、符合期刊投稿规范",
        "intent_description": "用户需要对学术论文、研究报告进行语言润色、批注修改、投稿前精修",
        "task_types": ["DOC_ANNOTATE"],
        "priority": 70,
        "prompt": (
            "\n\n## 📚 领域要求：学术批注\n"
            "- 采用「原文 → 修改建议 → 修改理由」三段式逐段批注\n"
            "- 重点关注：措辞精确性、被动语态使用、逻辑连接词、数据引用格式（APA/GB/MLA）\n"
            "- 标注可能引发审稿人质疑的逻辑断层或过度主张\n"
            "- 区分「必须修改」（❌）和「建议优化」（⚠️）两个层级\n"
            "- 最终附「总体评估」：语言水平评分、逻辑严密性评分、是否符合期刊投稿规范"
        ),
        "enabled": False,
    },
    {
        "id": "annotate_business",
        "name": "商务批注",
        "icon": "💼",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "对商业报告、合同、邮件进行专业批注：语气规范化、结构清晰化、歧义消除",
        "intent_description": "用户需要修改商务报告、合同文本、商务邮件或内部公文",
        "task_types": ["DOC_ANNOTATE", "FILE_GEN"],
        "priority": 65,
        "prompt": (
            "\n\n## 💼 领域要求：商务批注\n"
            "- 聚焦「语气是否专业」「结构是否清晰」「表达是否精准」三个维度\n"
            "- 标注过于口语化、产生歧义或可能引发误解的措辞，给出替代表达\n"
            "- 检查文档逻辑结构：开头目的是否明确、中间论据是否支撑结论、结尾行动项是否清晰\n"
            "- 对关键数据或承诺用「🔍 建议核实」标注，防止后期争议\n"
            "- 输出格式：批注表格（位置 | 原文 | 问题 | 建议） + 下方附修订版全文"
        ),
        "enabled": False,
    },
    {
        "id": "annotate_code_review",
        "name": "代码审查批注",
        "icon": "🔍",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "对代码文件进行专业审查批注：安全性、可读性、最佳实践、潜在 bug 分级标注",
        "intent_description": "用户需要 code review、代码批注、检查代码质量、找出 bug",
        "task_types": ["DOC_ANNOTATE", "CODER"],
        "priority": 75,
        "conflict_with": ["code_best_practices"],
        "executor_tools": ["read_file_snippet", "find_file", "list_directory"],
        "plan_template": [
            "读取目标代码文件或代码块",
            "分析安全性（注入风险/硬编码密钥/权限漏洞）",
            "检查可读性与最佳实践合规性",
            "识别潜在 bug 和边界条件缺失",
            "生成按严重等级分类的批注报告（🔴严重 / 🟡警告 / 🔵建议）",
            "输出整体代码质量评分与 Top 3 优先修复项",
        ],
        "prompt": (
            "\n\n## 🔍 领域要求：代码审查批注\n"
            "- 按严重等级分类：🔴 严重（安全漏洞/逻辑错误）、🟡 警告（可读性差/违反规范）、🔵 建议（优化空间）\n"
            "- 必查项：硬编码敏感信息、SQL/命令注入风险、未处理异常、资源未关闭、边界条件缺失\n"
            "- 每条批注格式：「第 N 行 / 函数名 → [等级] 问题描述 → 建议修改方式」\n"
            "- 识别重复代码块，标注「建议抽取为工具函数」\n"
            "- 最终给出「整体代码质量评分」(1-10) 和「Top 3 优先修复项」"
        ),
        "enabled": False,
    },
    {
        "id": "annotate_translation",
        "name": "翻译质检批注",
        "icon": "🌐",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "对翻译文本进行质检批注：忠实度、流畅性、术语一致性对照审查",
        "intent_description": "用户需要检查翻译质量、对照原文批注、纠正译文错误或术语不一致",
        "task_types": ["DOC_ANNOTATE"],
        "priority": 70,
        "prompt": (
            "\n\n## 🌐 领域要求：翻译质检批注\n"
            "- 采用「原文 | 现有译文 | 问题类型 | 建议译文」四列对照格式\n"
            "- 错误类型分类：漏译（MT）、误译（MIS）、术语不一致（TERM）、语法错误（GRAM）、不自然表达（NAT）\n"
            "- 专有名词和术语首次出现后建立「术语表」，后续逐一检查一致性\n"
            "- 区分「硬错误」（影响理解，必须修改）和「软错误」（风格偏好，建议修改）\n"
            "- 最终给出：忠实度评分 + 流畅度评分 + 术语一致性评分（各 1-5 分）"
        ),
        "enabled": False,
    },
    # ── 专项调试技能 ──────────────────────────────────────────────────────────
    {
        "id": "debug_python",
        "name": "Python 调试",
        "icon": "🐍",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "专项 Python 调试：精准定位 traceback、分析根因、给出最小复现和修复方案",
        "intent_description": "用户遇到 Python 报错、traceback、AttributeError、ImportError 等 bug 需要调试",
        "task_types": ["CODER"],
        "priority": 80,
        "executor_tools": ["read_file_snippet", "find_file", "execute_python"],
        "plan_template": [
            "读取错误信息或问题代码文件",
            "分析完整 traceback，从最末帧定位根本原因",
            "提炼最小可复现代码片段",
            "用 diff 格式或完整替换块给出修复代码",
            "说明如何防止同类问题复发（类型检查/防御性断言等）",
        ],
        "prompt": (
            "\n\n## 🐍 领域要求：Python 专项调试\n"
            "- 优先读取完整 traceback，从最末帧（实际出错位置）逆向分析调用链\n"
            "- 调试输出结构：\n"
            "  1. 🔎 **根本原因**：一句话定位问题核心\n"
            "  2. 📋 **最小复现**：提炼能独立运行的最小代码片段\n"
            "  3. 🔧 **修复代码**：diff 格式或完整替换块\n"
            "  4. 💡 **防止复发**：说明如何避免同类问题（类型检查/防御性断言等）\n"
            "- 特别关注：类型错误、None 解引用、编码（UTF-8/GBK）、异步/线程竞态、循环 import\n"
            "- 涉及第三方库时，指出版本差异可能的影响（附 `pip show <pkg>` 验证建议）"
        ),
        "enabled": False,
    },
    {
        "id": "debug_web_frontend",
        "name": "前端调试",
        "icon": "🖥️",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "专项前端调试：JS/TS/CSS 报错、控制台异常、DOM 问题、React/Vue 组件 bug",
        "intent_description": "用户遇到 JS/TS 报错、CSS 样式异常、React/Vue 组件 bug、浏览器控制台错误",
        "task_types": ["CODER"],
        "priority": 80,
        "prompt": (
            "\n\n## 🖥️ 领域要求：前端专项调试\n"
            "- 调试排查顺序：控制台错误 → Network 面板 → DOM 状态 → 组件 state/props\n"
            "- 区分错误来源：运行时错误（TypeError/ReferenceError）、异步错误（Promise rejection）、渲染错误（框架报错边界）\n"
            "- React/Vue 特有检查：key 唯一性、状态更新时机、副作用依赖数组、生命周期执行顺序\n"
            "- CSS 问题排查：选择器权重（specificity）、盒模型（box-sizing）、flexbox 轴方向、stacking context（z-index）\n"
            "- 修复输出：精确到文件行号的代码改动 + DevTools 验证步骤说明"
        ),
        "enabled": False,
    },
    {
        "id": "debug_api",
        "name": "API 联调调试",
        "icon": "🔌",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "专项 API 调试：HTTP 状态码分析、请求/响应体排查、认证失败和超时问题",
        "intent_description": "用户遇到 API 请求失败、HTTP 错误码、认证失败、跨域 CORS、超时等接口问题",
        "task_types": ["CODER", "SYSTEM"],
        "priority": 75,
        "prompt": (
            "\n\n## 🔌 领域要求：API 联调调试\n"
            "- 分析框架：状态码 → 请求头（Authorization/Content-Type） → 请求体格式 → 响应体 → 超时/重试策略\n"
            "- 常见状态码处理指引：\n"
            "  - 4xx：聚焦客户端参数（401=认证缺失、403=权限不足、404=路径错误、422=参数格式）\n"
            "  - 5xx：聚焦服务端日志（500=内部错误、502=网关/反代问题、503=服务过载）\n"
            "  - CORS：检查 Origin 匹配和 Preflight 请求头\n"
            "- 认证问题：给出脱敏后的 curl 验证命令，帮助隔离问题\n"
            "- 超时问题：建议指数退避重试策略，区分 connect_timeout / read_timeout / total_timeout\n"
            "- 输出可直接运行的 curl 命令，便于在终端独立验证"
        ),
        "enabled": False,
    },
    {
        "id": "debug_performance",
        "name": "性能分析",
        "icon": "⚡",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "代码/系统性能分析：定位瓶颈、复杂度评估、profiling 工具使用、分层优化建议",
        "intent_description": "用户遇到程序运行慢、内存占用高、接口延迟大、CPU 飙升等性能瓶颈问题",
        "task_types": ["CODER", "SYSTEM"],
        "priority": 70,
        "conflict_with": ["concise_mode"],
        "prompt": (
            "\n\n## ⚡ 领域要求：性能分析\n"
            "- 分析维度：时间复杂度（Big-O）、空间复杂度、I/O 密集型 vs CPU 密集型\n"
            "- 定位瓶颈优先级：数据库 N+1 查询 > 同步 I/O 阻塞 > 算法复杂度 > 内存分配与 GC\n"
            "- 输出结构：\n"
            "  1. 🎯 **性能瓶颈**：指明具体函数/行号/查询语句\n"
            "  2. 📊 **复杂度对比**：当前 vs 优化后的 Big-O\n"
            "  3. 🔧 **优化方案**：按收益从大到小排列，给出代码示例\n"
            "  4. 🧪 **验证方法**：Python → cProfile/timeit；前端 → Performance DevTools；DB → EXPLAIN ANALYZE\n"
            "- 将建议分层：「快速见效（< 1 天）」和「需要重构（> 3 天）」明确区分"
        ),
        "enabled": False,
    },
    # ── 记忆类 ────────────────────────────────────────────────────────────────
    {
        "id": "long_term_memory",
        "name": "长期记忆",
        "icon": "🧠",
        "category": "memory",
        "skill_nature": "system",
        "description": "跨会话记住用户偏好、项目背景和习惯，无需每次重复说明。通过「记忆」管理页面增删查看已存记忆。",
        "task_types": [],  # 所有任务类型均生效
        "prompt": "",  # prompt 由 inject_into_prompt 动态注入，此处留空
        "enabled": True,  # 默认开启，替代旧版设置里的独立开关
    },
    # ══════════════════════════════════════════════════════════════════════════
    # ── 办公工作流类 ──────────────────────────────────────────────────────────
    # ══════════════════════════════════════════════════════════════════════════
    # ── 邮件 / 沟通 ───────────────────────────────────────────────────────────
    {
        "id": "email_composer",
        "name": "邮件撰写",
        "icon": "📧",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "撰写结构清晰、语气得体的商务邮件：主题行、开头问候、正文分段、行动项、礼貌收尾",
        "intent_description": "用户需要写邮件、回复邮件、发通知、写感谢信或拒绝信",
        "task_types": ["CHAT", "FILE_GEN"],
        "priority": 65,
        "prompt": (
            "\n\n## 📧 工作流要求：邮件撰写\n"
            "- **主题行**：简短有力（≤ 50 字符），包含核心动作词（「确认/跟进/邀请/通知」）\n"
            "- **开头**：一句话说明写信目的，不要「您好，我是…」式的自我介绍铺垫\n"
            "- **正文结构**：背景（1-2句）→ 核心内容（要点化）→ 明确的行动项（截止日期+负责人）\n"
            "- **语气**：根据关系调整——内部同事用简洁直接语气，对外客户/合作方用礼貌专业语气\n"
            "- **收尾**：「期待您的回复」或「如有疑问请随时联系」，附上称呼署名格式\n"
            "- 输出完整邮件正文，同时给出可选的「更正式版本」和「更简短版本」对比"
        ),
        "enabled": False,
    },
    {
        "id": "meeting_minutes",
        "name": "会议纪要",
        "icon": "📋",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "将会议记录/要点整理为标准会议纪要：议题、讨论结论、决策、行动项（含负责人和截止日）",
        "intent_description": "用户需要整理会议记录、生成会议纪要、输出行动清单",
        "task_types": ["CHAT", "FILE_GEN", "DOC_ANNOTATE"],
        "priority": 65,
        "prompt": (
            "\n\n## 📋 工作流要求：会议纪要\n"
            "- **标准格式**：\n"
            "  ```\n"
            "  【会议纪要】\n"
            "  日期：  参会人：  主持人：\n"
            "  ─────────────────────\n"
            "  1. 议题与背景\n"
            "  2. 讨论要点（按议题分块）\n"
            "  3. 决策结论（✅ 已决定 / ⏳ 待定）\n"
            "  4. 行动项（| 任务 | 负责人 | 截止日 | 状态 |）\n"
            "  5. 下次会议安排\n"
            "  ```\n"
            "- 区分「结论性表述」和「讨论性意见」，前者加粗，后者正常字体\n"
            "- 行动项必须有明确责任人，「大家」「团队」不可作为负责人\n"
            "- 若信息不完整，用 `[待确认]` 占位并在末尾列出待确认项"
        ),
        "enabled": False,
    },
    {
        "id": "ppt_outline",
        "name": "PPT 大纲",
        "icon": "🖼️",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "快速生成结构化 PPT 大纲：封面、目录、分章节内容页、结尾页，附每页核心要点和演讲备注",
        "intent_description": "用户需要做演示文稿、做汇报、做提案、输出 PPT 结构",
        "task_types": ["CHAT", "FILE_GEN"],
        "priority": 65,
        "prompt": (
            "\n\n## 🖼️ 工作流要求：PPT 大纲\n"
            "- **输出结构**：\n"
            "  - Slide 1：封面（标题 / 副标题 / 演讲人 / 日期）\n"
            "  - Slide 2：目录页（3-5 个章节）\n"
            "  - 内容页：每页给出「标题 + 3-5 个要点 + 1 句演讲备注」\n"
            "  - 结尾页：核心结论 / 致谢 / Q&A\n"
            "- 每章节控制在 3-5 页，避免内容过密\n"
            "- 建议每页加「可视化提示」：饼图 / 时间轴 / 对比表 / 流程图\n"
            "- 用 `> 演讲备注: …` 格式写每页的口头要点"
        ),
        "enabled": False,
    },
    {
        "id": "report_writer",
        "name": "报告撰写",
        "icon": "📄",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "输出专业报告：执行摘要、背景、分析（含数据支撑）、结论/建议，适用于周报、调研报告、项目报告",
        "intent_description": "用户需要写工作报告、调研报告、项目总结、月报、周报",
        "task_types": ["CHAT", "FILE_GEN", "RESEARCH"],
        "priority": 65,
        "prompt": (
            "\n\n## 📄 工作流要求：报告撰写\n"
            "- **必备章节**：执行摘要（1/4 页）→ 背景与目标 → 方法/数据来源 → 分析发现 → 结论与建议 → 附录\n"
            "- **执行摘要**：用 3-5 句话概括全文最重要的结论，让读者不看正文也能决策\n"
            "- **分析段落**：每个论点用「论点 → 数据/证据 → 解读」三段式展开\n"
            "- **建议部分**：按「紧急程度 × 影响范围」矩阵排列，不超过 5 条，每条有可衡量的成功指标\n"
            "- 数字必须带单位，百分比保留 1 位小数\n"
            "- 若用户未提供数据，明确注明「[数据待填入]」而不是自行编造"
        ),
        "enabled": False,
    },
    # ── 沟通表达类 ────────────────────────────────────────────────────────────
    {
        "id": "feedback_polisher",
        "name": "反馈润色",
        "icon": "💬",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "将直接/负面的反馈改写为建设性、正向语言，保留核心意思但降低对抗性，适合绩效评价、同事 Review",
        "intent_description": "用户想给出批评意见、绩效反馈、代码 Review、评价建议但担心措辞太直接或伤感情",
        "task_types": ["CHAT", "DOC_ANNOTATE"],
        "priority": 60,
        "conflict_with": ["debate_mode"],
        "prompt": (
            "\n\n## 💬 工作流要求：反馈润色\n"
            "- 先输出「润色版本」，再可选附「原文对比」\n"
            "- **SBI 框架**：情境（Situation）→ 行为（Behavior）→ 影响（Impact），"
            "避免人身攻击，聚焦行为和结果\n"
            "- 将「你总是」→「在 X 场景中，我注意到…」；将「很差」→「有改进空间」\n"
            "- 负面内容采用「三明治结构」：正面认可 → 改进意见 → 鼓励展望\n"
            "- 保留核心意思，不可将严重问题软化/掩盖为无关紧要"
        ),
        "enabled": False,
    },
    {
        "id": "negotiation_assist",
        "name": "谈判话术",
        "icon": "🤝",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "为谈判/商务沟通提供话术框架：立场锚定、利益对齐、异议化解、条件交换技巧",
        "intent_description": "用户需要进行薪资谈判、合同谈判、砍价、资源申请、利益协调等场景的沟通准备",
        "task_types": ["CHAT"],
        "priority": 60,
        "prompt": (
            "\n\n## 🤝 工作流要求：谈判话术\n"
            "- 先分析双方「立场（Positions）」和「利益（Interests）」，找共赢空间\n"
            "- 提供「开场锚定」话术：先报出有利数字/条件，为后续让步留余地\n"
            "- 预设对方可能提出的 2-3 个反对意见，给出化解回应\n"
            "- 条件交换原则：让步时附加条件（「如果…那么…」），不单方面让步\n"
            "- 给出「BATNA（最佳替代方案）」提示：告知用户走不下去时的底线选项\n"
            "- 语气保持合作而非对抗，强调「共同目标」"
        ),
        "enabled": False,
    },
    # ── 分析决策类 ────────────────────────────────────────────────────────────
    {
        "id": "pros_cons",
        "name": "利弊分析",
        "icon": "⚖️",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "对任意方案/决策进行结构化利弊分析：优缺点、风险、机会、适用条件，辅助快速决策",
        "intent_description": "用户在多个选项中纠结、需要做决策、比较方案优劣、分析是否值得做某事",
        "task_types": ["CHAT", "RESEARCH"],
        "priority": 60,
        "conflict_with": ["concise_mode"],
        "prompt": (
            "\n\n## ⚖️ 工作流要求：利弊分析\n"
            "- **SWOT 简化框架**：优势 / 劣势 / 机会 / 风险，每项 2-4 点\n"
            "- 用表格对比多个方案（当方案 ≥ 2 个时）\n"
            "- 标明每项的「权重/重要性」：🔴 高 / 🟡 中 / 🟢 低\n"
            "- 给出「综合建议」段落：基于利弊平衡后推荐哪个选项，说明前提条件\n"
            "- 如果信息不足以判断，明确指出「还需知道 X 才能做出更准确判断」"
        ),
        "enabled": False,
    },
    {
        "id": "okr_builder",
        "name": "OKR 制定",
        "icon": "🎯",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "帮助制定 OKR（目标与关键结果）：将模糊目标转化为可量化 KR，检查 SMART 原则",
        "intent_description": "用户需要制定季度/年度目标、写 OKR、将模糊愿景变成可执行计划",
        "task_types": ["CHAT", "FILE_GEN"],
        "priority": 60,
        "prompt": (
            "\n\n## 🎯 工作流要求：OKR 制定\n"
            "- **O（目标）**：鼓舞人心、定性描述、1 个季度内可实现，不含数字\n"
            "- **KR（关键结果）**：2-5 个，必须可量化（含基线值、目标值、截止日）\n"
            "- **SMART 检查**：为每个 KR 注上 S/M/A/R/T 各维度是否满足\n"
            "- 如果用户给出的 KR 是「完成某事」型（任务型），帮助改写为结果型\n"
            "  示例：「完成用户调研」→「在 3 月底前完成 20 份用户深度访谈，沉淀 5 个核心痛点」\n"
            "- 输出格式：\n"
            "  ```\n"
            "  O：[目标描述]\n"
            "  KR1：[可量化结果] — 截止：[日期]\n"
            "  KR2：…\n"
            "  ```"
        ),
        "enabled": False,
    },
    {
        "id": "root_cause",
        "name": "根因分析",
        "icon": "🔍",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "用 5-Why / 鱼骨图等方法对问题进行结构化根因分析，找到可解决的根本原因而非表象",
        "intent_description": "用户遇到反复出现的问题、事故复盘、想找到真正原因而不是临时解决方案",
        "task_types": ["CHAT", "RESEARCH"],
        "priority": 60,
        "prompt": (
            "\n\n## 🔍 工作流要求：根因分析\n"
            "- **首选方法**：5-Why 追问（最少问 3 Why，停在「可控制的人/流程/系统」层级）\n"
            "- **辅助工具**：鱼骨图六维度——人（Man）、机（Machine）、料（Material）、"
            "法（Method）、环（Environment）、测（Measurement）\n"
            "- 区分「直接原因」（表象）和「根本原因」（系统性缺陷）\n"
            "- 对每个根因给出「预防措施」和「纠正措施」，并标明优先级\n"
            "- 输出时用树形缩进展示因果链，让追溯路径一目了然"
        ),
        "enabled": False,
    },
    # ── 文档处理类 ────────────────────────────────────────────────────────────
    {
        "id": "contract_reviewer",
        "name": "合同审阅",
        "icon": "📝",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "审阅合同/协议文本，标出关键条款、潜在风险点、不平等条款和缺失保护条款",
        "intent_description": "用户需要审查合同、协议、框架协议，检查是否有对己方不利的条款或缺失条款",
        "task_types": ["DOC_ANNOTATE", "CHAT"],
        "priority": 70,
        "conflict_with": ["creative_writing"],
        "prompt": (
            "\n\n## 📝 工作流要求：合同审阅\n"
            "- **必检条款**：付款条件 / 违约责任 / 知识产权归属 / 保密条款 / 争议解决方式 / 合同期限与续签\n"
            "- **风险标注等级**：🔴 高风险（需修改才能签）/ 🟡 中风险（建议谈判）/ 🟢 低风险（可接受）\n"
            "- 对每个风险条款：指出问题 → 解释为何有风险 → 给出建议改写措辞\n"
            "- 检查「缺失保护条款」：指出通常应有但本合同未包含的条款\n"
            "- ⚠️ 提示：本分析仅供参考，重要合同请咨询专业律师"
        ),
        "enabled": False,
    },
    {
        "id": "sop_writer",
        "name": "SOP 编写",
        "icon": "📐",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "将操作经验转化为标准操作规程（SOP）：触发条件、前提、步骤、检查点、异常处理",
        "intent_description": "用户需要将操作流程文档化、制定标准作业程序、写操作手册、固化业务流程",
        "task_types": ["CHAT", "FILE_GEN"],
        "priority": 65,
        "prompt": (
            "\n\n## 📐 工作流要求：SOP 编写\n"
            "- **SOP 结构**：\n"
            "  1. 文档信息（名称 / 版本 / 编写人 / 生效日期 / 适用范围）\n"
            "  2. 目的与背景\n"
            "  3. 前提条件（工具 / 权限 / 基础知识）\n"
            "  4. 操作步骤（编号，每步：操作 + 预期结果 + 截图/示例占位符）\n"
            "  5. 质量检查点（✅ 完成标志）\n"
            "  6. 异常情况及处理方法\n"
            "  7. 相关文档链接\n"
            "- 步骤语言用「动词 + 对象」格式（「点击 → 保存按钮」而非「按钮被点击」）\n"
            "- 估算每步耗时，帮助读者预期总用时"
        ),
        "enabled": False,
    },
    {
        "id": "doc_qa",
        "name": "文档问答",
        "icon": "❓",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "基于用户提供的文档内容进行精准问答，引用原文段落，不在文档外推测",
        "intent_description": "用户粘贴了一段文档/手册/规范/报告，希望从中找到特定问题的答案",
        "task_types": ["CHAT", "DOC_ANNOTATE"],
        "priority": 65,
        "prompt": (
            "\n\n## ❓ 工作流要求：文档问答\n"
            "- **严格基于文档**：所有回答必须引用原文，用「> 原文：…」格式标注来源段落\n"
            "- 如果文档中没有答案，明确说「文档中未提及此问题」，不进行推测性回答\n"
            "- 如果文档内容与问题有部分相关，给出「文档中最接近的信息」并注明差异\n"
            "- 多处文档内容共同回答一个问题时，综合引用并标注各段来源位置\n"
            "- 可以指出文档中的矛盾、模糊或过时之处"
        ),
        "enabled": False,
    },
    # ── 项目管理类 ────────────────────────────────────────────────────────────
    {
        "id": "risk_register",
        "name": "风险登记",
        "icon": "⚠️",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "为项目/计划生成风险登记册：识别风险、评估可能性×影响程度、制定应对策略",
        "intent_description": "用户需要做风险评估、识别项目风险、制定风险应对计划",
        "task_types": ["CHAT", "FILE_GEN"],
        "priority": 60,
        "prompt": (
            "\n\n## ⚠️ 工作流要求：风险登记\n"
            "- **风险矩阵**：可能性（1-5）× 影响（1-5）= 风险分数，≥15 为高风险\n"
            "- **输出表格**：| 风险 ID | 描述 | 可能性 | 影响 | 风险分 | 应对策略 | 负责人 |\n"
            "- **应对策略四选项**：规避（Avoid）/ 转移（Transfer）/ 减轻（Mitigate）/ 接受（Accept）\n"
            "- 对每个高风险（🔴）配套「触发信号」：当哪些表现出现时需要启动应对措施\n"
            "- 检查常见遗漏风险类：技术风险 / 人员风险 / 外部依赖风险 / 合规风险 / 预算风险"
        ),
        "enabled": False,
    },
    {
        "id": "sprint_planner",
        "name": "迭代规划",
        "icon": "🏃",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "将需求列表拆分为迭代/Sprint，估算工作量、排优先级、输出迭代目标和任务看板",
        "intent_description": "用户需要规划一个 Sprint/迭代、排期需求、拆解开发任务、做产品路线图",
        "task_types": ["CHAT", "FILE_GEN"],
        "priority": 60,
        "prompt": (
            "\n\n## 🏃 工作流要求：迭代规划\n"
            "- **优先级排序**：按「价值 × 紧迫性 ÷ 工作量」评分，用 MoSCoW 标注（Must/Should/Could/Won't）\n"
            "- **任务拆解**：每个功能拆到「1-3 天可完成」的粒度，配子任务 checklist\n"
            "- **工作量估算**：给出故事点（1/2/3/5/8/13），并注明估算假设\n"
            "- **Team Capacity**：提示用户填写本迭代可用人天，计算是否超负荷\n"
            "- 输出「迭代目标一句话」（Sprint Goal），便于对齐团队方向\n"
            "- 输出任务表格：| 任务 | 优先级 | 故事点 | 负责人 | 状态 |"
        ),
        "enabled": False,
    },
    # ── 数据与量化类 ──────────────────────────────────────────────────────────
    {
        "id": "kpi_designer",
        "name": "KPI 设计",
        "icon": "📊",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "为岗位/团队/业务线设计 KPI 指标体系：北极星指标、过程指标、质量指标，附评分规则",
        "intent_description": "用户需要制定绩效指标、设计考核体系、评估业务健康度",
        "task_types": ["CHAT", "FILE_GEN"],
        "priority": 60,
        "prompt": (
            "\n\n## 📊 工作流要求：KPI 设计\n"
            "- **三层指标**：\n"
            "  - L1 北极星指标（1个，反映终极业务目标）\n"
            "  - L2 驱动指标（3-5个，直接影响 L1 的可控变量）\n"
            "  - L3 健康指标（护栏，保证追求 L1 时不出现副作用，如用户满意度）\n"
            "- 每个 KPI 必须包含：公式定义 / 数据来源 / 统计周期 / 基准值 / 目标值\n"
            "- 检查「反激励」风险：该 KPI 是否可能导致员工做对指标但对业务有害的行为\n"
            "- 给出「及格线 / 良好线 / 优秀线」三档评分标准"
        ),
        "enabled": False,
    },
    {
        "id": "survey_designer",
        "name": "问卷设计",
        "icon": "📊",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "设计调研问卷：问题结构、量表类型、避免引导性问题、样本策略，输出可用问卷草稿",
        "intent_description": "用户需要做用户调研、员工满意度调查、市场调研、需求收集问卷",
        "task_types": ["CHAT", "FILE_GEN"],
        "priority": 55,
        "prompt": (
            "\n\n## 📊 工作流要求：问卷设计\n"
            "- **问卷结构**：引言（说明目的+预计耗时）→ 筛选题 → 核心问题 → 开放问题 → 感谢语\n"
            "- **题型匹配**：\n"
            "  - 态度/满意度 → Likert 5/7 点量表\n"
            "  - 行为频率 → 频率量表（从不/偶尔/经常/总是）\n"
            "  - 排名 → 强制排序题\n"
            "  - 原因 → 多选题 + 「其他」开放项\n"
            "- 检查「引导性问题」并改写（去掉情感词，改为中性表达）\n"
            "- 控制总量：≤ 15 题为佳，超过提示用户精简\n"
            "- 附「问卷说明」参考文案"
        ),
        "enabled": False,
    },
    # ── 学习与个人成长类 ──────────────────────────────────────────────────────
    {
        "id": "learning_roadmap",
        "name": "学习路线图",
        "icon": "🗺️",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "为技能/领域制定个人学习路线图：阶段划分、核心资源推荐、里程碑检验点",
        "intent_description": "用户想学一门技术/技能、转行、系统学习某个领域，需要路线规划",
        "task_types": ["CHAT", "RESEARCH"],
        "priority": 55,
        "prompt": (
            "\n\n## 🗺️ 工作流要求：学习路线图\n"
            "- **阶段划分**（3-4 阶段，每阶段配预估时长）：\n"
            "  - 入门（快速建立基本认知）\n"
            "  - 实践（完成 1-2 个真实项目）\n"
            "  - 进阶（深入原理，处理复杂场景）\n"
            "  - 精通（输出/教别人，参与社区）\n"
            "- 每阶段推荐：1 本书 / 1 个课程 / 1 个实战项目\n"
            "- 设置「里程碑检验题」：每阶段完成后用 2-3 个问题自测是否达标\n"
            "- 根据用户当前水平和每周可投入时长，给出个性化时间预估\n"
            "- 标注「常见陷阱」：哪些资源过时/质量差，应避开"
        ),
        "enabled": False,
    },
    {
        "id": "interview_prep",
        "name": "面试准备",
        "icon": "🎙️",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "为特定岗位准备面试：高频问题模拟、STAR 法则答题框架、技术考点梳理、公司研究要点",
        "intent_description": "用户即将参加面试，需要准备回答问题、模拟面试、了解面试套路",
        "task_types": ["CHAT"],
        "priority": 55,
        "prompt": (
            "\n\n## 🎙️ 工作流要求：面试准备\n"
            "- **高频问题分类**：自我介绍 / 过往经历 / 优缺点 / 职业规划 / 行为问题 / 技术问题\n"
            "- **STAR 答题框架**：情境（Situation）→ 任务（Task）→ 行动（Action）→ 结果（Result）\n"
            "  - 每条经历都要量化结果（百分比 / 绝对数值）\n"
            "- 提供 3-5 个该岗位最可能被问到的「刁钻问题」及化解思路\n"
            "- **反问环节**：给出 3 个高质量反问问题，展现候选人主动性\n"
            "- 针对用户描述的岗位，输出「一定要做的功课」（公司背景 / 近期新闻 / 产品体验）"
        ),
        "enabled": False,
    },
    # ── 创意与策略类 ──────────────────────────────────────────────────────────
    {
        "id": "brainstorm",
        "name": "头脑风暴",
        "icon": "🧩",
        "category": "workflow",
        "skill_nature": "model_hint",
        "description": "快速发散生成大量创意，使用 SCAMPER / 随机刺激 / 逆向思维等技巧，不过早批判想法",
        "intent_description": "用户需要想点子、产品创意、营销方案、命名灵感、任何创意发散场景",
        "task_types": ["CHAT"],
        "priority": 55,
        "conflict_with": ["strict_mode", "concise_mode"],
        "prompt": (
            "\n\n## 🧩 工作流要求：头脑风暴\n"
            "- **数量优先**：第一轮无条件输出 10-15 个想法，不自我审查\n"
            "- **SCAMPER 激发**：Substitute / Combine / Adapt / Modify / Put to other use / Eliminate / Reverse\n"
            "- **逆向提问**：「如果要让这件事完全失败，你会怎么做？」→ 反转为正向方案\n"
            "- 用「普通 / 有趣 / 疯狂」三档标记想法，避免过早删掉极端想法\n"
            "- 第二轮：从第一轮中挑选 3 个最有潜力的想法，深化成可执行方案\n"
            "- 鼓励跨界类比：「X 行业是怎么解决类似问题的？」"
        ),
        "enabled": False,
    },
    {
        "id": "social_copy",
        "name": "社媒文案",
        "icon": "📱",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "为微信公众号/小红书/LinkedIn/朋友圈等平台生成针对性文案，匹配各平台风格和用户习惯",
        "intent_description": "用户需要写公众号文章、小红书笔记、微博、LinkedIn 帖子、朋友圈文案等社交媒体内容",
        "task_types": ["CHAT", "FILE_GEN"],
        "priority": 60,
        "conflict_with": ["professional_tone"],
        "prompt": (
            "\n\n## 📱 工作流要求：社媒文案\n"
            "- **平台风格适配**：\n"
            "  - 微信/公众号：标题悬念感+数字，正文有故事性，结尾引导互动\n"
            "  - 小红书：口语化、亲切感、「姐妹/宝藏」等平台用语，多用 emoji，分点排版\n"
            "  - LinkedIn：专业叙事，第一人称故事+洞察，#标签3个以内\n"
            "  - 朋友圈：100字以内，金句开头，留白余味\n"
            "- **钩子公式**：痛点/惊喜/故事开头，前 2 行必须抓住人\n"
            "- 每次输出 2-3 个版本供选择，并标注适用场景\n"
            "- 附推荐 #话题标签（平台相关）"
        ),
        "enabled": False,
    },
    {
        "id": "prompt_engineer",
        "name": "Prompt 优化",
        "icon": "✏️",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "分析并改写用户的 AI 提示词，使其更清晰、有结构、产出更好的 AI 回复",
        "intent_description": "用户想优化自己写的 AI 提示词、Prompt 不好用、想让 AI 更准确地完成任务",
        "task_types": ["CHAT", "DOC_ANNOTATE"],
        "priority": 55,
        "prompt": (
            "\n\n## ✏️ 工作流要求：Prompt 优化\n"
            "- **诊断原 Prompt**：指出模糊之处（缺少角色/缺少格式/缺少约束/目标不清晰）\n"
            "- **改写原则**：角色设定 + 任务 + 背景 + 输出格式 + 约束条件\n"
            "- **输出结构**：\n"
            "  1. 📋 原 Prompt 问题分析（2-3点）\n"
            "  2. ✅ 改写后的 Prompt\n"
            "  3. 🔎 改写说明（解释每项改动的原因）\n"
            "  4. 💡 进阶变体（如需不同风格/用途，给出 2 个变体）\n"
            "- 对 Prompt 技巧给出简短说明：few-shot / chain-of-thought / 角色扮演 等"
        ),
        "enabled": False,
    },
    # ─────────── 文件读取与解析类 Skills ───────────
    {
        "id": "pdf_reader",
        "name": "PDF 深度解析",
        "icon": "📕",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "精准读取和解析 PDF 文档，使用正确的工具调用方式处理各种 PDF 场景",
        "intent_description": "用户要读取 PDF、解析 PDF 内容、提取 PDF 文字、分析 PDF 文档",
        "task_types": ["RESEARCH", "FILE_GEN", "DOC_ANNOTATE", "CHAT"],
        "priority": 60,
        "executor_tools": ["read_file_snippet", "find_file", "summarize_file"],
        "plan_template": [
            "定位目标 PDF 文件路径（未知时用 find_file 搜索）",
            "调用 read_file_snippet 读取前 8000 字内容",
            "若内容被截断则再调用 summarize_file 获取针对性摘要",
            "基于读取内容直接回答用户问题",
        ],
        "prompt": (
            "\n\n## 📕 文件工具：PDF 深度解析\n"
            "\n**工具调用规范：**\n"
            "1. **路径已知** → 直接调用 `read_file_snippet(path, max_chars=8000)` 读取前 8000 字\n"
            "   - 返回内容若末尾带 `...` 说明被截断，文件有更多内容\n"
            '   - 截断时：再调用 `summarize_file(path, focus="用户的具体问题")` 获取针对性摘要\n'
            '2. **路径未知** → 先调用 `find_file(query="文件名关键词", category="document", limit=5)` 定位文件\n'
            '3. **超长文档（>10页）** → 优先用 `summarize_file(path, focus="...")` 而非全文读取\n'
            "\n**边界情况处理：**\n"
            "- 若读取内容乱码（全是`?????`）→ 提示用户该 PDF 可能是扫描件无法直接提取文字\n"
            "- 若返回内容极短或空白 → 提示 PDF 可能加密或损坏，建议用户手动检查\n"
            "- 多页 PDF：`read_file_snippet` 只读前 8 页，如需更多页请用 `summarize_file`\n"
            "\n**输出要求：**\n"
            "- 读取完成后直接呈现内容，不要重复文件路径\n"
            "- 如用户有具体问题，基于读取内容直接回答，不要原文堆砌"
        ),
        "enabled": False,
    },
    {
        "id": "multi_format_reader",
        "name": "多格式文件读取",
        "icon": "📂",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "按文件类型自动选择正确的读取策略，支持 PDF/DOCX/XLSX/CSV/PPTX/TXT 等",
        "intent_description": "用户要读取或查看各种类型的文件，需要解析文件内容",
        "task_types": ["RESEARCH", "FILE_GEN", "CHAT", "DOC_ANNOTATE"],
        "priority": 61,
        "executor_tools": ["read_file_snippet", "find_file", "summarize_file", "list_directory"],
        "plan_template": [
            "识别文件扩展名确定读取策略",
            "路径不确定时用 find_file 定位文件",
            "按格式调用对应工具读取内容（PDF/DOCX/XLSX/CSV/TXT）",
            "文件过长时改用 summarize_file 获取针对性摘要",
            "整理并呈现文件内容，回答用户问题",
        ],
        "prompt": (
            "\n\n## 📂 文件工具：多格式文件读取规范\n"
            "\n**按扩展名选择读取方式（严格遵守）：**\n"
            "| 格式 | 工具调用 | 最大字数 | 说明 |\n"
            "|------|---------|---------|------|\n"
            "| `.pdf` | `read_file_snippet(path, 8000)` | PyPDF2 最多 8 页 | 扫描件可能乱码 |\n"
            "| `.docx` | `read_file_snippet(path, 8000)` | 全部段落 | python-docx 提取 |\n"
            "| `.xlsx` / `.xls` | `read_file_snippet(path, 6000)` | 最多 200 行 | openpyxl 提取 |\n"
            "| `.csv` | `read_file_snippet(path, 10000)` | 纯文本 | 逗号分隔 |\n"
            "| `.txt` / `.md` | `read_file_snippet(path, 10000)` | 全文 | UTF-8/GBK 自动检测 |\n"
            "| `.json` / `.xml` / `.yaml` | `read_file_snippet(path, 10000)` | 结构化文本 | 原样返回 |\n"
            "| `.py` / `.js` / `.ts` / `.sql` | `read_file_snippet(path, 10000)` | 代码文件 | 原样返回 |\n"
            "| `.pptx` | ⚠️ **不支持** read_file_snippet | — | 见下方特殊处理 |\n"
            "\n**PPTX 特殊处理（必须走 CODER 路径）：**\n"
            "```python\n"
            "from pptx import Presentation\n"
            "prs = Presentation(r'文件路径.pptx')\n"
            "for i, slide in enumerate(prs.slides, 1):\n"
            "    print(f'=== 第{i}页 ===')\n"
            "    for shape in slide.shapes:\n"
            "        if hasattr(shape, 'text') and shape.text.strip():\n"
            "            print(shape.text.strip())\n"
            "```\n"
            "\n**通用规则：**\n"
            '- 路径不确定时先调用 `find_file(query="文件名", limit=5)` 定位\n'
            '- 文件过长时改用 `summarize_file(path, focus="用户问题")` 获取摘要\n'
            "- 读取失败时告知用户具体失败原因（文件不存在/格式不对/编码问题）"
        ),
        "enabled": False,
    },
    {
        "id": "long_doc_parser",
        "name": "长文档分段解析",
        "icon": "📜",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "处理超出上下文窗口的长文档，使用分块读取+递进摘要的策略",
        "intent_description": "用户要分析很长的文档、文章、报告，内容超过一次能读取的范围",
        "task_types": ["RESEARCH", "FILE_GEN", "DOC_ANNOTATE"],
        "priority": 62,
        "prompt": (
            "\n\n## 📜 工作流：长文档分段解析策略\n"
            "\n**判断触发条件：** 当用户文档超过 10000 字 / 10 页时启动此工作流\n"
            "\n**三步解析法：**\n"
            "1. **快速预览** \n"
            "   ```\n"
            "   read_file_snippet(path, max_chars=3000)  # 读取文档开头，了解结构\n"
            "   ```\n"
            "   读完后：识别文档类型（报告/合同/论文/手册）、章节结构、主要话题\n"
            "\n"
            "2. **针对性摘要**（按需多次调用，每次聚焦不同维度）：\n"
            "   ```\n"
            '   summarize_file(path, focus="文档的核心结论和主要观点")\n'
            '   summarize_file(path, focus="具体的数据、统计和事实")\n'
            '   summarize_file(path, focus="风险、问题和注意事项")\n'
            "   ```\n"
            "\n"
            "3. **综合整理**：将多次摘要的结果合并，用结构化格式呈现\n"
            "\n**输出格式建议：**\n"
            "- 📌 文档概述（1句话总结）\n"
            "- 📋 核心内容（分章节条目）\n"
            "- 📊 关键数据/事实\n"
            "- ⚠️ 注意事项/风险点（如有）\n"
            "- 💡 我的分析与建议（基于内容）\n"
            "\n**注意：** `summarize_file` 内部最多处理 6000 字原文，不同的 `focus` 参数会引导 AI 提取不同方面"
        ),
        "enabled": False,
    },
    {
        "id": "spreadsheet_analyst",
        "name": "表格数据深度分析",
        "icon": "📊",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "读取 Excel/CSV 表格后进行结构化分析，识别数据类型、清洗问题和统计摘要",
        "intent_description": "用户要分析表格数据、Excel 文件、CSV 数据，进行统计或数据解读",
        "task_types": ["RESEARCH", "FILE_GEN", "CHAT"],
        "priority": 63,
        "executor_tools": ["read_file_snippet", "find_file", "execute_python"],
        "plan_template": [
            "定位并读取 Excel/CSV 数据文件",
            "分析数据结构（行列数、字段名、数据类型、缺失值）",
            "执行统计分析（均值/中位数/分布/趋势）",
            "识别数据清洗问题（重复行/异常值/格式不一致）",
            "生成结构化分析报告或可视化建议",
        ],
        "prompt": (
            "\n\n## 📊 工作流：表格数据深度分析\n"
            "\n**第一步：读取表格**\n"
            "```\n"
            "read_file_snippet(path, max_chars=6000)  # xlsx/xls 最多 200 行，csv 全文\n"
            "```\n"
            "\n**第二步：结构识别（读完立即做）**\n"
            "- 识别表头（第一行通常是列名）\n"
            "- 识别数据类型：数字列、文本列、日期列、分类列\n"
            "- 检查数据质量：空值比例、重复行、异常值（负数、0值、超大值）\n"
            "\n**第三步：按用户需求分析**\n"
            "| 分析类型 | 关注点 |\n"
            "|---------|-------|\n"
            "| 汇总统计 | 总计/合计/平均/最大最小 |\n"
            "| 趋势分析 | 日期列排序后的变化趋势 |\n"
            "| 分组对比 | 按分类列分组后各组差异 |\n"
            "| 数据清洗 | 列出存在问题的行和修复建议 |\n"
            "\n**输出规范：**\n"
            "- 优先用 Markdown 表格呈现数据\n"
            "- 数字保留 2 位小数，千位加逗号\n"
            "- 列出前 5-10 行数据作为样本展示\n"
            "- 最后给出 1-3 条数据洞察（用❗标注重要发现）\n"
            "\n**如果数据很大（>200行）**：用 CODER 路径写 pandas 代码进行完整分析"
        ),
        "enabled": False,
    },
    {
        "id": "multi_doc_synthesis",
        "name": "多文档综合分析",
        "icon": "🗂️",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "同时读取多个文件，进行对比、整合、矛盾识别和综合摘要",
        "intent_description": "用户要对比多个文件、综合分析多份文档、找出文件间的差异或共同点",
        "task_types": ["RESEARCH", "FILE_GEN"],
        "priority": 64,
        "prompt": (
            "\n\n## 🗂️ 工作流：多文档综合分析\n"
            "\n**执行步骤：**\n"
            "1. **逐一读取** 每个文件（每个文件 `max_chars=5000`，确保覆盖主要内容）\n"
            "   ```\n"
            "   read_file_snippet(file1_path, 5000)\n"
            "   read_file_snippet(file2_path, 5000)\n"
            "   # ... 重复直到所有文件读完\n"
            "   ```\n"
            '   如果某个文件过长 → 改用 `summarize_file(path, focus="与其他文件对比的核心内容")`\n'
            "\n"
            "2. **构建对比矩阵** — 读完所有文件后，按以下维度对比：\n"
            "   - 🟢 **共同点**：各文档都提到的核心内容\n"
            "   - 🔴 **矛盾点**：各文档中相互冲突的信息\n"
            "   - 🔵 **差异点**：某文档独有而其他没有的信息\n"
            "   - 📊 **数据差异**：不同文档中相同指标的不同数值\n"
            "\n"
            "3. **结构化输出**：\n"
            "   - 用 Markdown 表格做侧边对比\n"
            "   - 格式：| 维度 | 文件1 | 文件2 | 文件3 |\n"
            "   - 最后给出综合结论（整合所有来源后的最终判断）\n"
            "\n**注意：** 若文件路径不明，先用 `find_file(query=...)` 定位每个文件"
        ),
        "enabled": False,
    },
    # ─────────── 文件生成类 Skills ───────────
    {
        "id": "ppt_generator_pro",
        "name": "PPT 专业生成",
        "icon": "📽️",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "输出标准 Markdown 大纲，由系统自动生成专业幻灯片 PPTX 文件",
        "intent_description": "用户要生成 PPT、做幻灯片报告、制作演示文稿",
        "task_types": ["FILE_GEN"],
        "priority": 65,
        "prompt": (
            "\n\n## 📽️ PPT 生成规范：输出 Markdown 大纲"
            "\n\n**重要：不要调用任何工具。直接输出以下格式的 Markdown，系统将自动生成 PPTX 文件。**"
            "\n\n### 输出格式（必须严格遵守）"
            "\n```markdown\n"
            "# PPT 主标题\n"
            "\n"
            "## 幻灯片1标题\n"
            "- 要点一（15字以内）\n"
            "- 要点二\n"
            "- 要点三\n"
            "\n"
            "[过渡页]\n"
            "## 第二章：章节过渡标题\n"
            "过渡页说明文字（一句话）\n"
            "\n"
            "[对比]\n"
            "## 对比幻灯片标题\n"
            "### 左侧标题\n"
            "- 左侧要点1\n"
            "- 左侧要点2\n"
            "### 右侧标题\n"
            "- 右侧要点1\n"
            "- 右侧要点2\n"
            "```"
            "\n\n### 可用幻灯片类型标签（写在 ## 之前一行）"
            "\n- `[详细]` 普通内容页（默认，可省略）"
            "\n- `[重点]` / `[亮点]` / `[数据]` 强调页，突出显示"
            "\n- `[概览]` 带子章节的概述页"
            "\n- `[对比]` 左右对比页（用 ### 分左右两侧）"
            "\n- `[过渡页]` / `[过渡]` 章节分隔页（正文一行说明文字）"
            "\n\n### 内容规范"
            "\n- 第一行必须是 `# 主标题`（封面）"
            "\n- 每节用 `## 标题` 开始，不要跳级"
            "\n- 每页最多 **5 个要点**，超过 5 个拆分为两页"
            "\n- 典型结构：封面 → 目录 → 各章内容页 → 总结"
            "\n- 总页数：10-15页（简短汇报）/ 15-25页（详细方案）"
            "\n- 输出完 Markdown 后不要再说 '已调用工具' 或附加 JSON，系统自动处理"
        ),
        "enabled": False,
    },
    {
        "id": "excel_generator_pro",
        "name": "Excel 专业生成",
        "icon": "🟩",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "通过 Python/openpyxl 生成带格式的 Excel 表格，包含标题样式、数据行和自动列宽",
        "intent_description": "用户要生成 Excel 表格、制作统计表、创建工作表格",
        "task_types": ["FILE_GEN", "CODER"],
        "priority": 66,
        "prompt": (
            "\n\n## 🟩 工具规范：Excel 生成（openpyxl 标准模板）\n"
            "\n**必须使用 CODER 路径！生成完整可运行代码：**\n"
            "```python\n"
            "import openpyxl\n"
            "from openpyxl.styles import Font, PatternFill, Alignment, Border, Side\n"
            "import os\n"
            "\n"
            "wb = openpyxl.Workbook()\n"
            "ws = wb.active\n"
            'ws.title = "数据表"\n'
            "\n"
            "# === 标题行样式 ===\n"
            'headers = ["列名1", "列名2", "列名3"]  # 按实际需求替换\n'
            'header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")\n'
            'header_font = Font(bold=True, color="FFFFFF", size=12)\n'
            "border = Border(left=Side(style='thin'), right=Side(style='thin'),\n"
            "                top=Side(style='thin'), bottom=Side(style='thin'))\n"
            "for col_i, h in enumerate(headers, 1):\n"
            "    c = ws.cell(1, col_i, h)\n"
            "    c.fill = header_fill; c.font = header_font\n"
            "    c.alignment = Alignment(horizontal='center'); c.border = border\n"
            "\n"
            "# === 数据行 ===\n"
            "data = [\n"
            '    # ["值1", "值2", "值3"],  # 按实际数据填充\n'
            "]\n"
            'alt_fill = PatternFill(start_color="EBF3FB", end_color="EBF3FB", fill_type="solid")\n'
            "for row_i, row in enumerate(data, 2):\n"
            "    for col_i, val in enumerate(row, 1):\n"
            "        c = ws.cell(row_i, col_i, val)\n"
            "        c.border = border\n"
            "        if row_i % 2 == 0: c.fill = alt_fill  # 隔行填色\n"
            "\n"
            "# === 自动列宽 ===\n"
            "for col in ws.columns:\n"
            "    max_len = max((len(str(c.value or '')) for c in col), default=8)\n"
            "    ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)\n"
            "\n"
            "ws.freeze_panes = 'A2'  # 冻结标题行\n"
            "output_path = os.path.join(os.path.expanduser('~'), 'Desktop', 'output.xlsx')\n"
            "wb.save(output_path)\n"
            "print(f'✅ Excel 已保存: {output_path}')\n"
            "```\n"
            "\n**生成要求：**\n"
            "- 按用户需求替换表头和数据，不要留占位符\n"
            "- 数字列设置为数字格式；日期列设置日期格式\n"
            "- 最后一行可加合计行（SUM 公式）"
        ),
        "enabled": False,
    },
    {
        "id": "docx_generator_pro",
        "name": "Word 文档专业生成",
        "icon": "📝",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "通过 python-docx 生成带标题层级、正文、表格、样式的 Word 文档",
        "intent_description": "用户要生成 Word 文档、创建正式文档、制作报告",
        "task_types": ["FILE_GEN", "CODER"],
        "priority": 67,
        "prompt": (
            "\n\n## 📝 工具规范：Word 文档生成（python-docx 标准模板）\n"
            "\n**必须使用 CODER 路径！生成完整可运行代码：**\n"
            "```python\n"
            "from docx import Document\n"
            "from docx.shared import Pt, Cm, RGBColor\n"
            "from docx.enum.text import WD_ALIGN_PARAGRAPH\n"
            "import os\n"
            "\n"
            "doc = Document()\n"
            "\n"
            "# === 页面设置 ===\n"
            "section = doc.sections[0]\n"
            "section.page_width  = Cm(21.0)  # A4 宽\n"
            "section.page_height = Cm(29.7)  # A4 高\n"
            "section.left_margin = section.right_margin = Cm(2.54)\n"
            "\n"
            "# === 文档标题 ===\n"
            "title = doc.add_heading('文档标题', level=0)\n"
            "title.alignment = WD_ALIGN_PARAGRAPH.CENTER\n"
            "\n"
            "# === 一级标题 ===\n"
            "doc.add_heading('1. 第一章标题', level=1)\n"
            "\n"
            "# === 正文段落 ===\n"
            "para = doc.add_paragraph('正文内容。首行缩进两字符。')\n"
            "para.paragraph_format.first_line_indent = Pt(24)  # 首行缩进\n"
            "\n"
            "# === 列表项 ===\n"
            "doc.add_paragraph('• 要点一', style='List Bullet')\n"
            "doc.add_paragraph('• 要点二', style='List Bullet')\n"
            "\n"
            "# === 表格 ===\n"
            "table = doc.add_table(rows=1, cols=3)\n"
            "table.style = 'Table Grid'\n"
            "hdr = table.rows[0].cells\n"
            "hdr[0].text = '列1'; hdr[1].text = '列2'; hdr[2].text = '列3'\n"
            "row = table.add_row().cells\n"
            "row[0].text = '数据1'; row[1].text = '数据2'; row[2].text = '数据3'\n"
            "\n"
            "output_path = os.path.join(os.path.expanduser('~'), 'Desktop', 'output.docx')\n"
            "doc.save(output_path)\n"
            "print(f'✅ Word 文档已保存: {output_path}')\n"
            "```\n"
            "\n**生成要求：**\n"
            "- 按用户需求填写真实内容，不要留空白占位符\n"
            "- 标题层级：一级(level=1) 用于章节，二级(level=2) 用于小节\n"
            "- 中文文档设置首行缩进两字，英文文档不缩进"
        ),
        "enabled": False,
    },
    {
        "id": "pdf_generator_pro",
        "name": "PDF 文档生成",
        "icon": "🔴",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "通过 fpdf2 或 reportlab 生成带格式的 PDF，正确处理中文字符",
        "intent_description": "用户要生成 PDF 文件、创建不可编辑文档、导出 PDF 报告",
        "task_types": ["FILE_GEN", "CODER"],
        "priority": 68,
        "prompt": (
            "\n\n## 🔴 工具规范：PDF 生成（推荐 reportlab，支持中文）\n"
            "\n**必须使用 CODER 路径！推荐方案（reportlab，中文兼容）：**\n"
            "```python\n"
            "from reportlab.lib.pagesizes import A4\n"
            "from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle\n"
            "from reportlab.lib.units import cm\n"
            "from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle\n"
            "from reportlab.lib import colors\n"
            "from reportlab.pdfbase import pdfmetrics\n"
            "from reportlab.pdfbase.cidfonts import UnicodeCIDFont\n"
            "import os\n"
            "\n"
            "# 注册中文字体（STSong-Light 是 reportlab 内置的中文字体）\n"
            "pdfmetrics.registerFont(UnicodeCIDFont('STSong-Light'))\n"
            "\n"
            "output_path = os.path.join(os.path.expanduser('~'), 'Desktop', 'output.pdf')\n"
            "doc = SimpleDocTemplate(output_path, pagesize=A4,\n"
            "                        leftMargin=2.5*cm, rightMargin=2.5*cm,\n"
            "                        topMargin=2.5*cm, bottomMargin=2.5*cm)\n"
            "\n"
            "# 中文样式\n"
            "styles = getSampleStyleSheet()\n"
            "title_style = ParagraphStyle('Title', fontName='STSong-Light', fontSize=18,\n"
            "                             spaceAfter=12, alignment=1)  # 居中\n"
            "body_style  = ParagraphStyle('Body',  fontName='STSong-Light', fontSize=12,\n"
            "                             leading=20, spaceAfter=8)\n"
            "\n"
            "story = []\n"
            "story.append(Paragraph('文档标题', title_style))\n"
            "story.append(Spacer(1, 0.5*cm))\n"
            "story.append(Paragraph('正文内容，支持中文。', body_style))\n"
            "\n"
            "doc.build(story)\n"
            "print(f'✅ PDF 已保存: {output_path}')\n"
            "```\n"
            "\n**备用方案（fpdf2，简单文档）：**\n"
            "```python\n"
            "# 注意：fpdf2 需要系统有中文字体文件\n"
            "# 如果没有字体文件，只能生成英文 PDF，中文会显示为方块\n"
            "from fpdf import FPDF\n"
            "pdf = FPDF(); pdf.add_page(); pdf.set_font('Helvetica', size=12)\n"
            "pdf.cell(0, 10, 'English only content', new_x='LMARGIN', new_y='NEXT')\n"
            "pdf.output('output.pdf')\n"
            "```\n"
            "\n**选择建议：**\n"
            "- **含中文内容** → 必须用 reportlab（内置 STSong-Light 字体）\n"
            "- **纯英文简单文档** → fpdf2 更轻量\n"
            "- **复杂布局（多列/图文混排）** → reportlab + Platypus"
        ),
        "enabled": False,
    },
    {
        "id": "docx_translator",
        "name": "Word 文档翻译",
        "icon": "🌐",
        "category": "workflow",
        "skill_nature": "domain_skill",
        "description": "在保留原始格式（样式、字体、表格、页眉页脚）的前提下，将 Word 文档翻译为指定语言，输出新 .docx 文件",
        "intent_description": "用户要翻译 Word/docx 文档到其他语言，保持格式不变",
        "task_types": ["FILE_GEN", "CODER", "CHAT"],
        "priority": 70,
        "prompt": (
            "\n\n## 🌐 Word 文档翻译（格式保留）\n"
            "\n**此功能由系统服务器端自动执行，不需要你写代码或调用任何工具。**\n"
            "\n**你只需要：**\n"
            "1. 告知用户翻译已开始（例如：'正在将您的 Word 文档翻译成英文，格式将完整保留...'）\n"
            "2. 等待系统完成后，展示翻译结果摘要\n"
            "\n**系统已自动处理：**\n"
            "- ✅ 读取 .docx 文档的所有段落、表格、页眉页脚\n"
            "- ✅ 使用 AI 翻译所有文本内容\n"
            "- ✅ 保留原始格式（段落样式/加粗/斜体/字号/颜色/表格结构）\n"
            "- ✅ 输出同名文件加语言后缀（如 `文件_en.docx`），不覆盖原文\n"
            "\n**支持的目标语言：** 英文(English)、日文(Japanese)、韩文(Korean)、"
            "法文(French)、德文(German)、西班牙语(Spanish)、俄语(Russian)、阿拉伯语(Arabic)、"
            "简体中文(Chinese Simplified)、繁体中文(Chinese Traditional) 等\n"
            "\n⚠️ **禁止**：不要生成翻译代码，不要说'请运行以下代码'，系统会自动完成翻译。"
        ),
        "enabled": False,
    },
    {
        "id": "table_extractor",
        "name": "表格结构化提取",
        "icon": "🔲",
        "category": "domain",
        "skill_nature": "domain_skill",
        "description": "从 PDF/DOCX/XLSX 中精准提取表格数据，转换为结构化的 Markdown 或 JSON",
        "intent_description": "用户要提取文档中的表格、把表格转成结构化数据、读取表格内容",
        "task_types": ["RESEARCH", "FILE_GEN", "CODER"],
        "priority": 69,
        "prompt": (
            "\n\n## 🔲 工具规范：表格结构化提取\n"
            "\n**按来源选择提取方案：**\n"
            "\n**1. 来自 Excel (.xlsx/.xls)**\n"
            "```\n"
            "read_file_snippet(path, max_chars=6000)  # 返回前 200 行\n"
            "```\n"
            "读完后整理表头 → 输出 Markdown 表格 → 附统计摘要（行数/列数/空值数）\n"
            "\n**2. 来自 Word (.docx) 中的表格**\n"
            "```python\n"
            "# CODER 路径\n"
            "from docx import Document\n"
            "doc = Document('文件路径.docx')\n"
            "for t_idx, table in enumerate(doc.tables):\n"
            "    print(f'=== 第{t_idx+1}个表格 (共{len(table.rows)}行 × {len(table.columns)}列) ===')\n"
            "    for row in table.rows:\n"
            "        print('|', ' | '.join(c.text.strip() for c in row.cells), '|')\n"
            "```\n"
            "\n**3. 来自 PDF 中的表格**\n"
            "```python\n"
            "# CODER 路径 (需要 pdfplumber)\n"
            "import pdfplumber\n"
            "with pdfplumber.open('文件路径.pdf') as pdf:\n"
            "    for page_no, page in enumerate(pdf.pages, 1):\n"
            "        tables = page.extract_tables()\n"
            "        for t_idx, table in enumerate(tables):\n"
            "            print(f'=== 第{page_no}页 第{t_idx+1}个表格 ===')\n"
            "            for row in table:\n"
            "                print('|', ' | '.join(str(c or '') for c in row), '|')\n"
            "```\n"
            "注：若 pdfplumber 未安装，先用 `read_file_snippet` 读取，然后人工识别表格结构\n"
            "\n**输出规范：**\n"
            "- 始终用 Markdown 表格格式输出（第一行为表头）\n"
            "- 末尾注明：共 X 行 × Y 列，X 个空值"
        ),
        "enabled": False,
    },
]

# 所有合法的 category 和 task_type
SKILL_CATEGORIES = {
    "behavior": "⚙️ 行为",
    "style": "🎨 风格",
    "domain": "🔬 领域",
    "workflow": "🔄 工作流",
    "memory": "🧠 记忆",
    "custom": "🛠️ 自定义",
}

# skill_nature 枚举及其 UI 说明
SKILL_NATURE_LABELS = {
    "model_hint": "💬 模型行为调整",  # 通过 prompt 激活模型原生能力
    "domain_skill": "🔧 领域专项技能",  # 注入领域知识/专有模板
    "system": "⚙️ 系统功能",  # 记忆/工具等系统级功能
}

ALL_TASK_TYPES = [
    "CHAT",
    "CODER",
    "RESEARCH",
    "FILE_GEN",
    "SYSTEM",
    "AGENT",
    "WEB_SEARCH",
    "PAINTER",
    "DOC_ANNOTATE",
]


class SkillManager:
    """
    Skills 管理器 v2
    - 单例模式加载，首次访问时初始化
    - 启用/禁用状态持久化到 config/user_settings.json
    - 新增：_def_registry 存储原子化 SkillDefinition 对象（MCP 兼容）
    - 新增：register_custom / list_mcp_tools / load_custom_skills_dir
    """

    _registry: Dict[str, Dict] = {}  # id → 旧版 skill dict（向后兼容）
    _def_registry: Dict[str, SkillDefinition] = {}  # id → 新版 SkillDefinition（v2）
    _builtin_prompt_index: Dict[str, str] = (
        {}
    )  # id → original built-in prompt (O(1) lookup)
    _initialized: bool = False
    # 单轮注入的用户启用 Skill 数上限（系统 Skill 不计入），防止 token 膨胀
    _MAX_ACTIVE_INJECT: int = 20

    # ── 初始化 ─────────────────────────────────────────────────────────────────
    @classmethod
    def _ensure_init(cls):
        if cls._initialized:
            return
        cls._registry = {}
        cls._def_registry = {}
        cls._builtin_prompt_index = {s["id"]: s["prompt"] for s in BUILTIN_SKILLS}
        for skill in BUILTIN_SKILLS:
            s = dict(skill)  # shallow copy
            cls._registry[s["id"]] = s
            # 同步升级到 SkillDefinition
            cls._def_registry[s["id"]] = SkillDefinition.from_legacy_dict(s)
        cls._load_states_from_settings()
        # 从 config/skills/ 目录加载用户自定义 Skill（如果存在）
        cls._load_custom_skills_dir()
        cls._initialized = True

    @classmethod
    def _settings_path(cls) -> Path:
        """返回 config/user_settings.json 的绝对路径"""
        import sys

        if getattr(sys, "frozen", False):
            # 打包模式：config/ 紧邻 Koto.exe，不在 _internal/ 里
            project_root = Path(sys.executable).parent
        else:
            here = Path(__file__).resolve()
            # app/core/skills/skill_manager.py → project_root/config/user_settings.json
            project_root = here.parents[3]
        return project_root / "config" / "user_settings.json"

    @classmethod
    def _load_states_from_settings(cls):
        """从 user_settings.json 读取持久化的启用状态"""
        try:
            p = cls._settings_path()
            if not p.exists():
                return
            with open(p, "r", encoding="utf-8") as f:
                data = json.load(f)
            skills_state = data.get("skills", {})
            for skill_id, state in skills_state.items():
                if skill_id in cls._registry and isinstance(state, dict):
                    if "enabled" in state:
                        cls._registry[skill_id]["enabled"] = bool(state["enabled"])
                    if "prompt_override" in state and state["prompt_override"]:
                        cls._registry[skill_id]["prompt"] = state["prompt_override"]
        except Exception as e:
            print(f"[SkillManager] 加载设置失败: {e}")

    @classmethod
    def _save_states_to_settings(cls):
        """将当前启用状态写回 user_settings.json"""
        try:
            p = cls._settings_path()
            p.parent.mkdir(parents=True, exist_ok=True)
            data = {}
            if p.exists():
                with open(p, "r", encoding="utf-8") as f:
                    data = json.load(f)
            skills_state = {}
            for skill_id, skill in cls._registry.items():
                state: Dict = {"enabled": skill["enabled"]}
                # 如果有自定义 prompt，也保存
                builtin_prompt = cls._builtin_prompt_index.get(skill_id)
                if skill["prompt"] != builtin_prompt:
                    state["prompt_override"] = skill["prompt"]
                skills_state[skill_id] = state
            data["skills"] = skills_state
            with open(p, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"[SkillManager] 保存设置失败: {e}")

    # ── 公开 API ───────────────────────────────────────────────────────────────
    @classmethod
    def list_skills(cls) -> List[Dict]:
        """
        返回所有技能的完整信息列表（内置 + 自定义）
        """
        cls._ensure_init()
        result = []
        seen_ids: set = set()

        # 内置 Skill（保持原有顺序）
        for skill in BUILTIN_SKILLS:
            sid = skill["id"]
            seen_ids.add(sid)
            s = cls._registry.get(sid, skill)
            builtin_prompt = skill["prompt"]
            result.append(
                {
                    "id": s["id"],
                    "name": s["name"],
                    "icon": s["icon"],
                    "category": s["category"],
                    "skill_nature": s.get("skill_nature", "domain_skill"),
                    "description": s["description"],
                    "task_types": s["task_types"],
                    "enabled": s["enabled"],
                    "has_custom_prompt": s.get("prompt") != builtin_prompt,
                    "prompt": s["prompt"],
                    "is_builtin": True,
                }
            )

        # 自定义 Skill（不在内置列表中的）
        for skill_id, s in cls._registry.items():
            if skill_id in seen_ids:
                continue
            result.append(
                {
                    "id": s["id"],
                    "name": s["name"],
                    "icon": s["icon"],
                    "category": s["category"],
                    "description": s["description"],
                    "task_types": s["task_types"],
                    "enabled": s["enabled"],
                    "skill_nature": s.get("skill_nature", "domain_skill"),
                    "has_custom_prompt": False,
                    "prompt": s.get("prompt", ""),
                    "is_builtin": False,
                }
            )

        return result

    @classmethod
    def set_enabled(cls, skill_id: str, enabled: bool) -> bool:
        """启用或禁用一个技能，立即持久化"""
        cls._ensure_init()
        if skill_id not in cls._registry:
            return False
        cls._registry[skill_id]["enabled"] = enabled
        cls._save_states_to_settings()
        print(f"[SkillManager] {'✅ 启用' if enabled else '⏸️ 禁用'} skill: {skill_id}")
        return True

    @classmethod
    def update_prompt(cls, skill_id: str, prompt: str) -> bool:
        """更新某个技能的自定义 Prompt 内容（用户可自定义后保存）"""
        cls._ensure_init()
        if skill_id not in cls._registry:
            return False
        cls._registry[skill_id]["prompt"] = prompt
        cls._save_states_to_settings()
        return True

    @classmethod
    def reset_prompt(cls, skill_id: str) -> bool:
        """将某个技能的 Prompt 恢复为内置默认值"""
        cls._ensure_init()
        if skill_id not in cls._registry:
            return False
        builtin_prompt = cls._builtin_prompt_index.get(skill_id)
        if builtin_prompt is not None:
            cls._registry[skill_id]["prompt"] = builtin_prompt
            cls._save_states_to_settings()
        return True

    @classmethod
    def check_conflicts(cls, task_type: Optional[str] = None) -> List[Dict]:
        """
        检测当前所有已启用 Skills 之间的冲突关系。

        返回冲突列表，每项格式：
          {
            "winner_id":   "concise_mode",      # 优先级更高（保留）
            "winner_name": "精简模式",
            "loser_id":    "research_depth",    # 被压制（注入时跳过）
            "loser_name":  "深度研究",
            "reason":      "concise_mode 优先级(90) > research_depth 优先级(55)"
          }
        """
        cls._ensure_init()
        conflicts: List[Dict] = []
        seen_pairs: set = set()

        enabled_skills = {
            sid: s for sid, s in cls._registry.items() if s.get("enabled", False)
        }
        # 筛选适用于当前 task_type 的
        if task_type:
            tt = task_type.upper()
            enabled_skills = {
                sid: s
                for sid, s in enabled_skills.items()
                if not s.get("task_types") or tt in s.get("task_types", [])
            }

        for skill_id, s in enabled_skills.items():
            for conflict_id in s.get("conflict_with", []):
                if conflict_id not in enabled_skills:
                    continue
                pair = tuple(sorted([skill_id, conflict_id]))
                if pair in seen_pairs:
                    continue
                seen_pairs.add(pair)

                other = enabled_skills[conflict_id]
                pri_a = s.get("priority", 50)
                pri_b = other.get("priority", 50)

                if pri_a >= pri_b:
                    winner, loser = s, other
                    winner_id, loser_id = skill_id, conflict_id
                else:
                    winner, loser = other, s
                    winner_id, loser_id = conflict_id, skill_id

                conflicts.append(
                    {
                        "winner_id": winner_id,
                        "winner_name": winner.get("name", winner_id),
                        "winner_priority": winner.get("priority", 50),
                        "loser_id": loser_id,
                        "loser_name": loser.get("name", loser_id),
                        "loser_priority": loser.get("priority", 50),
                        "reason": (
                            f"「{winner.get('name', winner_id)}」优先级({winner.get('priority', 50)}) "
                            f"≥「{loser.get('name', loser_id)}」优先级({loser.get('priority', 50)})，"
                            f"后者 prompt 在本次请求中被抑制"
                        ),
                    }
                )
        return conflicts

    @classmethod
    def inject_into_prompt(
        cls,
        base_instruction: str,
        task_type: Optional[str] = None,
        user_input: Optional[str] = None,
        temp_skill_ids: Optional[List[str]] = None,
    ) -> str:
        """
        将当前启用的、适用于 task_type 的 Skills 注入到 base_instruction 末尾。

        冲突处理：当两个互相冲突的 Skill 同时启用时，优先级（priority）更高的
        skill 正常注入，低优先级的 skill 被静默抑制（但仍保持「启用」状态供用户查看）。

        Token 膨胀保护：单轮注入的用户启用 Skill 数（不含长期记忆）超过
        _MAX_ACTIVE_INJECT 时，按 priority 降序仅保留靠前的 skill，其余跳过。

        Args:
            base_instruction: 原始系统指令文本
            task_type:        当前任务类型（如 "CHAT"），None 表示通用
            user_input:       用户当前输入文本；传入后为长期记忆 skill
                              提供语义检索依据，精准命中相关记忆条目
            temp_skill_ids:   本轮临时激活的 Skill ID 列表（AutoMatcher 推荐）。
                              这些 Skill 即使 enabled=False 也会注入，且不修改
                              持久状态；注入后标注「自动匹配」供模型区分。

        Returns:
            注入后的系统指令文本
        """
        cls._ensure_init()
        active_prompts = []
        memory_block = ""
        seen_ids: set = set()  # 防止重复注入
        _inject_skill_count = 0  # 已注入的非系统 Skill 计数（用于上限保护）

        # ── 预计算冲突：找出所有因冲突被抑制的 skill_id ────────────────────
        suppressed_ids: set = set()
        all_enabled = {
            sid: s for sid, s in cls._registry.items() if s.get("enabled", False)
        }
        for sid, s in all_enabled.items():
            pri_a = s.get("priority", 50)
            for conflict_id in s.get("conflict_with", []):
                if conflict_id not in all_enabled:
                    continue
                pri_b = all_enabled[conflict_id].get("priority", 50)
                # 低优先级的那个被抑制
                if pri_a >= pri_b:
                    suppressed_ids.add(conflict_id)
                    logger.debug(
                        f"[SkillManager] 冲突抑制: {sid}(p={pri_a}) 抑制了 "
                        f"{conflict_id}(p={pri_b})"
                    )
                # 如果两者 priority 相等，按字母序决策（确保稳定性）
                elif pri_a == pri_b and sid < conflict_id:
                    suppressed_ids.add(conflict_id)

        # 遍历所有 registry 中的 Skill（内置 + 自定义），按 priority 排序
        all_skill_items = sorted(
            cls._registry.items(),
            key=lambda kv: kv[1].get("priority", 50),
            reverse=True,  # 高 priority 先注入
        )

        for skill_id, s in all_skill_items:
            if skill_id in seen_ids:
                continue

            if not s.get("enabled", False):
                continue

            seen_ids.add(skill_id)

            # 跳过被冲突抑制的 Skill
            if skill_id in suppressed_ids:
                logger.debug(f"[SkillManager] 跳过被抑制的 Skill: {skill_id}")
                continue

            applicable_types = s.get("task_types", [])
            if (
                applicable_types
                and task_type
                and task_type.upper() not in applicable_types
            ):
                continue

            # ── 长期记忆 skill：动态从 MemoryManager 检索相关记忆并注入 ────────
            if skill_id == "long_term_memory":
                try:
                    from web.memory_manager import MemoryManager

                    _mm = MemoryManager()
                    ctx = _mm.get_context_string(user_input or "")
                    if ctx.strip():
                        memory_block = ctx
                except Exception as _me:
                    logger.debug(f"[SkillManager] 长期记忆注入跳过: {_me}")
                continue  # 长期记忆不走普通 prompt 通道

            # ── 注入上限检测（长期记忆走了 continue，不计入此计数）─────────────
            if _inject_skill_count >= cls._MAX_ACTIVE_INJECT:
                logger.debug(
                    "[SkillManager] 注入上限 (%d) 已达，跳过低优先级 Skill: %s",
                    cls._MAX_ACTIVE_INJECT, skill_id,
                )
                continue

            # 优先使用新版 SkillDefinition 的 render_prompt()
            skill_def = cls._def_registry.get(skill_id)
            if skill_def and skill_def.system_prompt_template:
                p = skill_def.render_prompt().strip()
            else:
                p = s.get("prompt", "").strip()

            if p:
                # 注入 plan_template（仅在 prompt 中尚未包含执行步骤时追加，避免重复）
                pt = (
                    (getattr(skill_def, "plan_template", None) if skill_def else None)
                    or s.get("plan_template", [])
                )
                if pt and "执行步骤" not in p:
                    p = p + (
                        "\n\n### ⚙️ 执行步骤（必须严格按顺序完成）\n"
                        + "\n".join(f"{i+1}. {step}" for i, step in enumerate(pt))
                    )
                active_prompts.append(p)

            _inject_skill_count += 1

            # ── Word 模板 skill：追加模板字段说明 ─────────────────────────────
            tmpl_path_rel = s.get("template_path")
            if tmpl_path_rel:
                try:
                    import sys as _sys
                    from pathlib import Path as _Path

                    _base = (
                        _Path(_sys.executable).parent
                        if getattr(_sys, "frozen", False)
                        else _Path(__file__).resolve().parents[3]
                    )
                    tmpl_abs = _base / tmpl_path_rel
                    if not tmpl_abs.exists():
                        # 也尝试约定路径
                        tmpl_abs = (
                            _base
                            / "config"
                            / "skill_templates"
                            / skill_id
                            / "template.docx"
                        )
                    if tmpl_abs.exists():
                        from app.core.skills.template_engine import TemplateEngine

                        fields = TemplateEngine.parse_fields(tmpl_abs)
                        preview = TemplateEngine.get_raw_text(tmpl_abs)
                        tmpl_prompt = TemplateEngine.build_agent_prompt(
                            s.get("name", skill_id), fields, preview
                        )
                        active_prompts.append(tmpl_prompt)
                except Exception as _te:
                    logger.debug(f"[SkillManager] 模板提示注入跳过 ({skill_id}): {_te}")

        # ── 临时 Skill 注入（AutoMatcher 推荐，本轮生效，不改变持久状态）────
        auto_prompts = []
        _temp_ids = temp_skill_ids or []
        for skill_id in _temp_ids:
            if skill_id in seen_ids:
                continue  # 已作为用户启用 Skill 注入过，无需重复
            seen_ids.add(skill_id)
            s = cls._registry.get(skill_id)
            if not s:
                logger.debug(f"[SkillManager] 临时 Skill 不存在，跳过: {skill_id}")
                continue
            skill_def = cls._def_registry.get(skill_id)
            if skill_def and skill_def.system_prompt_template:
                p = skill_def.render_prompt().strip()
            else:
                p = s.get("prompt", "").strip()
            if p:
                # 临时 skill 也注入 plan_template（不注入 examples，避免 token 浪费）
                pt = getattr(skill_def, "plan_template", None) if skill_def else None
                if pt:
                    p = p + (
                        "\n\n### ⚙️ 执行步骤（必须严格按顺序完成）\n"
                        + "\n".join(f"{i+1}. {step}" for i, step in enumerate(pt))
                    )
                auto_prompts.append(p)
                logger.debug(f"[SkillManager] 🤖 临时注入 Auto-Skill: {skill_id}")

        # 组装注入块：记忆优先放在 skills 前面
        result = base_instruction
        if memory_block:
            result = (
                result + "\n\n─────────────────────────────────────────" + memory_block
            )
        if active_prompts:
            separator = "\n\n─────────────────────────────────────────"
            skills_block = (
                separator
                + "\n## 🎯 当前激活的 Skills（用户自定义行为）\n"
                + "\n".join(active_prompts)
            )
            result = result + skills_block
        if auto_prompts:
            separator = "\n\n─────────────────────────────────────────"
            auto_block = (
                separator
                + "\n## 🤖 自动匹配的 Skills（本轮智能推荐，仅本次生效）\n"
                + "\n".join(auto_prompts)
            )
            result = result + auto_block
        return result

    @classmethod
    def get_active_skill_names(cls, task_type: Optional[str] = None) -> List[str]:
        """返回当前启用的、适用于 task_type 的技能名称列表（含自定义 Skill）"""
        cls._ensure_init()
        names = []
        seen_ids: set = set()
        for skill_id, s in cls._registry.items():
            if skill_id in seen_ids:
                continue
            seen_ids.add(skill_id)
            if not s.get("enabled", False):
                continue
            applicable_types = s.get("task_types", [])
            if (
                applicable_types
                and task_type
                and task_type.upper() not in applicable_types
            ):
                continue
            names.append(s["name"])
        return names

    @classmethod
    def reload(cls):
        """强制重新加载（settings 文件被外部修改后调用）"""
        cls._initialized = False
        cls._ensure_init()

    # ═══════════════════════════════════════════════════════════════
    # ▼  v2 新增 API（SkillDefinition / MCP / 自定义 Skill 支持）
    # ═══════════════════════════════════════════════════════════════

    @classmethod
    def get_definition(cls, skill_id: str) -> Optional[SkillDefinition]:
        """
        返回指定 skill 的完整 SkillDefinition 对象。
        用于路由决策、变量渲染、输出验收。
        """
        cls._ensure_init()
        return cls._def_registry.get(skill_id)

    @classmethod
    def list_mcp_tools(cls) -> List[Dict]:
        """
        将所有 **已启用** 的 Skill 导出为 MCP (Model Context Protocol) 兼容的
        Tool 描述列表，可直接传给支持 MCP 的 LLM host 或外部系统。
        """
        cls._ensure_init()
        tools = []
        for skill_id, skill_def in cls._def_registry.items():
            # 同步最新启用状态
            legacy = cls._registry.get(skill_id)
            if legacy:
                skill_def.enabled = legacy.get("enabled", False)
            if skill_def.enabled:
                tools.append(skill_def.to_mcp_tool())
        return tools

    @classmethod
    def list_all_mcp_tools(cls) -> List[Dict]:
        """
        导出所有 Skill（不论是否启用）的 MCP Tool 描述列表。
        用于 Studio / Marketplace 展示。
        """
        cls._ensure_init()
        return [skill_def.to_mcp_tool() for skill_def in cls._def_registry.values()]

    @classmethod
    def register_custom(cls, skill_def: SkillDefinition) -> bool:
        """
        注册一个新的自定义 Skill（运行时动态注册）。
        同时写入 _registry 和 _def_registry，并持久化到 config/skills/{id}.json。

        Args:
            skill_def: 完整的 SkillDefinition 对象

        Returns:
            True 成功，False 失败（id 已存在于 builtin 且 author=="builtin" 时拒绝覆盖）
        """
        cls._ensure_init()
        existing = cls._def_registry.get(skill_def.id)
        if existing and existing.author == "builtin":
            logger.warning(
                f"[SkillManager] 拒绝覆盖内置 Skill: {skill_def.id}。"
                f"如需修改内置 Skill，请先 fork 并使用不同 id。"
            )
            return False

        skill_def.author = skill_def.author or "user"
        cls._def_registry[skill_def.id] = skill_def

        # 同步到旧版 _registry（保证 inject_into_prompt 等方法正常工作）
        cls._registry[skill_def.id] = {
            "id": skill_def.id,
            "name": skill_def.name,
            "icon": skill_def.icon,
            "category": (
                skill_def.category.value
                if hasattr(skill_def.category, "value")
                else skill_def.category
            ),
            "description": skill_def.description,
            "task_types": skill_def.task_types,
            "prompt": skill_def.render_prompt(),
            "enabled": skill_def.enabled,
            # 执行层增强字段（供 inject_into_prompt auto-skill 路径使用）
            "executor_tools": list(getattr(skill_def, "executor_tools", None) or []),
            "plan_template": list(getattr(skill_def, "plan_template", None) or []),
        }

        cls._apply_default_triggers(skill_def)

        # 持久化到 config/skills/{id}.json
        cls._persist_custom_skill(skill_def)
        logger.info(
            f"[SkillManager] ✅ 注册自定义 Skill: {skill_def.id} (v{skill_def.version})"
        )
        return True

    @classmethod
    def _apply_default_triggers(cls, skill_def: SkillDefinition):
        """Register manifest v2 default triggers once per skill/config pair."""
        default_triggers = list(getattr(skill_def, "default_triggers", None) or [])
        if not default_triggers:
            return

        try:
            from app.core.skills.skill_trigger_binding import get_skill_binding_manager

            binding_manager = get_skill_binding_manager()
            existing = binding_manager.list_bindings(
                skill_id=skill_def.id,
                binding_type="trigger",
            )
            existing_keys = {
                (
                    binding.trigger_type,
                    json.dumps(
                        binding.trigger_config or {}, sort_keys=True, ensure_ascii=False
                    ),
                )
                for binding in existing
            }

            for trigger in default_triggers:
                trigger_type = (
                    trigger.get("trigger_type") or trigger.get("type") or ""
                ).strip()
                trigger_config = trigger.get("config") or {}
                if not trigger_type:
                    continue

                trigger_key = (
                    trigger_type,
                    json.dumps(trigger_config, sort_keys=True, ensure_ascii=False),
                )
                if trigger_key in existing_keys:
                    continue

                binding_manager.bind_trigger(
                    skill_id=skill_def.id,
                    trigger_type=trigger_type,
                    trigger_config=trigger_config,
                    mode=trigger.get("mode", "execute"),
                    job_payload=trigger.get("job_payload")
                    or {
                        "skill_id": skill_def.id,
                        "query": trigger.get("query") or f"执行技能: {skill_def.name}",
                    },
                    name=trigger.get("name"),
                )
                existing_keys.add(trigger_key)
        except Exception as exc:
            logger.warning(f"[SkillManager] 应用默认触发器失败: {exc}")

    @classmethod
    def _persist_custom_skill(cls, skill_def: SkillDefinition):
        """将自定义 Skill 写入 config/skills/{id}.json"""
        try:
            skills_dir = cls._settings_path().parent / "skills"
            skills_dir.mkdir(parents=True, exist_ok=True)
            skill_file = skills_dir / f"{skill_def.id}.json"
            with open(skill_file, "w", encoding="utf-8") as f:
                json.dump(skill_def.to_dict(), f, ensure_ascii=False, indent=2)
        except Exception as e:
            logger.error(f"[SkillManager] 自定义 Skill 持久化失败: {e}")

    @classmethod
    def _load_custom_skills_dir(cls):
        """从 config/skills/ 目录加载所有自定义 Skill JSON 文件"""
        try:
            skills_dir = cls._settings_path().parent / "skills"
            if not skills_dir.exists():
                return
            for skill_file in skills_dir.glob("*.json"):
                try:
                    with open(skill_file, "r", encoding="utf-8") as f:
                        data = json.load(f)
                    skill_def = SkillDefinition.from_dict(data)
                    if skill_def.id in cls._def_registry:
                        # 已有内置注册：将 JSON 文件的增强字段合并进去
                        # （trigger_keywords / plan_template / executor_tools / examples / prompt）
                        # 保留 enabled 状态由 _load_states_from_settings 管理，不从 JSON 覆盖
                        existing = cls._def_registry[skill_def.id]
                        if skill_def.trigger_keywords:
                            existing.trigger_keywords = skill_def.trigger_keywords
                        if skill_def.plan_template:
                            existing.plan_template = skill_def.plan_template
                        if skill_def.executor_tools:
                            existing.executor_tools = skill_def.executor_tools
                        if skill_def.examples:
                            existing.examples = skill_def.examples
                        if skill_def.prompt:
                            existing.prompt = skill_def.prompt
                            # 同步到 _registry 的 prompt 字段
                            reg_entry = cls._registry.get(skill_def.id)
                            if reg_entry:
                                reg_entry["prompt"] = skill_def.render_prompt()
                                reg_entry["plan_template"] = skill_def.plan_template
                        logger.debug(f"[SkillManager] 合并自定义增强字段到内置 Skill: {skill_def.id}")
                    else:
                        cls._def_registry[skill_def.id] = skill_def
                        entry = {
                            "id": skill_def.id,
                            "name": skill_def.name,
                            "icon": skill_def.icon,
                            "category": (
                                skill_def.category.value
                                if hasattr(skill_def.category, "value")
                                else skill_def.category
                            ),
                            "description": skill_def.description,
                            "task_types": skill_def.task_types,
                            "prompt": skill_def.render_prompt(),
                            "enabled": skill_def.enabled,
                            "plan_template": skill_def.plan_template,
                        }
                        # 保留 template_path 和 bound_tools（若 JSON 中有）
                        if data.get("template_path"):
                            entry["template_path"] = data["template_path"]
                        if data.get("bound_tools"):
                            entry["bound_tools"] = data["bound_tools"]
                        cls._registry[skill_def.id] = entry
                        logger.info(f"[SkillManager] 加载自定义 Skill: {skill_def.id}")
                except Exception as e:
                    logger.warning(f"[SkillManager] 加载 {skill_file.name} 失败: {e}")
        except Exception as e:
            logger.warning(f"[SkillManager] 加载自定义 Skill 目录失败: {e}")

    @classmethod
    def validate_output(cls, skill_id: str, text: str) -> tuple:
        """
        使用指定 Skill 的 OutputSpec 验收文本。
        返回 (passed: bool, reason: str)
        若 Skill 不存在或无 OutputSpec，默认通过。
        """
        cls._ensure_init()
        skill_def = cls._def_registry.get(skill_id)
        if not skill_def:
            return True, "Skill 不存在，跳过验收"
        return skill_def.output_spec.validate(text)

    @classmethod
    def get_intent_descriptions(cls) -> Dict[str, str]:
        """
        返回所有已启用 Skill 的 {id: intent_description} 映射。
        供 Qwen Router 在意图识别时参考，提升路由准确性。
        """
        cls._ensure_init()
        result = {}
        for skill_id, skill_def in cls._def_registry.items():
            legacy = cls._registry.get(skill_id, {})
            if (
                legacy.get("enabled", skill_def.enabled)
                and skill_def.intent_description
            ):
                result[skill_id] = skill_def.intent_description
        return result

    # ═══════════════════════════════════════════════════════════════
    # ▼  智能推荐：根据用户输入建议相关 Skill
    # ═══════════════════════════════════════════════════════════════

    @classmethod
    def suggest_skills(
        cls,
        user_input: str,
        task_type: Optional[str] = None,
        top_k: int = 3,
        exclude_enabled: bool = True,
    ) -> List[Dict]:
        """
        根据用户输入和任务类型，推荐最相关的未启用 Skill。

        算法：
        1. 将 user_input 与每个 Skill 的 intent_description + description + tags 做关键词匹配打分
        2. 同时检查 task_type 适配性
        3. 返回得分最高的 top_k 个 Skill

        Args:
            user_input:     用户当前消息文本
            task_type:      当前任务类型（如 "CHAT", "CODER"）
            top_k:          返回建议数量
            exclude_enabled: True = 只推荐未启用的（让用户决定是否开启）

        Returns:
            [{"id", "name", "icon", "description", "score", "reason"}, ...]
        """
        import re as _re

        cls._ensure_init()
        user_lower = user_input.lower()
        scores: List[Dict] = []

        for skill_id, skill_def in cls._def_registry.items():
            s = cls._registry.get(skill_id, {})

            # 排除已启用的（可选）
            if exclude_enabled and s.get("enabled", skill_def.enabled):
                continue

            # 检查 task_type 适配性
            applicable = skill_def.task_types or []
            if applicable and task_type and task_type.upper() not in applicable:
                continue

            score = 0.0
            matched_reasons: List[str] = []

            # 计算与 intent_description 的相关性
            intent = (skill_def.intent_description or "").lower()
            desc = (skill_def.description or "").lower()
            tags = " ".join(skill_def.tags).lower()
            combined = f"{intent} {desc} {tags}"

            # 词频匹配 — 提取用户输入的关键词
            words = _re.findall(r"[\u4e00-\u9fff]{2,}|[a-zA-Z]{3,}", user_lower)
            for word in words:
                if word in combined:
                    score += 1.0
                    matched_reasons.append(word)

            # 精确短语匹配加权
            for phrase_len in (4, 3, 2):
                for i in range(len(user_lower) - phrase_len + 1):
                    phrase = user_lower[i : i + phrase_len]
                    if phrase in combined:
                        score += phrase_len * 0.5
                        break

            # 名称完全或部分匹配
            skill_name_lower = (skill_def.name or "").lower()
            if any(w in skill_name_lower for w in words if len(w) >= 2):
                score += 2.0
                matched_reasons.append(f"名称匹配: {skill_def.name}")

            if score > 0:
                reason = "与「{name}」相关：{r}".format(
                    name=skill_def.name,
                    r=(
                        "、".join(dict.fromkeys(matched_reasons))[:50]
                        if matched_reasons
                        else "语义相关"
                    ),
                )
                scores.append(
                    {
                        "id": skill_id,
                        "name": skill_def.name,
                        "icon": skill_def.icon,
                        "description": skill_def.description,
                        "score": round(score, 2),
                        "reason": reason,
                        "category": (
                            skill_def.category.value
                            if hasattr(skill_def.category, "value")
                            else skill_def.category
                        ),
                    }
                )

        # 按得分降序排列，取 top_k
        scores.sort(key=lambda x: x["score"], reverse=True)
        return scores[:top_k]

    # ═══════════════════════════════════════════════════════════════
    # ▼  冲突检测：防止相互矛盾的 Skill 同时启用
    # ═══════════════════════════════════════════════════════════════

    # 已知冲突组 — 同组内同时启用 2+ 个则报冲突
    _CONFLICT_GROUPS: List[tuple] = [
        # 详细 vs 精简
        ("step_by_step", "concise_mode"),
        # 正式 vs 随意（无强制，作为警告提示）
        # 高幽默 vs 严谨模式（警告级别）
        ("strict_mode", "emoji_assist"),
    ]
    # 软冲突（警告但不阻止）
    _SOFT_CONFLICTS: Dict[str, List[str]] = {
        "concise_mode": ["step_by_step", "teaching_mode", "proactive_suggestions"],
        "step_by_step": ["concise_mode"],
        "strict_mode": ["creative_writing", "emoji_assist"],
        "creative_writing": ["strict_mode", "professional_tone", "data_analysis"],
        "professional_tone": ["creative_writing"],
    }

    @classmethod
    def detect_conflicts(cls, skill_id: str) -> Dict:
        """
        检测启用 skill_id 后是否与当前其他已启用 Skill 产生冲突。
        综合检查：内置冲突规则表 + SkillDefinition.conflict_with 声明字段

        Returns:
            {
              "has_conflict": bool,
              "hard_conflicts": [{"id", "name", "reason"}, ...],
              "soft_conflicts": [{"id", "name", "reason"}, ...],
            }
        """
        cls._ensure_init()
        hard: List[Dict] = []
        soft: List[Dict] = []

        # 当前已启用集合
        enabled_ids = {
            sid
            for sid, s in cls._registry.items()
            if s.get("enabled", False) and sid != skill_id
        }

        this_def = cls._def_registry.get(skill_id)
        this_name = this_def.name if this_def else skill_id

        # 声明式 conflict_with 字段检查
        if this_def and this_def.conflict_with:
            for other_id in this_def.conflict_with:
                if other_id in enabled_ids:
                    other_def = cls._def_registry.get(other_id)
                    other_name = other_def.name if other_def else other_id
                    hard.append(
                        {
                            "id": other_id,
                            "name": other_name,
                            "reason": f"「{this_name}」声明与「{other_name}」不兼容",
                        }
                    )

        # 内置硬冲突规则
        hard_ids = {h["id"] for h in hard}
        for group in cls._CONFLICT_GROUPS:
            if skill_id in group:
                others = [g for g in group if g != skill_id]
                for other_id in others:
                    if other_id in enabled_ids and other_id not in hard_ids:
                        other_def = cls._def_registry.get(other_id)
                        other_name = other_def.name if other_def else other_id
                        hard.append(
                            {
                                "id": other_id,
                                "name": other_name,
                                "reason": f"「{this_name}」与「{other_name}」行为相反，同时启用会产生矛盾",
                            }
                        )
                        hard_ids.add(other_id)

        # 软冲突检测
        soft_list = cls._SOFT_CONFLICTS.get(skill_id, [])
        for other_id in soft_list:
            if other_id in enabled_ids and other_id not in hard_ids:
                other_def = cls._def_registry.get(other_id)
                other_name = other_def.name if other_def else other_id
                soft.append(
                    {
                        "id": other_id,
                        "name": other_name,
                        "reason": f"与「{other_name}」可能存在风格不一致，建议选其一",
                    }
                )

        return {
            "has_conflict": bool(hard) or bool(soft),
            "hard_conflicts": hard,
            "soft_conflicts": soft,
        }

    # ═══════════════════════════════════════════════════════════════
    # ▼  响应验收：对 LLM 回复批量验证所有激活 Skill 的 OutputSpec
    # ═══════════════════════════════════════════════════════════════

    @classmethod
    def validate_response(
        cls,
        text: str,
        task_type: Optional[str] = None,
    ) -> Dict:
        """
        对 LLM 生成的回复文本，批量检验所有当前激活 Skill 的 OutputSpec。

        Args:
            text:      LLM 生成的回复文本
            task_type: 当前任务类型

        Returns:
            {
              "all_passed": bool,
              "results": [{"skill_id", "skill_name", "passed", "reason"}, ...]
            }
        """
        cls._ensure_init()
        results = []
        all_passed = True

        for skill_id, skill_def in cls._def_registry.items():
            s = cls._registry.get(skill_id, {})
            if not s.get("enabled", skill_def.enabled):
                continue
            applicable = skill_def.task_types or []
            if applicable and task_type and task_type.upper() not in applicable:
                continue
            # 若 OutputSpec 约束为空（默认），跳过
            spec = skill_def.output_spec
            has_constraint = (
                spec.must_contain
                or spec.must_not_contain
                or spec.min_chars
                or spec.max_chars
                or spec.required_json_keys
            )
            if not has_constraint:
                continue

            passed, reason = spec.validate(text)
            if not passed:
                all_passed = False
            results.append(
                {
                    "skill_id": skill_id,
                    "skill_name": skill_def.name,
                    "passed": passed,
                    "reason": reason,
                }
            )

        return {"all_passed": all_passed, "results": results}

    # ═══════════════════════════════════════════════════════════════
    # ▼  Skill 摘要：供调试 / UI 状态面板使用
    # ═══════════════════════════════════════════════════════════════

    @classmethod
    def get_status_summary(cls) -> Dict:
        """
        返回当前 Skill 库的完整状态摘要，供 UI/日志使用。
        """
        cls._ensure_init()
        total = len(cls._def_registry)
        enabled = sum(1 for s in cls._registry.values() if s.get("enabled", False))
        builtin_count = sum(
            1 for s in cls._def_registry.values() if s.author == "builtin"
        )
        custom_count = total - builtin_count
        active_names = cls.get_active_skill_names()

        return {
            "total": total,
            "enabled": enabled,
            "builtin": builtin_count,
            "custom": custom_count,
            "active_skill_names": active_names,
            "version": "2.1",
        }
